#!/usr/bin/env python3
"""
Générateur de 50 présentations PowerPoint pour apprendre la programmation Python
Du débutant complet à avancé
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ═══════════════════════════════════════════════════════════════
# DESIGN SYSTEM
# ═══════════════════════════════════════════════════════════════

class Theme:
    DARK       = RGBColor(0x0A, 0x0A, 0x2E)
    PRIMARY    = RGBColor(0x16, 0x21, 0x3E)
    SECONDARY  = RGBColor(0x1A, 0x1A, 0x4E)
    ACCENT_OR  = RGBColor(0xFF, 0x6B, 0x35)  # Orange
    ACCENT_CY  = RGBColor(0x00, 0xB4, 0xD8)  # Cyan
    ACCENT_PU  = RGBColor(0x7B, 0x2F, 0xF7)  # Purple
    ACCENT_GR  = RGBColor(0x06, 0xD6, 0xA0)  # Green
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT      = RGBColor(0xF8, 0xF9, 0xFA)
    GRAY       = RGBColor(0x6C, 0x75, 0x7D)
    LIGHT_GRAY = RGBColor(0xDE, 0xE2, 0xE6)
    CODE_BG    = RGBColor(0x1E, 0x1E, 0x2E)
    CODE_TEXT  = RGBColor(0xCD, 0xD6, 0xF4)
    RED        = RGBColor(0xEF, 0x44, 0x44)
    YELLOW     = RGBColor(0xF9, 0xC7, 0x4C)

FONT_TITLE = 'Calibri'
FONT_BODY  = 'Calibri'
FONT_CODE  = 'Consolas'

# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def add_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=14,
                color=None, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color or Theme.WHITE
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_textbox(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox

def add_paragraph(tf, text, size=14, color=None, bold=False, name=FONT_BODY,
                  alignment=PP_ALIGN.LEFT, space_before=0, space_after=0,
                  italic=False, level=0):
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == '':
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color or Theme.WHITE
    p.font.bold = bold
    p.font.name = name
    p.font.italic = italic
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.level = level
    return p

def add_code_block(slide, left, top, width, height, code_lines, font_size=11):
    shape = add_rounded_rect(slide, left, top, width, height, Theme.CODE_BG)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                      width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in code_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.name = FONT_CODE
        p.font.size = Pt(font_size)
        p.font.color.rgb = Theme.CODE_TEXT
        p.space_before = Pt(0)
        p.space_after = Pt(0)
    return txBox

def make_title_slide(prs, title, subtitle, part=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_slide_bg(slide, Theme.DARK)

    # Decorative top bar
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), Theme.ACCENT_OR)
    
    # Decorative side accent
    add_shape(slide, Inches(0), Inches(0), Inches(0.06), prs.slide_height, Theme.ACCENT_CY)

    # Part number
    if part:
        add_textbox(slide, Inches(1.5), Inches(2.0), Inches(8), Inches(0.5),
                    part, 16, Theme.ACCENT_OR, bold=True)

    # Title
    add_textbox(slide, Inches(1.5), Inches(2.6), Inches(8), Inches(1.5),
                title, 40, Theme.WHITE, bold=True)
    
    # Subtitle
    add_textbox(slide, Inches(1.5), Inches(4.2), Inches(8), Inches(1),
                subtitle, 20, Theme.GRAY)
    
    # Decorative line
    add_shape(slide, Inches(1.5), Inches(5.2), Inches(3), Inches(0.04), Theme.ACCENT_OR)
    
    # Bottom info
    add_textbox(slide, Inches(1.5), Inches(6.5), Inches(8), Inches(0.5),
                "Cours de Programmation Python - Fondamentaux", 12, Theme.GRAY)

def make_section_slide(prs, section_num, title, description=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, Theme.PRIMARY)
    
    # Large number
    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(2), Inches(2),
                f"{section_num:02d}", 72, Theme.ACCENT_OR, bold=True)
    
    # Title
    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(8.5), Inches(1.2),
                title, 36, Theme.WHITE, bold=True)
    
    # Decorative line
    add_shape(slide, Inches(0.8), Inches(4.3), Inches(4), Inches(0.04), Theme.ACCENT_CY)
    
    if description:
        add_textbox(slide, Inches(0.8), Inches(4.6), Inches(8.5), Inches(1.5),
                    description, 18, Theme.LIGHT_GRAY)
    
    # Bottom decoration
    add_shape(slide, Inches(0), Inches(7.1), prs.slide_width, Inches(0.06), Theme.ACCENT_GR)

def make_content_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, Theme.LIGHT)
    
    # Top bar
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), Theme.SECONDARY)
    add_shape(slide, Inches(0), Inches(0.9), prs.slide_width, Inches(0.04), Theme.ACCENT_OR)
    
    # Title
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(9), Inches(0.7),
                title, 26, Theme.WHITE, bold=True)
    
    # Content area
    txBox = add_rich_textbox(slide, Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.5))
    tf = txBox.text_frame
    first = True
    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        if isinstance(bullet, tuple):
            text, level = bullet
            p.text = text
            p.level = level
        else:
            p.text = bullet
            p.level = 0
        
        p.font.name = FONT_BODY
        p.font.color.rgb = Theme.DARK
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        
        if p.level == 0:
            p.font.size = Pt(18)
            p.font.bold = True
        else:
            p.font.size = Pt(15)
            p.font.bold = False
    
    # Notes
    if notes:
        add_textbox(slide, Inches(0.6), Inches(6.8), Inches(8.8), Inches(0.5),
                    f"💡 {notes}", 11, Theme.GRAY)
    
    # Bottom accent
    add_shape(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.04), Theme.ACCENT_CY)
    return slide

def make_code_slide(prs, title, code_lines, explanation=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, Theme.LIGHT)
    
    # Top bar
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), Theme.SECONDARY)
    add_shape(slide, Inches(0), Inches(0.9), prs.slide_width, Inches(0.04), Theme.ACCENT_GR)
    
    # Title with code icon
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(9), Inches(0.7),
                f"  {title}", 24, Theme.WHITE, bold=True)
    
    # Code block
    add_code_block(slide, Inches(0.6), Inches(1.2), Inches(8.8), Inches(4.5), code_lines, 13)
    
    # Explanation
    if explanation:
        txBox = add_rich_textbox(slide, Inches(0.6), Inches(5.9), Inches(8.8), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"💡 {explanation}"
        p.font.size = Pt(13)
        p.font.color.rgb = Theme.GRAY
        p.font.name = FONT_BODY
    
    # Bottom accent
    add_shape(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.04), Theme.ACCENT_GR)
    return slide

def make_exercise_slide(prs, title, exercise_text, solution_code=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, Theme.LIGHT)
    
    # Top bar - Purple theme
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), Theme.ACCENT_PU)
    add_shape(slide, Inches(0), Inches(0.9), prs.slide_width, Inches(0.04), Theme.ACCENT_OR)
    
    # Title
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(9), Inches(0.7),
                f"  {title}", 24, Theme.WHITE, bold=True)
    
    # Exercise content
    txBox = add_rich_textbox(slide, Inches(0.6), Inches(1.3), Inches(8.8), Inches(2.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = exercise_text
    p.font.size = Pt(16)
    p.font.color.rgb = Theme.DARK
    p.font.name = FONT_BODY
    
    # Solution code if provided
    if solution_code:
        add_code_block(slide, Inches(0.6), Inches(4.0), Inches(8.8), Inches(3.2), solution_code, 11)
        add_textbox(slide, Inches(0.6), Inches(3.7), Inches(3), Inches(0.3),
                    "Solution :", 12, Theme.ACCENT_PU, bold=True)
    
    # Bottom accent
    add_shape(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.04), Theme.ACCENT_PU)

def make_objectives_slide(prs, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, Theme.LIGHT)
    
    # Top bar
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), Theme.SECONDARY)
    add_shape(slide, Inches(0), Inches(0.9), prs.slide_width, Inches(0.04), Theme.ACCENT_OR)
    
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(9), Inches(0.7),
                "  Objectifs", 24, Theme.WHITE, bold=True)
    
    txBox = add_rich_textbox(slide, Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.5))
    tf = txBox.text_frame
    first = True
    for obj in objectives:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"✅  {obj}"
        p.font.size = Pt(17)
        p.font.color.rgb = Theme.DARK
        p.font.name = FONT_BODY
        p.space_before = Pt(8)
        p.space_after = Pt(4)
    
    add_shape(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.04), Theme.ACCENT_CY)

def make_summary_slide(prs, points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, Theme.DARK)
    
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.04), Theme.ACCENT_OR)
    
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(9), Inches(0.7),
                "  Ce qu'il faut retenir", 28, Theme.WHITE, bold=True)
    
    add_shape(slide, Inches(0.8), Inches(1.2), Inches(3), Inches(0.04), Theme.ACCENT_CY)
    
    txBox = add_rich_textbox(slide, Inches(0.8), Inches(1.6), Inches(8.4), Inches(5.0))
    tf = txBox.text_frame
    first = True
    for pt in points:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = f"✦  {pt}"
        p.font.size = Pt(16)
        p.font.color.rgb = Theme.LIGHT_GRAY
        p.font.name = FONT_BODY
        p.space_before = Pt(6)
        p.space_after = Pt(4)
    
    add_shape(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.04), Theme.ACCENT_GR)

def make_two_col_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, Theme.LIGHT)
    
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), Theme.SECONDARY)
    add_shape(slide, Inches(0), Inches(0.9), prs.slide_width, Inches(0.04), Theme.ACCENT_OR)
    
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(9), Inches(0.7),
                title, 24, Theme.WHITE, bold=True)
    
    # Left column
    add_shape(slide, Inches(0.4), Inches(1.2), Inches(4.3), Inches(0.04), Theme.ACCENT_CY)
    add_textbox(slide, Inches(0.4), Inches(1.3), Inches(4.3), Inches(0.4),
                left_title, 16, Theme.ACCENT_CY, bold=True)
    
    txBox = add_rich_textbox(slide, Inches(0.4), Inches(1.8), Inches(4.3), Inches(5.0))
    tf = txBox.text_frame
    first = True
    for item in left_items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = Theme.DARK
        p.font.name = FONT_BODY
        p.space_before = Pt(3)
        p.space_after = Pt(2)
    
    # Right column
    add_shape(slide, Inches(5.3), Inches(1.2), Inches(4.3), Inches(0.04), Theme.ACCENT_OR)
    add_textbox(slide, Inches(5.3), Inches(1.3), Inches(4.3), Inches(0.4),
                right_title, 16, Theme.ACCENT_OR, bold=True)
    
    txBox = add_rich_textbox(slide, Inches(5.3), Inches(1.8), Inches(4.3), Inches(5.0))
    tf = txBox.text_frame
    first = True
    for item in right_items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = Theme.DARK
        p.font.name = FONT_BODY
        p.space_before = Pt(3)
        p.space_after = Pt(2)
    
    add_shape(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.04), Theme.ACCENT_GR)

def make_definition_slide(prs, term, definition, example=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, Theme.LIGHT)
    
    # Top bar
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), Theme.SECONDARY)
    add_shape(slide, Inches(0), Inches(0.9), prs.slide_width, Inches(0.04), Theme.ACCENT_OR)
    
    add_textbox(slide, Inches(0.6), Inches(0.12), Inches(9), Inches(0.7),
                "  Définition", 24, Theme.WHITE, bold=True)
    
    # Term in big letters
    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(8.8), Inches(0.8),
                term, 30, Theme.ACCENT_OR, bold=True)
    
    add_shape(slide, Inches(0.6), Inches(2.3), Inches(2), Inches(0.03), Theme.ACCENT_CY)
    
    # Definition
    txBox = add_rich_textbox(slide, Inches(0.6), Inches(2.6), Inches(8.8), Inches(2.0))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = definition
    p.font.size = Pt(17)
    p.font.color.rgb = Theme.DARK
    p.font.name = FONT_BODY
    
    if example:
        add_rounded_rect(slide, Inches(0.6), Inches(4.8), Inches(8.8), Inches(2.2), RGBColor(0xE8, 0xF4, 0xFD))
        add_textbox(slide, Inches(0.8), Inches(4.9), Inches(1.5), Inches(0.3),
                    "Exemple :", 13, Theme.ACCENT_CY, bold=True)
        add_textbox(slide, Inches(0.8), Inches(5.3), Inches(8.4), Inches(1.5),
                    example, 15, Theme.DARK)
    
    add_shape(slide, Inches(0), Inches(7.35), prs.slide_width, Inches(0.04), Theme.ACCENT_CY)

def save_presentation(prs, filename):
    path = os.path.join(OUTPUT_DIR, filename)
    prs.save(path)
    print(f"  ✅ {filename}")

OUTPUT_DIR = "/home/akaletekoffilevis/Bureau/Coach/presentations"

# ═══════════════════════════════════════════════════════════════
# PARTIE 1 : CULTURE INFORMATIQUE & SETUP
# ═══════════════════════════════════════════════════════════════

def p01_histoire_informatique():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "C'est quoi l'informatique ?",
        "Histoire, importance et impact dans notre quotidien",
        "Partie 1 - Culture Informatique")

    make_objectives_slide(prs, [
        "Comprendre ce qu'est l'informatique simplement",
        "Découvrir l'histoire des ordinateurs (des machines géantes aux smartphones)",
        "Comprendre pourquoi l'informatique est partout autour de nous",
        "Se motiver pour apprendre la programmation !"
    ])

    make_content_slide(prs, "C'est quoi l'informatique ?", [
        ("Définition simple", 0),
        "L'informatique = INFORMATION + AUTOMATIQUE",
        "C'est la science du traitement automatique de l'information",
        ("Ça inclut :", 0),
        ("Les ordinateurs (PC, Mac, etc.)", 1),
        ("Les smartphones et tablettes", 1),
        ("Internet et les sites web", 1),
        ("Les jeux vidéo", 1),
        ("Les applications (Instagram, TikTok, WhatsApp...)", 1),
        ("Les voitures modernes, les robots, l'IA...", 1),
        "",
        "💡 En gros : tout ce qui utilise un processeur pour traiter des données !"
    ])

    make_content_slide(prs, "Un peu d'histoire - Les débuts", [
        ("1940s : Les premiers ordinateurs", 0),
        "ENIAC (1945) : 30 tonnes, 167 m², 17 468 tubes à vide",
        "Il occupait une pièce entière et consommait autant qu'une petite ville !",
        ("1960s : Les transistors et circuits intégrés", 0),
        "Les ordinateurs deviennent plus petits et plus fiables",
        "Invention du microprocesseur (Intel 4004, 1971)",
        ("1970s-80s : L'ère des ordinateurs personnels", 0),
        "Apple II (1977), IBM PC (1981), Commodore 64",
        "Les ordinateurs arrivent dans les maisons",
    ])

    make_content_slide(prs, "L'évolution rapide", [
        ("1990s : Internet et le Web", 0),
        "World Wide Web inventé par Tim Berners-Lee en 1989",
        "Les ordinateurs se connectent entre eux → révolution !",
        ("2000s : L'ère mobile", 0),
        "Smartphones (iPhone 2007), tablettes",
        "L'informatique devient portable et omniprésente",
        ("2010s-2020s : Intelligence Artificielle et Cloud", 0),
        "IA, Machine Learning, ChatGPT, assistants vocaux",
        "Le cloud computing, le streaming (Netflix, Spotify)",
        ("Aujourd'hui :", 0),
        "Plus de 5 milliards d'internautes dans le monde",
        "L'informatique est dans toutes les industries"
    ])

    make_content_slide(prs, "Pourquoi l'informatique est partout ?", [
        ("Dans la vie quotidienne :", 0),
        ("Communication : SMS, emails, réseaux sociaux, visioconférence", 1),
        ("Divertissement : Netflix, YouTube, jeux vidéo, musique", 1),
        ("Éducation : cours en ligne, recherche, outils scolaires", 1),
        ("Santé : dossiers médicaux, imagerie, robots chirurgicaux", 1),
        ("Transports : GPS, billets, voitures autonomes", 1),
        ("Commerce : Amazon, paiements en ligne, banque", 1),
        ("Travail : ordinateurs, logiciels, réunions à distance", 1),
        "",
        "🔑 L'informatique est devenue aussi importante que l'écriture ou le calcul !"
    ])

    make_content_slide(prs, "Les métiers de l'informatique", [
        ("Développeur / Programmeur : crée des logiciels, sites web, applications", 1),
        ("Ingénieur réseau : gère les connexions entre ordinateurs", 1),
        ("Data Scientist : analyse des données pour en tirer des informations", 1),
        ("Cybersécurité : protège les systèmes contre les pirates", 1),
        ("Designer UX/UI : conçoit l'apparence et l'ergonomie des interfaces", 1),
        ("Administrateur système : gère les serveurs et infrastructures", 1),
        ("Chef de projet tech : organise le travail des équipes techniques", 1),
        "",
        "🎯 Et toi, quel domaine te fait rêver ?"
    ])

    make_definition_slide(prs, "Programmer, c'est quoi ?",
        "Programmer = donner des instructions à un ordinateur pour qu'il exécute une tâche.",
        "Comme une recette de cuisine : tu donnes des étapes précises dans l'ordre pour obtenir un résultat. L'ordinateur suit scrupuleusement tes instructions !")

    make_summary_slide(prs, [
        "L'informatique est la science du traitement automatique de l'information",
        "Elle a évolué des énormes machines des années 40 aux smartphones d'aujourd'hui",
        "L'informatique est présente dans tous les aspects de notre vie",
        "Programmer = donner des instructions à un ordinateur",
        "C'est une compétence super utile et recherchée !"
    ])

    save_presentation(prs, "01_Cest_quoi_informatique.pptx")
    return prs


def p02_fonctionnement_ordinateur():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Comment fonctionne un ordinateur ?",
        "Les composants matériels (Hardware) expliqués simplement",
        "Partie 1 - Culture Informatique")

    make_objectives_slide(prs, [
        "Connaître les principales parties d'un ordinateur",
        "Comprendre le rôle de chaque composant",
        "Savoir comment ils travaillent ensemble",
        "Être capable de choisir un ordinateur adapté"
    ])

    make_definition_slide(prs, "Hardware vs Software",
        "Hardware = la partie physique (ce qu'on touche). Software = les programmes (ce qu'on ne touche pas).",
        "Hardware : écran, clavier, souris, processeur, disque dur...\nSoftware : Windows, Chrome, Python, jeux, applications...")

    make_content_slide(prs, "Les 4 fonctions d'un ordinateur", [
        ("Tout ordinateur fait 4 choses :", 0),
        "",
        ("1️⃣  ENTRÉE (Input)", 1),
        "  Récupérer des informations depuis l'extérieur (clavier, souris, micro, caméra)",
        "",
        ("2️⃣  TRAITEMENT (Processing)", 1),
        "  Calculer, analyser, transformer les données (le processeur = le cerveau)",
        "",
        ("3️⃣  SORTIE (Output)", 1),
        "  Afficher/envoyer le résultat (écran, haut-parleurs, imprimante)",
        "",
        ("4️⃣  STOCKAGE (Storage)", 1),
        "  Sauvegarder les informations pour plus tard (disque dur, SSD)",
    ])

    make_two_col_slide(prs, "Les composants internes",
        "Composant", [
            "🖥️  Carte mère (Motherboard)",
            "🧠  Processeur (CPU)",
            "📝  RAM (Mémoire vive)",
            "💾  Disque dur / SSD",
            "🎮  Carte graphique (GPU)",
            "🔌  Alimentation (PSU)",
        ],
        "Rôle",
        [
            "Connecte tous les composants entre eux",
            "Exécute les calculs et instructions = le cerveau",
            "Mémoire temporaire pour les tâches en cours",
            "Stockage permanent de tes fichiers et programmes",
            "Génère les images, vidéos et jeux (indispensable pour les gamers)",
            "Fournit l'électricité à tous les composants",
        ]
    )

    make_content_slide(prs, "Le Processeur (CPU) - Le cerveau", [
        ("CPU = Central Processing Unit", 0),
        "",
        "C'est le composant le plus important !",
        "Il exécute des milliards d'opérations par seconde (GHz = gigahertz)",
        "",
        ("Comment il fonctionne :", 0),
        ("Cycle : Récupérer → Décoder → Exécuter → Stocker", 1),
        ("1. Il récupère une instruction depuis la RAM", 1),
        ("2. Il décode ce qu'il doit faire", 1),
        ("3. Il exécute l'instruction (calcul, déplacement...)", 1),
        ("4. Il stocke le résultat", 1),
        "",
        "Marques principales : Intel (Core i5, i7, i9) et AMD (Ryzen 5, 7, 9)",
        "Plus le GHz est élevé, plus le processeur est rapide !"
    ])

    make_content_slide(prs, "La RAM - La mémoire de travail", [
        ("RAM = Random Access Memory", 0),
        "",
        "C'est la mémoire TEMPORAIRE (elle s'efface quand on éteint)",
        "Elle stocke ce que l'ordinateur est en train de faire MAINTENANT",
        "",
        ("Analogies :", 0),
        ("C'est comme ton bureau de travail : plus il est grand, plus tu peux étaler de documents", 1),
        ("Le disque dur c'est l'armoire de rangement (pour plus tard)", 1),
        "",
        ("Conseils pratiques :", 0),
        "8 Go minimum pour un usage normal (navigation, bureautique)",
        "16 Go recommandé pour le gaming ou la programmation",
        "Plus de RAM = pouvoir faire plus de choses en même temps !",
        "",
        "⚡ La RAM est BEAUCOUP plus rapide que le disque dur, mais elle est volatile"
    ])

    make_content_slide(prs, "Le stockage : HDD vs SSD", [
        ("Disque Dur (HDD) :", 0),
        "Ancienne technologie (plateaux qui tournent et tête de lecture)",
        "Moins cher, plus lent, sensible aux chocs",
        "Capacité : 500 Go à 4 To typiquement",
        "",
        ("SSD (Solid State Drive) :", 0),
        "Technologie moderne (puce mémoire, sans parties mobiles)",
        "Plus cher, mais BEAUCOUP plus rapide",
        "Ton PC démarre en 10 secondes au lieu de 2 minutes !",
        "Capacité : 128 Go à 2 To typiquement",
        "",
        ("Lequel choisir ?", 0),
        "👉 SSD pour le système et les logiciels (rapidité au quotidien)",
        "👉 HDD pour les fichiers volumineux (photos, vidéos) si besoin de place",
        "👉 Le meilleur des deux : SSD + HDD combinés !"
    ])

    make_content_slide(prs, "Comment tout ça communique ?", [
        ("Le chemin d'une action simple :", 0),
        "",
        ("Tu appuies sur une touche du clavier", 1),
        ("→ Le signal va à la carte mère via un câble USB", 1),
        ("→ Le processeur récupère l'info et la traite", 1),
        ("→ Il va chercher les programmes nécessaires dans la RAM", 1),
        ("→ Si besoin, il lit/écrit sur le SSD pour charger des données", 1),
        ("→ La carte graphique génère l'image à afficher", 1),
        ("→ L'écran affiche le résultat (la lettre apparaît !)", 1),
        "",
        "⏱️ Tout ça se produit en quelques MILLISECONDES (millièmes de seconde) !",
        "Le processeur peut faire ça des milliards de fois par seconde."
    ])

    make_content_slide(prs, "Les périphériques essentiels", [
        ("Entrée (Input) :", 0),
        ("Clavier, souris, pavé tactile, micro, webcam, scanner, manette", 1),
        ("Les capteurs (tactiles, accéléromètres, GPS...)", 1),
        "",
        ("Sortie (Output) :", 0),
        ("Écran / moniteur, haut-parleurs / casque, imprimante", 1),
        ("Casque VR, projecteur", 1),
        "",
        ("Entrée-Sortie (Input/Output) :", 0),
        ("Écran tactile, disque dur externe, clé USB, carte réseau", 1),
        ("Imprimante multifonction, casque avec micro", 1),
        "",
        ("💡 Chaque périphérique a besoin d'un PILOTE (driver) pour communiquer avec l'OS")
    ])

    make_summary_slide(prs, [
        "Un ordinateur = entrée + traitement + sortie + stockage",
        "CPU = le cerveau, RAM = mémoire de travail, SSD = mémoire permanente",
        "La carte mère connecte tous les composants entre eux",
        "SSD est beaucoup plus rapide que HDD",
        "Plus de RAM et un meilleur CPU = ordinateur plus performant",
        "Tout fonctionne ensemble en quelques millisecondes !"
    ])

    save_presentation(prs, "02_Fonctionnement_ordinateur.pptx")
    return prs


def p03_systemes_exploitation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Les Systèmes d'Exploitation",
        "Windows, Linux, macOS - Le chef d'orchestre de l'ordinateur",
        "Partie 1 - Culture Informatique")

    make_objectives_slide(prs, [
        "Comprendre le rôle d'un système d'exploitation (OS)",
        "Connaître les OS principaux : Windows, macOS, Linux",
        "Savoir quel OS choisir pour la programmation",
        "Apprendre à interagir avec son OS"
    ])

    make_definition_slide(prs, "Qu'est-ce qu'un Système d'Exploitation ?",
        "Un OS (Operating System) est le logiciel principal qui fait fonctionner l'ordinateur. C'est l'intermédiaire entre le matériel (hardware) et les logiciels (software).",
        "Sans OS, un ordinateur n'est qu'un tas de métal et de plastique inutile !\nExemples : Windows 11, macOS Ventura, Ubuntu Linux, Android, iOS")

    make_content_slide(prs, "Le rôle d'un OS", [
        ("1️⃣ GÉRER LE MATÉRIEL (Hardware)", 0),
        ("Le processeur : décider quel programme utilise le CPU à quel moment", 1),
        ("La RAM : allouer la mémoire aux programmes qui en ont besoin", 1),
        ("Les périphériques : gérer les entrées/sorties (clavier, souris, écran)", 1),
        "",
        ("2️⃣ FOURNIR UNE INTERFACE UTILISATEUR", 0),
        ("Graphique (GUI) : fenêtres, icônes, boutons, souris", 1),
        ("En ligne de commande (CLI) : terminal, shell", 1),
        "",
        ("3️⃣ GÉRER LES FICHIERS", 0),
        ("Organisation en dossiers/arborescence", 1),
        ("Permissions de lecture/écriture/exécution", 1),
        "",
        ("4️⃣ GÉRER LES UTILISATEURS", 0),
        ("Comptes séparés, mots de passe, permissions", 1),
    ])

    make_two_col_slide(prs, "Les 3 OS principaux",
        "Windows", [
            "✅ Le plus répandu sur PC (80% des ordinateurs)",
            "✅ Interface graphique simple et intuitive",
            "✅ Compatible avec beaucoup de jeux et logiciels",
            "❌ Payant (licence Windows)",
            "❌ Moins utilisé pour la programmation serveur",
        ],
        "macOS",
        [
            "✅ Design soigné, très stable et sécurisé",
            "✅ Écosystème Apple (iPhone, iPad, Mac)",
            "✅ Basé sur Unix (compatible Linux pour la programmation)",
            "❌ Uniquement sur les Mac (chers)",
            "❌ Moins de jeux que Windows",
        ]
    )

    make_content_slide(prs, "Linux - L'OS des programmeurs", [
        ("C'est quoi Linux ?", 0),
        "C'est un OS libre et gratuit créé par Linus Torvalds en 1991",
        "Le noyau (kernel) est open source → tout le monde peut voir le code !",
        "",
        ("Pourquoi les programmeurs aiment Linux ?", 0),
        ("100% gratuit - pas de licence à payer", 1),
        ("Open source - tu peux voir et modifier le code source", 1),
        ("Le terminal est très puissant - idéal pour développer", 1),
        ("Très stable et sécurisé - peu de virus !", 1),
        ("C'est l'OS des serveurs - la plupart des sites web tournent sur Linux", 1),
        "",
        ("Les distributions (versions) populaires :", 0),
        "Ubuntu (la plus connue, facile pour débuter)",
        "Debian, Fedora, Linux Mint, Arch Linux"
    ])

    make_content_slide(prs, "Interface Graphique (GUI) vs Ligne de Commande (CLI)", [
        ("GUI - Graphical User Interface", 0),
        "Tu utilises la souris pour cliquer sur des icônes, fenêtres, boutons",
        "Intuitif, facile à apprendre",
        "Tu utilises ça tout le temps sans t'en rendre compte !",
        "",
        ("CLI - Command Line Interface", 0),
        "Tu tapes des commandes textuelles dans un terminal",
        "Plus puissant et plus rapide pour certaines tâches",
        "Un peu intimidant au début mais très efficace !",
        "",
        ("Exemple : Créer un dossier 'projet'", 0),
        ("GUI : Clic droit → Nouveau dossier → renommer", 1),
        ("CLI : mkdir projet", 1),
        "",
        "💡 En programmation, la CLI est ton amie ! On l'apprendra bientôt."
    ])

    make_content_slide(prs, "Naviguer dans le système de fichiers", [
        ("Structure en arborescence :", 0),
        "",
        "C: (Windows) ou / (Linux/Mac)",
        "├── Utilisateurs (Windows) / home (Linux/Mac)",
        "│   ├── Documents",
        "│   ├── Images",
        "│   ├── Téléchargements",
        "│   └── Musique",
        "├── Programmes / Applications",
        "├── Système (System32 / etc, bin, lib...)",
        "└── Temp",
        "",
        ("📁 Les chemins :", 0),
        ("Windows : C:\\Users\\TonNom\\Documents", 1),
        ("Linux/Mac : /home/tonnom/Documents", 1),
        "",
        "💡 Les dossiers sont séparés par \\ sur Windows et / sur Linux/Mac"
    ])

    make_summary_slide(prs, [
        "L'OS est le chef d'orchestre qui gère tout l'ordinateur",
        "Windows : grand public, macOS : design et stabilité, Linux : liberté et puissance",
        "Linux est l'OS préféré des développeurs (gratuit, open source, puissant)",
        "GUI = cliquer, CLI = taper des commandes (plus efficace)",
        "Les fichiers sont organisés en arborescence",
        "Nous utiliserons principalement le terminal en programmation"
    ])

    save_presentation(prs, "03_Systemes_exploitation.pptx")
    return prs


def p04_internet_web():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Internet et le Web",
        "Comment ça marche ? Sites web, DNS, HTTP, navigateurs",
        "Partie 1 - Culture Informatique")

    make_objectives_slide(prs, [
        "Comprendre la différence entre Internet et le Web",
        "Savoir comment les données voyagent sur Internet",
        "Comprendre le rôle du DNS, HTTP, des navigateurs",
        "Savoir ce qui se passe quand tu tapes une URL"
    ])

    make_definition_slide(prs, "Internet ≠ Web",
        "Internet est le RÉSEAU (l'infrastructure physique). Le Web est un SERVICE qui utilise Internet (les sites web).",
        "Internet c'est comme le système routier (les autoroutes).\nLe Web c'est comme les magasins, restaurants, maisons accessibles par ces routes.\nAutres services sur Internet : emails, jeux en ligne, streaming, WhatsApp...")

    make_content_slide(prs, "Comment Internet fonctionne ?", [
        ("Un réseau de réseaux", 0),
        "Des milliers de kilomètres de câbles (sous la mer, sous terre)",
        "Des serveurs partout dans le monde (data centers)",
        "Des routeurs qui dirigent le trafic",
        "",
        ("Les données voyagent en paquets", 0),
        "Quand tu envoies un message ou ouvres une page, les données sont découpées en PETITS PAQUETS",
        "Chaque paquet prend le chemin le plus rapide (peuvent prendre des routes différentes !)",
        "Les paquets sont reassemblés à l'arrivée",
        "",
        ("Le protocole IP (Internet Protocol)", 0),
        "Chaque appareil connecté a une adresse IP unique (comme une adresse postale)",
        "Exemple : 192.168.1.1 ou 216.58.214.206 (Google)"
    ])

    make_content_slide(prs, "Le DNS - L'annuaire d'Internet", [
        ("Problème :", 0),
        "Les adresses IP sont difficiles à retenir (172.217.18.14)",
        "On préfère les noms comme google.com ou youtube.com",
        "",
        ("Solution : le DNS (Domain Name System)", 0),
        "C'est l'annuaire téléphonique d'Internet !",
        "Il traduit les noms de domaine en adresses IP",
        "",
        ("Exemple : quand tu tapes youtube.com", 0),
        ("1. Ton navigateur demande au DNS : \"C'est quoi l'IP de youtube.com ?\"", 1),
        ("2. Le DNS répond : \"C'est 216.58.214.46\"", 1),
        ("3. Ton navigateur se connecte à cette IP", 1),
        ("4. Le serveur de YouTube envoie la page web", 1),
        "",
        "⚡ Tout ça en quelques millisecondes !"
    ])

    make_content_slide(prs, "HTTP et HTTPS - Le langage du Web", [
        ("HTTP = HyperText Transfer Protocol", 0),
        "C'est le protocole de communication entre le navigateur et le serveur",
        "",
        ("Quand tu visit un site :", 0),
        ("Ton navigateur envoie une requête HTTP : \"Donne-moi la page d'accueil\"", 1),
        ("Le serveur répond avec la page HTML + les images/CSS/JS", 1),
        ("Ton navigateur affiche la page", 1),
        "",
        ("HTTPS = HTTP + Sécurité (SSL/TLS)", 0),
        "Les données sont chiffrées (cryptées) entre toi et le serveur",
        "Personne ne peut espionner ce que tu fais",
        "🔒 Reconnaissable au cadenas vert dans la barre d'adresse",
        "",
        "💡 HTTP utilise le port 80, HTTPS utilise le port 443"
    ])

    make_two_col_slide(prs, "Les méthodes HTTP",
        "Requêtes du navigateur",
        [
            "GET : Demander une page ou un fichier",
            "POST : Envoyer des données (formulaire, login)",
            "PUT : Mettre à jour une ressource",
            "DELETE : Supprimer une ressource",
        ],
        "Codes de réponse du serveur",
        [
            "200 OK : Tout va bien !",
            "301/302 : Redirection (la page a bougé)",
            "403 : Interdit (tu n'as pas le droit)",
            "404 : Page non trouvée (Page Not Found)",
            "500 : Erreur interne du serveur (le site bug)",
        ]
    )

    make_content_slide(prs, "Les navigateurs web", [
        ("Le rôle du navigateur :", 0),
        "Interpréter le code (HTML, CSS, JavaScript) pour afficher une page",
        "Gérer les onglets, l'historique, les favoris",
        "Assurer la sécurité (bloquer les sites dangereux)",
        "",
        ("Les principaux navigateurs :", 0),
        ("Google Chrome : le plus utilisé (67% de parts de marché)", 1),
        ("Safari : le navigateur d'Apple (sur Mac, iPhone)", 1),
        ("Firefox : navigateur open source (respecte ta vie privée)", 1),
        ("Edge : le navigateur de Microsoft (basé sur Chrome)", 1),
        ("Brave, Opera, Vivaldi : alternatives", 1),
        "",
        ("🔧 Outils développeur (F12)", 0),
        "Tous les navigateurs modernes ont des outils pour développeurs",
        "Tu peux inspecter le code d'une page, voir les requêtes réseau, déboguer du JavaScript..."
    ])

    make_content_slide(prs, "URL - Une adresse web décortiquée", [
        ("https://www.youtube.com/watch?v=dQw4w9WgXcQ", 0),
        "",
        ("Protocole : https", 1),
        ("Sous-domaine : www", 1),
        ("Nom de domaine : youtube", 1),
        ("Domaine de premier niveau (TLD) : .com", 1),
        ("Chemin : /watch", 1),    
        ("Paramètre : ?v=dQw4w9WgXcQ", 1),
        "",
        ("Autres TLD courants :", 0),
        ".fr (France), .org (organisations), .net (réseau), .gov (gouvernement), .io (tech)",
        ".com (commercial, le plus répandu), .edu (éducation)"
    ])

    make_summary_slide(prs, [
        "Internet = infrastructure réseau, Web = service qui utilise Internet",
        "Les données voyagent en paquets via des câbles et routeurs",
        "Le DNS traduit noms de domaine en adresses IP (l'annuaire)",
        "HTTP/HTTPS est le protocole du Web (HTTPS = sécurisé)",
        "Le navigateur interprète le code pour afficher les pages",
        "Une URL se décompose en : protocole, domaine, chemin, paramètres",
        "La programmation web permet de créer des sites et applications"
    ])

    save_presentation(prs, "04_Internet_et_Web.pptx")
    return prs


def p05_programmation_algorithmes():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Qu'est-ce que la programmation ?",
        "Algorithmes, langages, compilateurs et interpréteurs",
        "Partie 1 - Culture Informatique")

    make_objectives_slide(prs, [
        "Comprendre ce qu'est un algorithme",
        "Connaître la différence entre compilateur et interpréteur",
        "Savoir ce qui distingue les langages de programmation",
        "Comprendre le processus de création d'un programme"
    ])

    make_definition_slide(prs, "Algorithme",
        "Un algorithme est une suite d'instructions étape par étape pour résoudre un problème. C'est la recette, le plan, avant d'écrire le code.",
        "Exemple : Algorithme pour faire un café\n1. Prendre une tasse\n2. Mettre le café dans la machine\n3. Ajouter de l'eau\n4. Allumer la machine\n5. Attendre que ça coule\n6. Servir")

    make_content_slide(prs, "Algorithme du quotidien", [
        ("Tu utilises des algorithmes tous les jours sans le savoir !", 0),
        "",
        ("Se brosser les dents :", 1),
        "1. Prendre la brosse à dents",
        "2. Mettre du dentifrice sur la brosse",
        "3. Ouvrir le robinet",
        "4. Mouiller la brosse",
        "5. Fermer le robinet",
        "6. Brosser les dents pendant 2 minutes",
        "7. Cracher",
        "8. Rincer la brosse",
        "",
        ("C'est exactement ça, un algorithme !", 0),
        "Des étapes précises, dans le bon ordre, pour arriver à un résultat."
    ])

    make_content_slide(prs, "Les concepts clés d'un algorithme", [
        ("1. SÉQUENCE", 0),
        "Les instructions s'exécutent dans l'ordre, une par une",
        "Exemple : Faire A, puis B, puis C",
        "",
        ("2. SÉLECTION (Condition)", 0),
        "Faire quelque chose seulement SI une condition est vraie",
        "Exemple : SI il pleut ALORS prendre un parapluie",
        "",
        ("3. ITÉRATION (Boucle)", 0),
        "Répéter une action plusieurs fois",
        "Exemple : POUR chaque dent, brosser 10 fois",
        "",
        ("Ces 3 concepts sont UNIVERSELS", 0),
        "Tous les langages de programmation les utilisent !"
    ])

    make_code_slide(prs, "Algorithme → Code Python",
        ["# Algorithme : Deviner un nombre",
         "# 1. Choisir un nombre secret",
         "nombre_secret = 42",
         "",
         "# 2. Demander un nombre à l'utilisateur",
         "proposition = int(input(\"Devine le nombre : \"))",
         "",
         "# 3. Vérifier la proposition",
         "if proposition == nombre_secret:",
         "    print(\"Bravo ! Tu as deviné !\")",
         "elif proposition < nombre_secret:",
         "    print(\"Trop petit !\")",
         "else:",
         "    print(\"Trop grand !\")"],
        "Chaque ligne de code correspond à une étape de l'algorithme. On voit les 3 concepts : séquence, condition (if/elif/else), et bientôt les boucles !")

    make_content_slide(prs, "Langages de programmation", [
        ("Pourquoi tant de langages ?", 0),
        "Chaque langage est conçu pour un type de tâche",
        "Certains sont plus faciles, d'autres plus puissants ou plus rapides",
        "",
        ("Les grandes familles :", 0),
        ("Python : facile à lire, idéal pour débuter, data science, IA", 1),
        ("JavaScript : le langage du web (sites internet)", 1),
        ("Java : applications d'entreprise, Android", 1),
        ("C/C++ : jeu vidéo, systèmes, performances", 1),
        ("Rust : sécurité, systèmes modernes", 1),
        ("Go : serveurs, outils réseaux", 1),
        ("Kotlin : Android moderne", 1),
        ("Swift : applications Apple (iOS, Mac)", 1),
    ])

    make_content_slide(prs, "Compilateur vs Interpréteur", [
        ("Problème : Le processeur ne comprend que le binaire (0 et 1)", 0),
        "On ne va pas écrire en 0 et 1 ! Il faut traduire notre code en langage machine.",
        "",
        ("Deux approches :", 0),
        "",
        ("🔨 COMPILATEUR (C, C++, Rust, Go)", 0),
        ("Traduit TOUT le code en une fois en fichier exécutable (.exe)", 1),
        ("Avantage : programme très rapide, prêt à l'emploi", 1),
        ("Inconvénient : si tu changes le code, tu dois recompiler", 1),
        "",
        ("⚡ INTERPRÉTEUR (Python, JavaScript, Ruby, PHP)", 0),
        ("Traduit et exécute le code ligne par ligne", 1),
        ("Avantage : facile de tester, pas besoin de recompiler", 1),
        ("Inconvénient : un peu plus lent que du code compilé", 1),
        "",
        "💡 Python est un langage interprété → on peut tester immédiatement !"
    ])

    make_content_slide(prs, "Le cycle de vie d'un programme", [
        ("Du problème à la solution :", 0),
        "",
        ("1️⃣ ANALYSER le problème", 1),
        ("  Comprendre ce qu'on doit faire, poser des questions", 1),
        "",
        ("2️⃣ CONCEVOIR un algorithme", 1),
        ("  Trouver la solution étape par étape", 1),
        "",
        ("3️⃣ ÉCRIRE le code", 1),
        ("  Traduire l'algorithme dans un langage (Python)", 1),
        "",
        ("4️⃣ TESTER le programme", 1),
        ("  Vérifier que ça marche, chercher les bugs", 1),
        "",
        ("5️⃣ CORRIGER et AMÉLIORER", 1),
        ("  Répéter les étapes 3-4 jusqu'à ce que ça marche !", 1),
        "",
        ("🔄 C'est normal de faire plusieurs allers-retours entre 3, 4 et 5 !")
    ])

    make_content_slide(prs, "Pourquoi apprendre Python en premier ?", [
        ("🔷 PYTHON EST PARFAIT POUR DÉBUTER :", 0),
        "",
        ("✅   Syntaxe simple et lisible", 1),
        ("   On dirait presque de l'anglais ! Pas de accolades partout", 1),
        "",
        ("✅   Pas de compilation", 1),
        ("   On écrit et on exécute directement, idéal pour tester", 1),
        "",
        ("✅   Très polyvalent", 1),
        ("   Sites web, jeux, data science, IA, applications...", 1),
        "",
        ("✅   Grande communauté", 1),
        ("   Des millions de développeurs, des tonnes de tutoriels et librairies", 1),
        "",
        ("✅   Très demandé en entreprise", 1),
        ("   Un des langages les mieux payés et les plus utilisés !", 1),
    ])

    make_summary_slide(prs, [
        "Un algorithme est une suite d'étapes pour résoudre un problème",
        "3 concepts clés : séquence, sélection (condition), itération (boucle)",
        "Un langage de programmation permet d'écrire des algorithmes",
        "Compilateur = traduction en une fois, Interpréteur = ligne par ligne",
        "Python est idéal pour débuter : simple, lisible, polyvalent",
        "Programmer = problème → algorithme → code → test → correction"
    ])

    save_presentation(prs, "05_Programmation_algorithmes.pptx")
    return prs


def p06_installation_python_vscode():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Installation de Python et VS Code",
        "Préparer ton environnement de développement pas à pas",
        "Partie 1 - Culture Informatique")

    make_objectives_slide(prs, [
        "Installer Python sur ton ordinateur",
        "Installer VS Code (l'éditeur de code)",  
        "Configurer VS Code pour Python",
        "Vérifier que tout fonctionne correctement"
    ])

    make_definition_slide(prs, "Environnement de développement",
        "C'est l'ensemble des outils que tu utilises pour écrire, tester et exécuter du code. Un bon environnement rend le développement plus facile et efficace.",
        "Notre environnement :\n🐍 Python (le langage)\n📝 VS Code (l'éditeur)\n⚡ Le terminal (pour exécuter les commandes)\n🔧 Les extensions VS Code (pour les fonctionnalités supplémentaires)")

    make_content_slide(prs, "Télécharger et installer Python", [
        ("Étape 1 : Aller sur python.org", 0),
        "Ouvre ton navigateur et va sur https://www.python.org/downloads/",
        "Clique sur le bouton \"Download Python 3.xx\" (la dernière version)",
        "",
        ("Étape 2 : Installer Python", 0),
        ("Lance le fichier téléchargé", 1),
        ("IMPORTANT : Coche \"Add Python to PATH\" (sinon ça ne marchera pas !)", 1),
        ("Clique sur \"Install Now\" (ou \"Customize\" si tu veux changer le dossier)", 1),
        ("Attends la fin de l'installation", 1),
        "",
        ("Étape 3 : Vérifier l'installation", 0),
        "Ouvre un terminal (CMD sur Windows, Terminal sur Mac/Linux)",
        "Tape : python --version  (ou python3 --version)",
        "Tu devrais voir : Python 3.xx.xx"
    ])

    make_content_slide(prs, "Télécharger et installer VS Code", [
        ("Étape 1 : Aller sur code.visualstudio.com", 0),
        "Va sur https://code.visualstudio.com/",
        "Clique sur \"Download for ...\" (selon ton système)",
        "",
        ("Étape 2 : Installer VS Code", 0),
        ("Lance l'installateur", 1),
        ("Accepte les conditions", 1),
        ("Coche les options : \"Add to PATH\", \"Open with Code\" (menu contextuel)", 1),
        ("Clique sur \"Install\" puis \"Finish\"", 1),
        "",
        ("Étape 3 : Découvrir VS Code", 0),
        "Ouvre VS Code → tu vois l'écran d'accueil",
        "Barre latérale gauche : fichiers, recherche, extensions, git",
        "Zone centrale : ton éditeur de code",
        "Barre du bas : terminal, erreurs, ligne actuelle"
    ])

    make_content_slide(prs, "Les extensions Python VS Code", [
        ("Pourquoi des extensions ?", 0),
        "VS Code est un éditeur généraliste. Les extensions ajoutent des fonctionnalités pour des langages spécifiques.",
        "",
        ("Extensions à installer :", 0),
        "",
        ("1️⃣ Python (par Microsoft)", 1),
        ("  C'est l'extension OFFICIELLE pour Python", 1),
        ("  Donne : coloration syntaxique, IntelliSense, débogueur, tests", 1),
        ("  Comment : clic sur l'icône Extensions (🧩) dans la barre de gauche", 1),
        ("  → Cherche \"Python\" → clique sur \"Install\"", 1),
        "",
        ("2️⃣ Pylance", 1),
        ("  Améliore l'autocomplétion et l'analyse de code", 1),
        ("  Installé automatiquement avec l'extension Python", 1),
        "",
        ("3️⃣ Python Debugger", 1),
        ("  Pour déboguer ton code (trouver les erreurs)", 1),
    ])

    make_content_slide(prs, "Configurer VS Code pour Python", [
        ("Sélectionner l'interpréteur Python", 0),
        ("Ouvre VS Code", 1),
        ("Appuie sur Ctrl+Shift+P (ou Cmd+Shift+P sur Mac)", 1),
        ("Tape \"Python: Select Interpreter\"", 1),
        ("Choisis l'interpréteur Python que tu viens d'installer", 1),
        "",
        ("Créer ton premier fichier Python", 0),
        ("Fais Fichier → Nouveau fichier (Ctrl+N)", 1),
        ("Sauvegarde (Ctrl+S) sous le nom : test.py", 1),
        ("VS Code reconnaît que c'est du Python et active les fonctionnalités", 1),
        "",
        ("Le thème (pour faire joli)", 0),
        ("Fais Ctrl+K Ctrl+T (ou va dans Fichier → Préférences → Thème)", 1),
        ("Choisis \"Dark Modern\" (thème sombre) → plus confortable pour les yeux", 1),
    ])

    make_content_slide(prs, "Ton premier vrai programme", [
        ("Maintenant que tout est installé, testons !", 0),
        "",
        ("Dans VS Code, tape ce code :", 0),
        "",
        "print(\"Bonjour le monde !\")",
        "print(\"J'apprends Python !\")",
        "",
        ("Pour exécuter :", 0),
        ("Clique sur le triangle ▶ (Run, en haut à droite)", 1),
        ("OU fais clic droit → \"Run Python\"", 1),
        ("OU utilise Ctrl+F5", 1),
        ("OU tape dans le terminal : python test.py", 1),
        "",
        ("Résultat attendu :", 0),
        "Bonjour le monde !",
        "J'apprends Python !",
        "",
        "🎉 Félicitations ! Tu viens d'écrire et exécuter ton premier programme !"
    ])

    make_summary_slide(prs, [
        "Python se télécharge sur python.org (cocher Add to PATH !)",
        "VS Code est l'éditeur recommandé avec l'extension Python",
        "L'interpréteur Python doit être sélectionné dans VS Code",
        "On écrit du code dans un fichier .py et on l'exécute",
        "L'environnement est prêt → on peut commencer à coder !",
        "N'hésite pas à personnaliser VS Code (thème, polices, etc.)"
    ])

    save_presentation(prs, "06_Installation_Python_VSCode.pptx")
    return prs


def p07_terminal_commandes():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Premier contact avec le Terminal",
        "Les commandes essentielles pour naviguer comme un pro",
        "Partie 1 - Culture Informatique")

    make_objectives_slide(prs, [
        "Comprendre ce qu'est un terminal (shell)",
        "Connaître les commandes de base pour se déplacer",
        "Savoir créer, copier, déplacer et supprimer des fichiers",
        "Se sentir à l'aise dans le terminal"
    ])

    make_definition_slide(prs, "Le Terminal (ou Shell)",
        "Le terminal est une interface en texte où tu tapes des commandes pour interagir avec l'ordinateur. C'est plus puissant que la souris pour beaucoup de tâches.",
        "💡 Terminal = interface en mode texte\n🐚 Shell = le programme qui interprète tes commandes (bash, zsh, PowerShell)\n📟 Console = l'écran où ça s'affiche")

    make_content_slide(prs, "Ouvrir le terminal", [
        ("Sur Windows :", 0),
        "Ouvre VS Code, puis Terminal → Nouveau terminal (ou Ctrl+ù)",
        "Tu peux aussi ouvrir PowerShell/CMD directement",
        "",
        ("Sur macOS :", 0),
        "Terminal (ou iTerm2) dans Applications/Utilitaires",
        "",
        ("Sur Linux :", 0),
        "Ctrl+Alt+T ou chercher \"Terminal\" dans les applications",
        "",
        ("L'invite de commande (prompt) :", 0),
        "Quand tu ouvres le terminal, tu vois :",
        "tonutilisateur@ordinateur:~$",
        "",
        "Le $ (ou > sur Windows) signifie que le terminal t'attend !",
        "~ signifie que tu es dans ton dossier personnel (home)"
    ])

    make_code_slide(prs, "Se déplacer dans les dossiers",
        ["# Afficher le dossier actuel",
         "pwd",
         "  → /home/tom/Documents",
         "",
         "# Lister les fichiers et dossiers",
         "ls",
         "  → Documents  Images  Musique  test.py",
         "",
         "# Se déplacer dans un dossier",
         "cd Documents",
         "",
         "# Revenir au dossier parent",
         "cd ..",
         "",
         "# Aller directement dans un sous-dossier",
         "cd Documents/Projets/Python",
         "",
         "# Revenir au dossier personnel (home)",
         "cd ~   (ou juste cd)"],
        "💡 pwd = print working directory\nls = list\ncd = change directory\nUtilise la touche Tab pour l'autocomplétion !")

    make_code_slide(prs, "Créer et manipuler des fichiers/dossiers",
        ["# Créer un dossier",
         "mkdir projet_python",
         "",
         "# Créer un fichier vide",
         "touch mon_fichier.py",
         "",
         "# Copier un fichier",
         "cp mon_fichier.py copie.py",
         "",
         "# Déplacer / Renommer un fichier",
         "mv mon_fichier.py ../autre_dossier/",
         "mv ancien_nom.py nouveau_nom.py",
         "",
         "# Supprimer un fichier (attention !)",
         "rm fichier_a_supprimer.txt",
         "",
         "# Supprimer un dossier et tout son contenu",
         "rm -rf dossier_a_supprimer/"],
        "⚠️  ATTENTION : rm -rf supprime définitivement, sans passer par la corbeille !")

    make_code_slide(prs, "Lire et écrire dans des fichiers",
        ["# Voir le contenu d'un fichier",
         "cat mon_fichier.py",
         "",
         "# Voir le contenu page par page",
         "less mon_fichier.py",
         "  (appuie sur q pour quitter)",
         "",
         "# Voir les premières lignes (10 par défaut)",
         "head mon_fichier.py",
         "",
         "# Voir les dernières lignes",
         "tail mon_fichier.py",
         "",
         "# Chercher un mot dans un fichier",
         "grep \"mot_recherche\" mon_fichier.py",
         "",
         "# Compter les lignes, mots, caractères",
         "wc mon_fichier.py"],
        "💡 Toutes ces commandes sont essentielles pour un développeur !")

    make_code_slide(prs, "Commandes système utiles",
        ["# Effacer l'écran du terminal",
         "clear",
         "",
         "# Voir la date et l'heure",
         "date",
         "",
         "# Voir qui est connecté",
         "whoami",
         "",
         "# Voir les processus en cours",
         "ps",
         "",
         "# Voir l'espace disque",
         "df -h",
         "",
         "# Voir l'arborescence d'un dossier",
         "ls -R",
         "  (affiche TOUS les sous-dossiers récursivement)",
         "",
         "# Voir les fichiers cachés",
         "ls -a"],
        "💡 Sur Windows, certaines commandes diffèrent. Dans VS Code, le terminal PowerShell accepte dir au lieu de ls.")

    make_content_slide(prs, "Les options des commandes", [
        ("La plupart des commandes acceptent des OPTIONS (flags)", 0),
        "",
        ("Structure : commande -option(s) argument", 0),
        "",
        ("Exemples :", 0),
        "ls → liste simple",
        "ls -l → liste détaillée (avec permissions, taille, date)",
        "ls -a → liste TOUS les fichiers (y compris cachés)",
        "ls -la → combine -l et -a (liste détaillée de tous les fichiers)",
        "ls -lh → liste détaillée avec tailles lisibles (Ko, Mo)",
        "",
        ("Commandes avec plusieurs options :", 0),
        "rm -rf dossier/   (-r = récursif, -f = force)",
        "cp -r dossier/ backup/   (-r = copie récursive pour les dossiers)",
        "",
        "💡 ls --help ou man ls → affiche toutes les options d'une commande"
    ])

    make_content_slide(prs, "Les raccourcis clavier du terminal", [
        ("Raccourcis INDISPENSABLES :", 0),
        "",
        ("↑ / ↓ : parcourir l'historique des commandes", 1),
        ("Tab : autocomplétion (magique ! tape le début et appuie sur Tab)", 1),
        ("Ctrl+C : annuler la commande en cours / arrêter un programme", 1),
        ("Ctrl+D : fermer le terminal", 1),
        ("Ctrl+L : effacer l'écran (comme clear)", 1),
        ("Ctrl+Shift+C : copier (dans le terminal)", 1),
        ("Ctrl+Shift+V : coller (dans le terminal)", 1),
        ("Ctrl+A : aller au début de la ligne", 1),
        ("Ctrl+E : aller à la fin de la ligne", 1),
        ("Ctrl+U : effacer toute la ligne", 1),
        ("Ctrl+W : effacer le mot avant le curseur", 1),
        "",
        "💡 Ces raccourcis te feront gagner BEAUCOUP de temps !"
    ])

    make_summary_slide(prs, [
        "Le terminal permet d'interagir avec l'ordinateur via des commandes textuelles",
        "pwd, ls, cd sont les commandes de base pour se déplacer",
        "mkdir, touch, cp, mv, rm créent/copient/déplacent/suppriment des fichiers",
        "cat, less, head, tail permettent de lire des fichiers",
        "Les options (-l, -a, -r, -f, -h) modifient le comportement des commandes",
        "Tab = autocomplétion, ↑/↓ = historique, Ctrl+C = annuler",
        "Entraîne-toi, le terminal deviendra ton meilleur ami !"
    ])

    save_presentation(prs, "07_Terminal_commandes.pptx")
    return prs


def p08_hello_world():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Hello World !",
        "Ton tout premier programme en Python",
        "Partie 1 - Culture Informatique")

    make_objectives_slide(prs, [
        "Écrire et exécuter son premier programme Python",
        "Comprendre ce qu'est la fonction print()",
        "Découvrir le rôle des commentaires",
        "Apprendre à gérer les erreurs (bugs)"
    ])

    make_definition_slide(prs, "La tradition Hello World",
        "Quand on apprend un nouveau langage, le premier programme qu'on écrit affiche \"Hello, World!\" C'est une tradition depuis les années 1970 !",
        "Ce programme simple permet de vérifier que tout l'environnement fonctionne correctement.")

    make_code_slide(prs, "Ton premier programme",
        ["# Mon premier programme en Python",
         "print(\"Bonjour le monde !\")",
         "print(\"Je commence mon aventure Python !\")"],
        "Exécute ce code dans VS Code (▶ en haut à droite)\n→ Tu verras le texte s'afficher dans le terminal.")

    make_content_slide(prs, "La fonction print() en détail", [
        ("print() est une FONCTION intégrée à Python", 0),
        "Elle permet d'afficher du texte dans le terminal (la console)",
        "",
        ("Syntaxe :", 0),
        "print(\"ce qu'on veut afficher\")",
        "",
        ("Règles :", 0),
        "Le texte doit être entre guillemets \" \" ou apostrophes ' '",
        "Les parenthèses () sont obligatoires (c'est comme ça qu'on appelle une fonction)",
        "",
        ("Exemples :", 0),
        "print(\"Bonjour\")    → Bonjour",
        "print('Salut')       → Salut  (les guillemets simples marchent aussi)",
        "print(42)            → 42  (les nombres sans guillemets)",
        "print(3 + 4)         → 7  (Python calcule avant d'afficher)",
    ])

    make_code_slide(prs, "Afficher plusieurs choses",
        ["# On peut afficher plusieurs valeurs avec une virgule",
         "prenom = \"Lucas\"",
          "age = 25",
          "print(\"Je m'appelle\", prenom, \"et j'ai\", age, \"ans\")",
          "",
          "# Résultat :",
          "# Je m'appelle Utilisateur et j'ai 25 ans",
         "",
         "# On peut aussi utiliser les f-strings (plus tard)",
         "print(f\"Je m'appelle {prenom} et j'ai {age} ans\")"],
        "💡 Avec print(), les éléments séparés par des virgules sont automatiquement espacés.")

    make_code_slide(prs, "Les commentaires - expliquer son code",
        ["# Les commentaires commencent par #",
         "# Python ignore tout ce qui est après # sur la même ligne",
         "# Ça sert à EXPLIQUER ce que fait le code",
         "",
         "print(\"Bonjour\")  # Ceci est aussi un commentaire",
         "",
         "# Les commentaires sont TRÈS importants quand on travaille en équipe",
         "# Et aussi quand on revoit son code 6 mois plus tard !",
         "",
         "# On peut aussi commenter plusieurs lignes",
         "# en mettant un # au début de chaque ligne",
         "",
         "# Pour les longues explications, on utilise les triples guillemets :",
         "\"\"\"",
         "Ceci est un commentaire",
         "sur plusieurs lignes.",
         "Très utile pour la documentation !",
         "\"\"\""],
        "💡 Prends l'habitude de commenter TON code dès le début. Tu me remercieras plus tard !")

    make_content_slide(prs, "Les erreurs : pas de panique !", [
        ("Même les meilleurs programmeurs font des erreurs !", 0),
        "L'important est de savoir les lire et les comprendre.",
        "",
        ("Type d'erreur les plus fréquentes :", 0),
        "",
        ("SyntaxError : faute de frappe dans la syntaxe", 1),
        ("  Exemple : print(\"Bonjour\")  →  parenthèse fermante manquante", 1),
        "",
        ("NameError : nom de variable inconnu", 1),
        ("  Exemple : print(Bonjour)  →  \"Bonjour\" sans guillemets", 1),
        "",
        ("TypeError : mauvais type de donnée", 1),
        ("   Exemple : print(\"Bonjour\" + 5)  →  on ne peut pas additionner texte et nombre", 1),
        "",
        ("IndentationError : problème d'espacement", 1),
        ("  (on verra ça plus tard avec les conditions et boucles)", 1),
        "",
        ("💡 LIS LE MESSAGE D'ERREUR ! Python te dit EXACTEMENT où est le problème !")
    ])

    make_content_slide(prs, "Exemple d'erreur pas à pas", [
        ("Imaginons ce code :", 0),
        "",
        "print(\"Bonjour\"",
        "",
        ("Message d'erreur :", 0),
        '  File \"test.py\", line 1',
        '    print(\"Bonjour\"',
        '         ^',
        "SyntaxError: '(' was never closed",
        "",
        ("Lecture de l'erreur :", 0),
        ("1. Fichier : test.py (c'est dans ton fichier)", 1),
        ("2. Ligne : line 1 (l'erreur est sur la ligne 1)", 1),
        ("3. ^ montre où le problème a été détecté", 1),
        ("4. SyntaxError : '(' was never closed → il manque une parenthèse)", 1),
        "",
        ("Correction :", 0),
        "print(\"Bonjour\")  →  ajouter la parenthèse fermante",
        "",
        "🔑 Lire les erreurs est une compétence qui s'apprend !"
    ])

    make_exercise_slide(prs, "Exercice 1 : Hello World personnalisé",
        "1. Crée un fichier hello.py dans VS Code\n"
        "2. Écris un programme qui affiche :\n"
        "   - \"Bonjour, je m'appelle [TonPrénom]\"\n"
        "   - \"J'ai [TonAge] ans\"\n"
        "   - \"J'apprends Python !\"\n"
        "3. Exécute-le et vérifie que ça marche\n"
        "4. Essaie de faire une erreur exprès pour voir le message\n"
        "5. Corrige l'erreur",
        ["print(\"Bonjour, je m'appelle Lucas\")",
         "print(\"J'ai 14 ans\")",
         "print(\"J'apprends Python !\")"]
    )

    make_exercise_slide(prs, "Exercice 2 : Le compilateur humain",
        "Sans exécuter le code, trouve ce que chaque ligne va afficher :\n\n"
        "a) print(10 + 5)\n"
        "b) print(\"10\" + \"5\")\n"
        "c) print(\"Salut\" * 3)\n"
        "d) print(\"La réponse est\", 42)\n\n"
        "Écris tes réponses sur un papier, puis exécute le code pour vérifier !\n"
        "Astuce : Python fait parfois des choses surprenantes... 🧐",
        ["# Réponses :",
         "a) 15          (10 + 5 = 15, addition normale)",
         "b) 105         (concaténation de chaînes \"10\" + \"5\" = \"105\")",
         "c) SalutSalutSalut  (répétition de chaîne)",
         "d) La réponse est 42  (print ajoute un espace entre les éléments)"]
    )

    make_summary_slide(prs, [
        "print() affiche du texte dans la console",
        "Le texte doit être entre guillemets ou apostrophes",
        "Les nombres sans guillemets sont interprétés comme des valeurs numériques",
        "Les commentaires (#) expliquent le code - Python les ignore",
        "Les erreurs sont normales ! Lis le message, il t'indique où est le problème",
        "Hello World est la première étape de ton voyage de programmeur !"
    ])

    save_presentation(prs, "08_Hello_World.pptx")
    return prs


# ═══════════════════════════════════════════════════════════════
# PARTIE 2 : LES BASES DE PYTHON
# ═══════════════════════════════════════════════════════════════

def p09_variables_types():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Variables et Types de données",
        "Stocker et manipuler des informations en Python",
        "Partie 2 - Les Bases de Python")

    make_objectives_slide(prs, [
        "Comprendre ce qu'est une variable",
        "Connaître les types de base : int, float, str, bool",
        "Savoir nommer correctement une variable",
        "Découvrir la fonction type()"
    ])

    make_definition_slide(prs, "Une variable, c'est quoi ?",
        "Une variable est une boîte qui contient une valeur. Tu donnes un nom à cette boîte pour pouvoir récupérer son contenu plus tard.",
        "age = 25      → la variable 'age' contient la valeur 25\nnom = \"Utilisateur\"  → la variable 'nom' contient le texte \"Utilisateur\"\nC'est comme une étiquette sur une boîte !")

    make_content_slide(prs, "Créer et utiliser une variable", [
        ("En Python, c'est très simple :", 0),
        "",
        "nom_de_la_variable = valeur",
        "",
        ("Le signe = est l'opérateur d'affectation (pas \"égal\" mathématique)", 0),
        "Il prend la valeur à droite et la stocke dans la variable à gauche",
        "",
        ("Exemples :", 0),
        "age = 25                 # age vaut 25",
        "prenom = \"Lucas\"         # prenom vaut \"Lucas\"",
        "prix = 19.99             # prix vaut 19.99",
        "est_etudiant = True      # est_etudiant vaut True",
        "",
        ("Ensuite, on peut utiliser la variable :", 0),
        "print(age)            # affiche 25",
        "print(\"Je m'appelle\", prenom)  # affiche \"Je m'appelle Utilisateur\"",
        "",
        "💡 Python déduit tout seul le type de la variable selon la valeur !"
    ])

    make_content_slide(prs, "Les types de base en Python", [
        ("int (entier) : nombres sans virgule", 0),
        "age = 14 / quantite = 100 / temperature = -5",
        "",
        ("float (décimal) : nombres à virgule", 0),
        "prix = 19.99 / pi = 3.14159 / note = 14.5",
        "(en anglais, le point = la virgule !)",
        "",
        ("str (string / chaîne) : texte", 0),
        "nom = \"Lucas\" / phrase = 'Bonjour !'",
        "(guillemets doubles ou simples, au choix mais cohérent)",
        "",
        ("bool (booléen) : vrai/faux", 0),
        "est_connecte = True / est_adulte = False",
        "(Toujours True ou False, avec majuscule !)",
    ])

    make_code_slide(prs, "Voir le type d'une variable",
        ["# On utilise la fonction type()",
         "",
         "age = 14",
         "prix = 19.99",
         "nom = \"Lucas\"",
         "est_ok = True",
         "",
         "print(type(age))    # <class 'int'>",
         "print(type(prix))   # <class 'float'>",
         "print(type(nom))    # <class 'str'>",
         "print(type(est_ok)) # <class 'bool'>",
         "",
         "# C'est super utile pour déboguer !",
         "valeur = input(\"Tape quelque chose : \")",
         "print(type(valeur))  # <class 'str'> (TOUJOURS une chaîne avec input !)"],
        "💡 type() est très utile pour comprendre pourquoi ton code ne marche pas comme prévu !")

    make_content_slide(prs, "Les règles pour nommer une variable", [
        ("Règles OBLIGATOIRES :", 0),
        "",
        "✅  Commencer par une lettre ou _ (pas par un chiffre)",
        "✅  Contenir seulement : lettres (a-z, A-Z), chiffres, _",
        "✅  Sensible à la casse : age, Age, AGE sont 3 variables différentes",
        "",
        ("⚠️  Mots réservés (ne pas utiliser) :", 0),
        "False, True, if, else, for, while, import, class, def, return...",
        "",
        ("Conventions RECOMMANDÉES (PEP 8) :", 0),
        "🐍  Utiliser le snake_case : age_utilisateur, mon_nom, est_connecte",
        "🐍  Noms en minuscules avec des underscores",
        "🐍  Utiliser des noms explicites (pas a, b, x mais age, prix_total)",
        "🐍  Éviter les noms trop longs (mais assez longs pour être clairs)",
    ])

    make_code_slide(prs, "Bons vs mauvais noms de variables",
        ["# ❌ MAUVAIS noms",
         "x = 14",
         "a = \"Lucas\"",
         "d = True",
         "t = \"Paris\"",
         "var1 = 42",
         "",
         "# ✅ BONS noms",
         "age = 14",
         "prenom = \"Lucas\"",
         "est_inscrit = True",
         "ville_residence = \"Paris\"",
         "nombre_mystere = 42",
         "",
         "# Le code doit être lisible comme un livre !",
         "# Imagine si tu devais comprendre ce code dans 1 an..."],
        "💡 Un bon code est un code qui se lit comme une phrase !")

    make_content_slide(prs, "Réassigner une variable", [
        ("Une variable peut changer de valeur :", 0),
        "",
        "score = 0",
        "print(score)  # 0",
        "",
        "score = 10",
        "print(score)  # 10",
        "",
        "score = score + 5   # prend la valeur actuelle, ajoute 5, stocke le résultat",
        "print(score)  # 15",
        "",
        "score += 5  # raccourci pour score = score + 5",
        "print(score)  # 20",
        "",
        ("Autres raccourcis :", 0),
        "score -= 3   # score = score - 3",
        "score *= 2   # score = score * 2",
        "score /= 4   # score = score / 4",
        "",
        "💡 La partie droite est calculée EN PREMIER, puis stockée dans la variable de gauche"
    ])

    make_content_slide(prs, "Variables multiples et échange", [
        ("On peut assigner plusieurs variables en une ligne :", 0),
        "x, y, z = 10, 20, 30",
        "print(x, y, z)  # 10 20 30",
        "",
        ("Même valeur à plusieurs variables :", 0),
        "a = b = c = 0",
        "print(a, b, c)  # 0 0 0",
        "",
        ("Échanger deux variables (spécifique à Python !) :", 0),
        "a = 5",
        "b = 10",
        "print(a, b)  # 5 10",
        "",
        "a, b = b, a  # échange magique !",
        "print(a, b)  # 10 5",
        "",
        ("Dans d'autres langages, il faudrait une variable temporaire :", 0),
        "temp = a",
        "a = b",
        "b = temp",
    ])

    make_content_slide(prs, "Le typage dynamique de Python", [
        ("Python est à TYPAGE DYNAMIQUE", 0),
        "Une même variable peut changer de type en cours de route !",
        "",
        ("Exemple :", 0),
        "ma_variable = 42           # int",
        "print(type(ma_variable))  # <class 'int'>",
        "",
        "ma_variable = \"Bonjour\"   # maintenant c'est str !",
        "print(type(ma_variable))  # <class 'str'>",
        "",
        "ma_variable = 3.14        # maintenant c'est float !",
        "print(type(ma_variable))  # <class 'float'>",
        "",
        ("⚠️  À ÉVITER dans la pratique !", 0),
        "C'est source de confusion et de bugs.",
        "Une variable devrait garder le même type tout au long du code.",
    ])

    make_exercise_slide(prs, "Exercice : Créer son profil",
        "Crée un programme qui :\n"
        "1. Stocke ton prénom dans une variable\n"
        "2. Stocke ton âge dans une variable\n"
        "3. Stocke ta ville dans une variable\n"
        "4. Stocke si tu es étudiant (True/False)\n"
        "5. Affiche : \"Je m'appelle [prénom], j'ai [âge] ans, j'habite à [ville]\"\n"
        "6. Affiche le type de chaque variable\n\n"
        "BONUS : Modifie ton âge et réaffiche la phrase",
        ["prenom = \"Lucas\"",
         "age = 14",
         "ville = \"Paris\"",
         "est_etudiant = True",
         "print(f\"Je m'appelle {prenom}, j'ai {age} ans, j'habite à {ville}\")",
         "print(type(prenom))",
         "print(type(age))",
         "print(type(ville))",
         "print(type(est_etudiant))"]
    )

    make_summary_slide(prs, [
        "Une variable stocke une valeur sous un nom",
        "Types de base : int (entier), float (décimal), str (texte), bool (vrai/faux)",
        "type() permet de connaître le type d'une variable",
        "Noms de variables : snake_case, explicites, sans mots réservés",
        "Python est typé dynamiquement (le type peut changer)",
        "On peut échanger des variables facilement : a, b = b, a"
    ])

    save_presentation(prs, "09_Variables_et_types.pptx")
    return prs


def p10_strings():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    make_title_slide(prs,
        "Les Chaînes de caractères (Strings)",
        "Tout sur le texte en Python",
        "Partie 2 - Les Bases de Python")

    make_objectives_slide(prs, [
        "Comprendre comment créer et manipuler des chaînes",
        "Maîtriser la concaténation et la répétition",
        "Utiliser les f-strings pour formater du texte",
        "Découvrir les méthodes principales des chaînes"
    ])

    make_definition_slide(prs, "String (chaîne de caractères)",
        "Une chaîne de caractères (str) est une séquence de caractères (lettres, chiffres, symboles) encadrée par des guillemets ou des apostrophes.",
        "message = \"Bonjour !\"    → guillemets doubles\nmessage = 'Bonjour !'     → guillemets simples\nmessage = \"J'ai un ordi\"   → guillemets doubles (contient une apostrophe)\nmessage = 'Il a dit \"Bonjour\"' → guillemets simples (contient des guillemets)")

    make_code_slide(prs, "Créer des chaînes",
        ["# Différentes façons de créer des chaînes",
         "",
         "str1 = \"Bonjour\"        # guillemets doubles",
         "str2 = 'Salut'           # guillemets simples",
         "str3 = \"\"\"Long texte",
         "sur plusieurs",
         "lignes.\"\"\"             # triples guillemets",
         "",
         "# Guillemets à l'intérieur :",
         "s = \"J'ai un chien\"     # OK",  
         "s = 'Il a dit \"oui\"'   # OK",
         "",
         "# Caractères d'échappement :",
         "s = \"Ligne1\\nLigne2\"    # \\n = saut de ligne",
         "s = \"Tabulation\\there\"  # \\t = tabulation",
         "s = \"C:\\\\Users\\\\Nom\" # \\\\ = un vrai backslash"],
        "💡 Les triples guillemets sont parfaits pour les longs textes et la documentation !")

    make_code_slide(prs, "Concaténation et répétition",
        ["# CONCATÉNATION : coller des chaînes avec +",
         "prenom = \"Lucas\"",
         "message = \"Bonjour \" + prenom",
         "print(message)  # Bonjour Lucas",
         "",
         "# Attention aux espaces !",
         "a = \"Bonjour\"",
         "b = \"Lucas\"",
         "c = a + b",
         "print(c)  # BonjourLucas  (pas d'espace automatique)",
         "",
         "# Solution : ajouter un espace",
         "c = a + \" \" + b",
         "print(c)  # Bonjour Lucas",
         "",
         "# RÉPÉTITION : multiplier une chaîne avec *",
         "ligne = \"-\" * 20",
         "print(ligne)  # --------------------",
         "",
         "print(\"Ha \" * 3 + \"! \")  # Ha Ha Ha !"],
        "💡 La répétition avec * est très pratique pour créer des séparateurs, barres de progression, etc.")

    make_code_slide(prs, "Les f-strings (formatage moderne)",
        ["# f-strings : la façon RECOMMANDÉE de formater du texte en Python",
         "# (disponible depuis Python 3.6)",
         "",
         "prenom = \"Lucas\"",
         "age = 14",
         "ville = \"Paris\"",
         "",
         "# Sans f-strings (avant 3.6) :",
         "print(\"Je m'appelle \" + prenom + \" et j'ai \" + str(age) + \" ans\")",
         "  # Pas très lisible...",
         "",
         "# Avec f-strings (MAGIQUE) :",
         "print(f\"Je m'appelle {prenom} et j'ai {age} ans\")",
         "  # Les {variables} sont remplacées automatiquement !",
         "",
         "# On peut même mettre des expressions :",
         "print(f\"Dans 5 ans, j'aurai {age + 5} ans\")",
         "  # Dans 5 ans, j'aurai 19 ans"],
        "💡 Les f-strings rendent le code BEAUCOUP plus lisible. On les utilisera TOUT le temps !")

    make_code_slide(prs, "Méthodes utiles sur les chaînes (1/2)",
        ["# Les chaînes ont plein de méthodes intégrées !",
         "texte = \"  Bonjour tout le Monde !  \"",
         "",
         "# Mise en forme :",
         "print(texte.upper())      # BONJOUR TOUT LE MONDE !",
         "print(texte.lower())      # bonjour tout le monde !",
         "print(texte.title())      # Bonjour Tout Le Monde !",
         "print(texte.capitalize()) # Bonjour tout le monde !",
         "print(texte.swapcase())   # bONJOUR TOUT LE mONDE !",
         "",
         "# Supprimer les espaces :",
         "print(texte.strip())      # \"Bonjour tout le Monde !\"",
         "print(texte.lstrip())     # \"Bonjour tout le Monde !  \"",
         "print(texte.rstrip())     # \"  Bonjour tout le Monde !\""],
        "💡 Les méthodes ne modifient PAS la chaîne originale (les chaînes sont immutables), elles retournent une nouvelle chaîne.")

    make_code_slide(prs, "Méthodes utiles sur les chaînes (2/2)",
        ["texte = \"Bonjour tout le Monde\"",
         "",
         "# Recherche et remplacement :",
         "print(len(texte))            # 22 (longueur, fonction intégrée)",
         "print(texte.count(\"o\"))      # 4 (compte les occurrences)",
         "print(texte.find(\"Monde\"))   # 15 (position de départ)",
         "print(texte.replace(\"Monde\", \"Python\"))  # Bonjour tout le Python",
         "",
         "# Vérifications :",
         "print(texte.startswith(\"Bon\"))  # True",
         "print(texte.endswith(\"!\"))      # False",
         "print(\"abc\".isalpha())          # True (que des lettres)",
         "print(\"123\".isdigit())          # True (que des chiffres)",
         "print(\"abc123\".isalnum())      # True (lettres et/ou chiffres)",
         "",
         "# Découpage et jointure :",
         "mots = texte.split(\" \")",     
         "print(mots)  # ['Bonjour', 'tout', 'le', 'Monde']",
         "print(\"-\".join(mots))  # Bonjour-tout-le-Monde"],
        "💡 split() et join() sont très utilisés pour traiter du texte !")

    make_content_slide(prs, "Indexation et slicing (tranches)", [
        ("Chaque caractère a une position (index) :", 0),
        "",
        "texte = \"PYTHON\"",
        "        012345   (index 0 → 5)",
        "",
        "tex     te[0] → 'P'      # premier caractère",
        "texte[1] → 'Y'",
        "texte[-1] → 'N'      # dernier caractère (-1 = dernier)",
        "texte[-2] → 'O'      # avant-dernier",
        "",
        ("Slicing (tranches) :", 0),
        "texte[0:3]  → 'PYT'   # de l'index 0 à 2 (3 exclu)",
        "texte[:3]   → 'PYT'   # depuis le début jusqu'à 3",
        "texte[2:]   → 'THON'  # de l'index 2 jusqu'à la fin",
        "texte[::2]  → 'PTO'   # un caractère sur 2",
        "texte[::-1] → 'NOHTYP'  # inverse la chaîne !",
        "",
        "💡 Le slicing est une fonctionnalité PUISSANTE de Python !"
    ])

    make_content_slide(prs, "Chaînes immutables",
        [("Les chaînes sont IMMUTABLES", 0),
        "Une fois créée, une chaîne ne peut PAS être modifiée.",
        "Les méthodes retournent NOUVELLE chaîne, elles ne modifient pas l'originale.",
        "",
        ("Exemple :", 0),
        "texte = \"Bonjour\"",
        "texte[0] = \"b\"  # ❌ ERREUR ! TypeError",
        "",
        ("Pour \"modifier\" une chaîne :", 0),
        "texte = \"b\" + texte[1:]  # On crée une NOUVELLE chaîne",
        "print(texte)  # \"bonjour\"",
        "",
        ("Autre possibilité : convertir en liste", 0),
        "liste = list(texte)    # ['b', 'o', 'n', 'j', 'o', 'u', 'r']",
        "liste[0] = 'B'",
        "texte = \"\".join(liste)  # \"Bonjour\"",
        "",
        "💡 L'immutabilité permet à Python d'optimiser la mémoire et la vitesse"
    ])

    make_exercise_slide(prs, "Exercice : Manipulation de texte",
        "1. Crée une variable phrase = \"j'apprends la programmation python\"\n"
        "2. Affiche la phrase en majuscules\n"
        "3. Affiche la longueur de la phrase\n"
        "4. Remplace \"python\" par \"Python\"\n"
        "5. Découpe la phrase en mots et affiche chaque mot sur une ligne\n"
        "6. Affiche le 3ème caractère\n"
        "7. Affiche les 5 derniers caractères\n"
        "8. Affiche la phrase inversée\n\n"
        "BONUS : Demande le prénom de l'utilisateur et affiche-le à l'envers",
        ["phrase = \"j'apprends la programmation python\"",
         "print(phrase.upper())",
         "print(len(phrase))",
         "print(phrase.replace(\"python\", \"Python\"))",
         "for mot in phrase.split():",
         "    print(mot)",
         "print(phrase[2])",
         "print(phrase[-5:])",
         "print(phrase[::-1])"]
    )

    make_summary_slide(prs, [
        "Les chaînes sont des séquences de caractères entre \" \" ou ' '",
        "On peut concaténer (+) et répéter (*) les chaînes",
        "Les f-strings (f\"...{variable}...\") sont le meilleur moyen de formater",
        "split() découpe en liste, join() rassemble en chaîne",
        "Les chaînes sont immutables (on ne peut pas les modifier, on en crée de nouvelles)",
        "Le slicing [début:fin:pas] permet d'extraire des parties de chaîne"
    ])

    save_presentation(prs, "10_Strings.pptx")
    return prs


# ─── Génération ─────────────────────────────────────────────────

if __name__ == "__main__":
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    
    presentations = [
        ("Partie 1", p01_histoire_informatique),
        ("Partie 1", p02_fonctionnement_ordinateur),
        ("Partie 1", p03_systemes_exploitation),
        ("Partie 1", p04_internet_web),
        ("Partie 1", p05_programmation_algorithmes),
        ("Partie 1", p06_installation_python_vscode),
        ("Partie 1", p07_terminal_commandes),
        ("Partie 1", p08_hello_world),
        ("Partie 2", p09_variables_types),
        ("Partie 2", p10_strings),
    ]
    
    for part, func in presentations:
        print(f"\n📊 Génération : {func.__name__}")
        func()
    
    print(f"\n{'='*50}")
    print(f"✅ {len(presentations)} présentations générées !")
    print(f"📁 Dossier : {OUTPUT_DIR}")
    print(f"{'='*50}")
