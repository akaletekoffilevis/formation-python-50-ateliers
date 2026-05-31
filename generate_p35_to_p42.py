#!/usr/bin/env python3
"""Generate presentations 35 to 42 - Part 4."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

class Th:
    D=RGBColor(0x0A,0x0A,0x2E); P=RGBColor(0x16,0x21,0x3E)
    S=RGBColor(0x1A,0x1A,0x4E); O=RGBColor(0xFF,0x6B,0x35)
    C=RGBColor(0x00,0xB4,0xD8); PU=RGBColor(0x7B,0x2F,0xF7)
    G=RGBColor(0x06,0xD6,0xA0); W=RGBColor(0xFF,0xFF,0xFF)
    L=RGBColor(0xF8,0xF9,0xFA); GR=RGBColor(0x6C,0x75,0x7D)
    LG=RGBColor(0xDE,0xE2,0xE6); CB=RGBColor(0x1E,0x1E,0x2E)
    CT=RGBColor(0xCD,0xD6,0xF4); R=RGBColor(0xEF,0x44,0x44)

OUT = "/home/akaletekoffilevis/Bureau/Coach/presentations"
os.makedirs(OUT, exist_ok=True)

def bg(sl,c):
    f=sl.background.fill; f.solid(); f.fore_color.rgb=c

def sh(sl,l,t,w,h,fc,bc=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=fc
    s.line.fill.background()
    if bc: s.line.color.rgb=bc; s.line.width=Pt(1)
    return s

def rr(sl,l,t,w,h,fc):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=fc; s.line.fill.background()
    return s

def tb(sl,l,t,w,h,text,fs=14,c=None,b=False,a=PP_ALIGN.LEFT):
    tx=sl.shapes.add_textbox(l,t,w,h); tf=tx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(fs); p.font.color.rgb=c or Th.W
    p.font.bold=b; p.font.name='Calibri'; p.alignment=a
    return tx

def tb_rich(sl,l,t,w,h):
    tx=sl.shapes.add_textbox(l,t,w,h); tf=tx.text_frame; tf.word_wrap=True; return tx

def ap(tf,text,sz=14,c=None,b=False,al=PP_ALIGN.LEFT,sb=0,sa=0,name='Calibri'):
    if len(tf.paragraphs)==1 and tf.paragraphs[0].text=='': p=tf.paragraphs[0]
    else: p=tf.add_paragraph()
    p.text=text; p.font.size=Pt(max(sz,8)); p.font.color.rgb=c or Th.W; p.font.bold=b
    p.font.name=name; p.alignment=al; p.space_before=Pt(sb); p.space_after=Pt(sa)
    return p

def code(sl,l,t,w,h,lines,fs=11):
    rr(sl,l,t,w,h,Th.CB)
    tx=tb_rich(sl,l+Inches(0.3),t+Inches(0.15),w-Inches(0.6),h-Inches(0.3))
    for line in lines:
        ap(tx.text_frame,line,sz=fs,c=Th.CT,name='Consolas',al=PP_ALIGN.LEFT)

def make_title(prs,title,sub,part):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    rr(sl,Inches(1),Inches(1.5),Inches(11.333),Inches(4.5),Th.P)
    tb(sl,Inches(1.5),Inches(1.8),Inches(10.333),Inches(1),part,fs=12,c=Th.O,b=True)
    tb(sl,Inches(1.5),Inches(2.5),Inches(10.333),Inches(1.8),title,fs=36,b=True)
    tb(sl,Inches(1.5),Inches(4.2),Inches(10.333),Inches(1),sub,fs=18,c=Th.GR)

def make_obj(prs,items):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(0.08),Inches(7.5),Th.C)
    tb(sl,Inches(0.8),Inches(0.5),Inches(11),Inches(0.8),"Objectifs",fs=28,b=True,c=Th.C)
    for i,item in enumerate(items):
        rr(sl,Inches(1),Inches(1.6+i*1.1),Inches(11),Inches(0.8),Th.S)
        tb(sl,Inches(1.3),Inches(1.7+i*1.1),Inches(10.5),Inches(0.7),f"✦ {item}",fs=16)

def make_content(prs,title,lines):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.7),title,fs=26,b=True,c=Th.O)
    tx=tb_rich(sl,Inches(0.8),Inches(1.4),Inches(11.5),Inches(5.5))
    for line in lines:
        t,sz=line if isinstance(line,tuple) else (line,15)
        ap(tx.text_frame,t,sz=sz)

def make_code_content(prs,title,lines):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    tb(sl,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),title,fs=24,b=True,c=Th.O)
    code(sl,Inches(0.8),Inches(1.2),Inches(11.5),Inches(5.5),lines)

def make_code_tip(prs,title,lines,tip):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    tb(sl,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),title,fs=24,b=True,c=Th.O)
    code(sl,Inches(0.8),Inches(1.2),Inches(11.5),Inches(3.8),lines)
    rr(sl,Inches(0.8),Inches(5.3),Inches(11.5),Inches(1.2),Th.S)
    tb(sl,Inches(1.2),Inches(5.4),Inches(10.8),Inches(1),tip,fs=14,c=Th.C)

def make_ex(prs,title,desc,code_lines=None):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.G)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.6),title,fs=26,b=True,c=Th.G)
    rr(sl,Inches(0.8),Inches(1.2),Inches(11.5),Inches(3.5),Th.S)
    tx=tb_rich(sl,Inches(1.3),Inches(1.4),Inches(10.5),Inches(3.2))
    if isinstance(desc, list) and desc and isinstance(desc[-1], list):
        code_lines = desc[-1]
        desc = desc[:-1]
    txt = '\n'.join(desc) if isinstance(desc, list) else desc
    ap(tx.text_frame,txt,sz=15)
    if code_lines:
        code(sl,Inches(0.8),Inches(5),Inches(11.5),Inches(2.2),code_lines,fs=10)

def make_summary(prs,items):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(0.08),Inches(7.5),Th.O)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.7),"À retenir",fs=28,b=True,c=Th.O)
    for i,item in enumerate(items):
        rr(sl,Inches(1),Inches(1.5+i*0.9),Inches(11),Inches(0.65),Th.P)
        tb(sl,Inches(1.3),Inches(1.55+i*0.9),Inches(10.5),Inches(0.6),f"✓ {item}",fs=15,c=Th.LG)

def save(prs,name):
    p=os.path.join(OUT,name); prs.save(p); print(f"  ✅ {name}")

# ═══════════════════════════════════════════════════════════════
# PRESENTATIONS
# ═══════════════════════════════════════════════════════════════

def p35_exceptions_avancees():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Exceptions Avancées et Débogage","Assertions, logs, et debugging","Partie 4 - POO & Git")
    make_obj(prs,["Créer ses propres exceptions","Utiliser les assertions","Déboguer avec print et logging","Comprendre les stacks traces"])
    make_content(prs,"Créer ses propres exceptions",[
        ("Parfois, les exceptions Python ne sont pas assez précises.",15),
        "","Solution : créer SES PROPRES classes d'exception","",
        "class SoldeInsuffisantError(Exception):", '    """Exception quand le solde est insuffisant."""', "    pass",
        "","class CompteInexistantError(Exception):", '    """Exception quand le compte nexiste pas."""', "    def __init__(self, numero):", "        self.numero = numero", '        super().__init__(f"Compte {numero} inexistant")',
        "","Avantages :","  • Messages clairs et spécifiques","  • On peut les attraper SÉLECTIVEMENT","  • On peut ajouter des attributs personnalisés"
    ])
    make_code_content(prs,"Exceptions personnalisées en action",[
        "class SoldeInsuffisantError(Exception):", "    pass",
        "", "class CompteBancaire:", "    def __init__(self, solde=0):", "        self.solde = solde",
        "    def retirer(self, montant):", "        if montant > self.solde:", "            raise SoldeInsuffisantError(", '                f"Solde {self.solde} < {montant}")', "        self.solde -= montant",
        "", "try:", "    c = CompteBancaire(50)", "    c.retirer(100)", "except SoldeInsuffisantError as e:", "    print(f\"❌ {e}\")  # Solde 50 < 100",
        "", "# Avantage : on peut capturer UNIQUEMENT cette erreur", "# et laisser les autres exceptions normales passer"
    ])
    make_code_content(prs,"Assertions - assert",[
        "# assert condition, \"message si faux\"",
        "", "# Vérifier une précondition", "def diviser(a, b):", '    assert b != 0, "Division par zéro !"', "    return a / b",
        "", "# Vérifier une hypothèse", "def calculer_moyenne(notes):", '    assert len(notes) > 0, "Liste vide"', '    assert all(0 <= n <= 20 for n in notes), "Note invalide"', "    return sum(notes) / len(notes)",
        "", "# Test invariant", "def ajouter_element(liste, element):", "    taille_avant = len(liste)", "    liste.append(element)", "    assert len(liste) == taille_avant + 1",
        "", "# ASTUCE : les assertions sont DÉSACTIVÉES avec -O", "# python -O mon_programme.py", "# Utile en production : ne PAS utiliser assert pour la validation utilisateur !",
        "", "# Pour valider des entrées : privilégier if + raise ValueError"
    ])
    make_content(prs,"Déboguer efficacement",[
        ("Techniques de débogage :",0),
        "","1. print() - la méthode du pauvre (mais efficace)","   print(f\"DEBUG: x={x}, type={type(x)}\")",
        "","2. Le module logging (plus professionnel)","   import logging","   logging.basicConfig(level=logging.DEBUG)","   logging.debug(\"Message de debug\")","   logging.info(\"Information\")","   logging.warning(\"Attention !\")","   logging.error(\"Erreur !\")",
        "","3. pdb - Python Debugger","   import pdb; pdb.set_trace()  # point d'arrêt",
        "","4. breakpoint() - Python 3.7+ (recommandé)","   breakpoint()  # s'arrête ici, console interactive",
        "","5. Lecture de la stack trace (la plus IMPORTANTE !)"
    ])
    make_content(prs,"Lire une Stack Trace",[
        ("Quand une erreur arrive, Python affiche une stack trace :",15),
        "","Traceback (most recent call last):","  File \"main.py\", line 10, in <module>","    resultat = diviser(10, 0)","  File \"main.py\", line 5, in diviser","    return a / b","ZeroDivisionError: division by zero",
        "","🔑 À lire DU BAS VERS LE HAUT :","1. En bas : le TYPE d'erreur et le MESSAGE","2. Au-dessus : la LIGNE qui a causé l'erreur","3. Encore au-dessus : la FONCTION qui appelait...","4. Tout en haut : le point d'entrée",
        "","Exercice : trouve l'erreur dans :", "def traiter(data): return data['valeur'] + 10", "config = {'nom': 'test'}", "print(traiter(config))"
    ])
    make_ex(prs,"Exercice : Validation robuste",[
        "Crée un système de gestion de bibliothèque avec :\n\n1. Exception LivreNonTrouveError (personnalisée)\n2. Exception EmpruntImpossibleError (avec raison)\n3. Classe Livre avec : titre, auteur, emprunté\n4. Classe Bibliotheque avec :\n   - ajouter(titre, auteur)\n   - emprunter(titre) → lève LivreNonTrouveError ou EmpruntImpossibleError\n   - retourner(titre)\n   - lister()\n5. Utilise assert pour valider que les titres ne sont pas vides\n6. Ajoute logging pour tracer les emprunts/retours\n7. Ajoute un bloc try/except dans le menu principal"
    ])
    make_summary(prs,["Créer ses exceptions : class MonErreur(Exception)","assert condition, message → vérifier des hypothèses","logging : DEBUG, INFO, WARNING, ERROR","breakpoint() → console interactive","Stack trace se lit du bas vers le haut","assert désactivé en production (-O), logging toujours actif"])
    save(prs,"35_Exceptions_avancees.pptx")
    return prs

def p36_mini_projet_jeu_cartes():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Mini-Projet : Jeu de Cartes","POO appliquée à un jeu","Partie 4 - POO & Git")
    make_obj(prs,["Appliquer la POO à un projet concret","Créer Carte, Paquet, et Jeu","Implémenter la logique d'un jeu de cartes","Utiliser random.shuffle et les stratégies"])
    make_content(prs,"Le projet : Bataille (War)",[
        ("On va créer le jeu de la BATAILLE !",15),
        "","Règles simplifiées :","  1. On distribue toutes les cartes aux 2 joueurs","  2. Chaque joueur retourne la carte du dessus","  3. Le plus fort remporte les 2 cartes","  4. En cas d'égalité → on rejoue (double la mise)",
        "","Classes à créer :","  • Carte : valeur, couleur, points","  • Paquet : liste de cartes, mélanger, distribuer","  • Joueur : nom, main de cartes, points","  • Jeu : logique de la bataille",
        "","Structure :","jeu_bataille/","├── carte.py","├── paquet.py","├── joueur.py","├── jeu.py","└── main.py"
    ])
    make_code_content(prs,"Classe Carte",[
        "# carte.py", "class Carte:", "    VALEURS = ['2','3','4','5','6','7','8','9','10','J','Q','K','A']", '    COULEURS = ["♠","♣","♥","♦"]',
        "    def __init__(self, valeur, couleur):", "        self.valeur = valeur", "        self.couleur = couleur", "        self.points = Carte.VALEURS.index(valeur) + 2",
        "    def __str__(self):", '        return f"{self.valeur}{self.couleur}"',
        "    def __repr__(self):", "        return self.__str__()",
        "    def __lt__(self, other):", "        return self.points < other.points",
        "", "c = Carte('Roi', '♠')", "print(c)            # Roi♠", "print(c.points)      # 13 (Roi = forte valeur)", "print(c < Carte('As','♠'))  # True"
    ])
    make_code_content(prs,"Classes Paquet et Joueur",[
        "# paquet.py", "import random", "from carte import Carte",
        "class Paquet:", "    def __init__(self):", "        self.cartes = [Carte(v,c) for v in Carte.VALEURS for c in Carte.COULEURS]",
        "    def melanger(self):", "        random.shuffle(self.cartes)",
        "    def distribuer(self, nb_joueurs):", "        return [self.cartes[i::nb_joueurs] for i in range(nb_joueurs)]",
        "", "# joueur.py", "class Joueur:", "    def __init__(self, nom):", "        self.nom = nom", "        self.main = []  # liste de cartes", "        self.score = 0",
        "    def piocher(self):", "        return self.main.pop(0) if self.main else None",
        "    def gagner_cartes(self, cartes):", "        self.main.extend(cartes)", "        self.score += len(cartes)",
        "    @property", "    def nb_cartes(self):", "        return len(self.main)"
    ])
    make_code_content(prs,"Classe Jeu (Bataille)",[
        "# jeu.py", "from paquet import Paquet", "from joueur import Joueur",
        "class Bataille:", "    def __init__(self, nom1, nom2):", "        paquet = Paquet()", "        paquet.melanger()", "        main1, main2 = paquet.distribuer(2)", "        self.j1 = Joueur(nom1); self.j1.main = main1", "        self.j2 = Joueur(nom2); self.j2.main = main2",
        "    def jouer_tour(self):", "        c1 = self.j1.piocher(); c2 = self.j2.piocher()", "        if not c1 or not c2: return False",
        "        print(f\"{self.j1.nom}: {c1}  vs  {self.j2.nom}: {c2}\")",
        "        if c1 > c2:", "            self.j1.gagner_cartes([c1, c2])", '            print(f"{self.j1.nom} gagne !")', "        elif c2 > c1:", "            self.j2.gagner_cartes([c1, c2])", '            print(f"{self.j2.nom} gagne !")', "        else:", "            print(\"Bataille !\")", "            self.bataille([c1, c2])", "        return True"
    ])
    make_code_content(prs,"Gestion de la bataille (égalité)",[
        "# Dans la classe Bataille",
        "    def bataille(self, cartes_au_centre):", "        # Chaque joueur met 2 cartes face cachée + 1 face visible", "        for _ in range(2):", "            c = self.j1.piocher()", "            if c: cartes_au_centre.append(c)", "            c = self.j2.piocher()", "            if c: cartes_au_centre.append(c)",
        "        c1 = self.j1.piocher(); c2 = self.j2.piocher()", "        if not c1 or not c2: return",
        "        cartes_au_centre.extend([c1, c2])",
        "        if c1 > c2:", "            self.j1.gagner_cartes(cartes_au_centre)", "        elif c2 > c1:", "            self.j2.gagner_cartes(cartes_au_centre)", "        else:", "            self.bataille(cartes_au_centre)  # récursif !",
        "    def jouer(self):", "        tour = 0", "        while self.j1.nb_cartes > 0 and self.j2.nb_cartes > 0:", "            tour += 1", "            if not self.jouer_tour(): break",
        '        print(f"\\n🏆 {self.j1.nom}: {self.j1.score} cartes")', '        print(f"🏆 {self.j2.nom}: {self.j2.score} cartes")'
    ])
    make_ex(prs,"À toi de jouer !",[
        "Étapes :\n1. Crée tous les fichiers (carte.py, paquet.py, joueur.py, jeu.py, main.py)\n2. Teste d'abord chaque classe séparément\n3. Joue une partie complète\n4. Améliorations possibles :\n   • Afficher le nombre de cartes restantes à chaque tour\n   • Limiter à 1000 tours max (sinon boucle infinie)\n   • Ajouter des statistiques (tours joués, bataille gagnées par qui)\n   • Sauvegarder les scores dans scores.json\n   • Permettre de choisir le nombre de joueurs\n\n💡 Un jeu de cartes est PARFAIT pour apprendre la POO"
    ])
    make_summary(prs,["POO concrète : Carte, Paquet, Joueur, Jeu","random.shuffle() pour mélanger","__lt__ pour comparer les cartes","Récursivité pour la bataille","Projet modulaire : un fichier par classe","Un jeu = excellent exercice de programmation"])
    save(prs,"36_Mini_Projet_Jeu_Cartes.pptx")
    return prs

def p37_git_intro():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Git - Introduction","Pourquoi versionner son code ?","Partie 4 - POO & Git")
    make_obj(prs,["Comprendre le versionnage","Savoir ce qu'est Git","Installer et configurer Git","Créer son premier dépôt"])
    make_content(prs,"Pourquoi versionner son code ?",[
        ("Tu as déjà eu ce problème ?",15),
        "","❌ \"mon_projet_final.py\" puis \"mon_projet_final2.py\"","❌ \"J'ai supprimé une fonction par erreur et je ne peux pas revenir\"","❌ \"Je veux essayer une nouvelle idée sans casser mon code\"","❌ \"Comment travailler à plusieurs sur le même projet ?\"",
        "","✅ Git résout TOUS ces problèmes !",
        "","Git = un système de contrôle de version (VCS)","  • Sauvegarde l'HISTORIQUE de chaque fichier","  • Permet de REVENIR en arrière","  • Permet de CRÉER DES BRANCHES pour expérimenter","  • Facilite le TRAVAIL EN ÉQUIPE","",
        "Créé par Linus Torvalds (créateur de Linux) en 2005","Aujourd'hui utilisé par presque TOUS les développeurs"
    ])
    make_content(prs,"Concepts clés",[
        ("Dépôt (Repository/Repo) : dossier suivi par Git",15),
        "","  • LOCAL : sur ton ordinateur","  • DISTANT (remote) : sur GitHub, GitLab, etc.",
        "","Commit : un INSTANTANÉ (snapshot) de ton projet à un moment T","  • Comme une sauvegarde d'un jeu vidéo","  • Chaque commit a un message qui décrit les changements",
        "","  📸 Commit 1 → 📸 Commit 2 → 📸 Commit 3 → ...",
        "",
        "Index (Staging Area) : zone intermédiaire","  • Tu prépares quels fichiers seront dans le prochain commit","  • Tu peux ajouter des fichiers UN PAR UN",
        "","  📝 Modifier → ➕ Stage → 📸 Commit → 🔄 Push (distant)"
    ])
    make_code_content(prs,"Installer et configurer Git",[
        "# 1. Vérifier si Git est installé", "git --version  # git version 2.x.x",
        "",
        "# Si pas installé :", "# sudo apt install git  (Ubuntu/Debian)", "# sudo dnf install git (Fedora)", "# https://git-scm.com/  (Windows/macOS)",
        "",
        "# 2. Configurer son identité (À FAIRE UNE FOIS)", "git config --global user.name \"Ton Pseudo\"", "git config --global user.email \"ton.email@example.com\"",
        "",
        "# 3. Configurer l'éditeur (optionnel)", "git config --global core.editor \"code --wait\"  # VS Code",
        "",
        "# 4. Voir la configuration", "git config --list",
        "",
        "# 5. Couleurs (plus lisible)", "git config --global color.ui auto",
        "",
        "# Ces infos sont incluses dans chaque commit (d'où l'importance)"
    ])
    make_code_content(prs,"Premier dépôt",[
        "# 1. Créer un dossier", "mkdir mon_premier_projet", "cd mon_premier_projet",
        "", "# 2. Initialiser Git", "git init", "# Crée un dossier caché .git/ qui contient tout l'historique",
        "", "# 3. Créer un fichier", 'echo "print(\"Hello Git!\")" > main.py',
        "", "# 4. Voir l'état", "git status  # main.py est en rouge = non suivi",
        "", "# 5. Ajouter au stage", "git add main.py  # main.py passe en vert",
        "", "# 6. Premier commit", 'git commit -m "Premier commit : Hello Git"',
        "", "# 7. Voir l'historique", "git log  # affiche le commit avec auteur, date, message"
    ])
    make_content(prs,".gitignore - Ignorer des fichiers",[
        ("Parfois, il ne faut PAS versionner certains fichiers :",15),
        "","❌ Dossier __pycache__/ (cache Python)","❌ .env (clés API, mots de passe)","❌ *.pyc (fichiers compilés)","❌ venv/ (environnement virtuel)","❌ .DS_Store (macOS)","❌ node_modules/ (JavaScript)","❌ Fichiers volumineux (>100 Mo)",
        "","Solution : fichier .gitignore à la racine du projet",
        "","  # .gitignore", "  __pycache__/", "  *.pyc", "  .env", "  venv/", "  *.log",
        "","Git ignore automatiquement ces fichiers/dossiers."
    ])
    make_ex(prs,"Exercice : Premier dépôt Git",[
        "1. Crée un dossier \"apprentissage_git\"\n2. Initialise un dépôt Git\n3. Crée 3 fichiers Python : hello.py, calcul.py, et .gitignore\n4. Dans hello.py : print(\"Hello!\")\n5. Dans calcul.py : print(2+2)\n6. Dans .gitignore : ignore __pycache__/ et *.pyc\n7. Fait git status → observe les fichiers non suivis\n8. Fait git add . (ajoute tout)\n9. Fait git commit -m \"Premier commit\"\n10. Fait git log pour voir l'historique\n11. Modifie hello.py et recommence : add + commit"
    ])
    make_summary(prs,["Git = système de contrôle de version","Dépôt (repo) = dossier suivi par Git","Commit = instantané du projet","git init → git add → git commit : cycle de base",".gitignore = fichiers à ne PAS versionner","git log pour voir l'historique"])
    save(prs,"37_Git_Introduction.pptx")
    return prs

def p38_git_commandes():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Git - Commandes Quotidiennes","Les commandes essentielles","Partie 4 - POO & Git")
    make_obj(prs,["Maîtriser le cycle de vie des fichiers","Comprendre git diff et git log","Utiliser git reset et git restore","Comparer les versions avec git diff"])
    make_content(prs,"Cycle de vie des fichiers",[
        ("Chaque fichier peut être dans 4 états :",15),
        "","┌──────────┐   git add   ┌──────────┐   git commit   ┌──────────┐","│ UNTRACKED │ ──────────→ │  STAGED  │ ──────────────→ │ COMMITTED │","│  (non     │             │  (index) │                 │  (sauvé)  │","│   suivi)  │             └──────────┘                 └──────────┘","└──────────┘","    ↑   git add . ────────────────────┘","    │                                     ","    └──── git checkout ──→ retour arrière",
        "","MODIFIED (modifié) : le fichier a changé depuis le dernier commit","STAGED (indexé) : prêt à être commité","COMMITTED (commit) : sauvegardé dans l'historique",
        "","Commande magique : git status à chaque étape !"
    ])
    make_code_content(prs,"Afficher les différences - git diff",[
        "# 1. Modifications NON encore indexées (working dir)", "git diff",
        "# Affiche les lignes ajoutées (+) et supprimées (-)",
        "",
        "# 2. Modifications DÉJÀ indexées (staged)", "git diff --staged",
        "",
        "# 3. Comparer deux commits", "git log --oneline  # voir les hash", "git diff abc123..def456",
        "",
        "# Exemple de sortie :", "# diff --git a/main.py b/main.py", "# --- a/main.py", "# +++ b/main.py", '# @@ -1 +1,2 @@', '#  print("Hello")', '# +print("Nouvelle ligne")',
        "",
        "# 4. git diff avec un fichier spécifique", "git diff main.py"
    ])
    make_code_content(prs,"Visualiser l'historique - git log",[
        "# 1. Log simple", "git log",
        "",
        "# 2. Log compact (1 ligne par commit)", "git log --oneline",
        "# a1b2c3d Ajout de la calculatrice", "# e4f5g6h Ajout du menu principal", "# i7j8k9l Premier commit",
        "",
        "# 3. Log graphique (joli !)", "git log --oneline --graph --all",
        "# * a1b2c3d (HEAD -> main) Ajout calculatrice", "# * e4f5g6h Ajout menu", "# * i7j8k9l Premier commit",
        "",
        "# 4. Log avec dates et auteurs", "git log --oneline --graph --all --decorate",
        "",
        "# 5. Log avec stats", "git log --stat  # fichiers modifiés et lignes",
        "",
        "# 6. Qui a écrit quoi ?", "git blame main.py  # chaque ligne avec auteur et commit"
    ])
    make_code_content(prs,"Annuler des changements",[
        "# 1. Modifié mais PAS encore indexé → restaurer la version commitée", "git restore main.py",
        "# Alternative: git checkout -- main.py (ancienne syntaxe)",
        "",
        "# 2. DÉJÀ indexé (staged) → enlever du stage", "git restore --staged main.py",
        "# Le fichier redevient MODIFIED (non staged)",
        "",
        "# 3. Modifier le dernier commit (oups, message ou oubli)", "# Si pas encore pushé :", "git commit --amend -m \"Nouveau message\"",
        "# Ajoute aussi les fichiers staged au dernier commit",
        "",
        "# 4. Supprimer un fichier", "rm fichier.txt", "git rm fichier.txt  # supprime + stage la suppression",
        "",
        "# 5. Déplacer/renommer", "git mv ancien.py nouveau.py  # renomme + stage",
        "",
        "# ATTENTION : ne JAMAIS faire --amend sur un commit déjà pushé !"
    ])
    make_content(prs,"Bonnes pratiques Git",[
        ("Règles d'or pour bien utiliser Git :",15),
        "","📝 Commits ATOMIQUES : un commit = une modification logique","   ❌ \"plein de trucs\"","   ✅ \"Ajout de la fonction de tri\"",
        "","✍️ Messages clairs en français ou anglais :","   ❌ \"fix\"","   ✅ \"Correction du bug de division par zéro\"",
        "","🔄 Commiter SOUVENT (après chaque petite étape fonctionnelle)","   → Plus facile de revenir en arrière",
        "","📄 Commiter le code qui MARCHE (pas de code cassé)","📁 .gitignore avant le premier commit",
        "","🔑 git status → git add → git commit : le cycle de vie"
    ])
    make_ex(prs,"Exercice : Manipuler l'historique",[
        "Dans le dépôt de l'exercice précédent :\n\n1. Ajoute une fonction addition(a,b) dans calcul.py\n2. Fait git diff (sans commit) → observe\n3. Fait git add calcul.py\n4. Fait git diff --staged → observe\n5. Fait git commit -m \"Ajout addition\"\n6. Ajoute soustraction(a,b) dans calcul.py\n7. Ajoute aussi multiplication(a,b) mais NE STAGE QUE soustraction\n8. Fait git status → un staged, un modified\n9. Commit avec \"Ajout soustraction\"\n10. Fait git log --oneline --graph\n11. Fait git show HEAD (dernier commit en détail)\n12. BONUS : utilise git restore pour annuler multiplication non commitée"
    ])
    make_summary(prs,["4 états : UNTRACKED → STAGED → COMMITTED (+ MODIFIED)","git diff = changements non indexés, git diff --staged = indexés","git log --oneline --graph pour l'historique","git restore = annuler, git commit --amend = corriger dernier commit","Commits atomiques, messages clairs, commiter souvent","git status est ton meilleur ami !"])
    save(prs,"38_Git_commandes.pptx")
    return prs

def p39_git_branches_github():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Git - Branches et GitHub","Collaborer et partager son code","Partie 4 - POO & Git")
    make_obj(prs,["Comprendre les branches","Créer et fusionner des branches","Découvrir GitHub","Push/Pull et travail collaboratif"])
    make_content(prs,"Qu'est-ce qu'une branche ?",[
        ("Une BRANCHE = une version PARALLÈLE de ton projet.",15),
        "","Imagine un arbre :","  • Le tronc = la branche principale (main ou master)","  • Les branches = des expériences, nouvelles fonctionnalités","  • Tu peux couper une branche, y travailler, et la rattacher",
        "","Pourquoi des branches ?","  • Tester une nouvelle idée SANS casser le code principal","  • Travailler à plusieurs sur des fonctionnalités différentes","  • Corriger un bug urgent pendant qu'on développe ailleurs",
        "","  main:    ●──●─────────●──────────●──●","                ╲              ╱","  feature:       ●──●──●──●──●      (fusion = merge)",
        "","Par défaut, tu es sur la branche main/master"
    ])
    make_code_content(prs,"Créer et utiliser des branches",[
        "# 1. Voir les branches existantes", "git branch  # * = branche active",
        "",
        "# 2. Créer une nouvelle branche", "git branch nouvelle-fonctionnalite",
        "",
        "# 3. Changer de branche", "git checkout nouvelle-fonctionnalite",
        "# Ou en une commande (créer + basculer):", "git checkout -b autre-fonctionnalite",
        "",
        "# 4. Travailler sur la branche (add + commit normal)", "echo 'print(\"Nouveau code\")' > test.py", "git add test.py", "git commit -m \"Ajout test sur feature\"",
        "",
        "# 5. Revenir sur main", "git checkout main",
        "# Le fichier test.py a DISPARU ! (il est sur l'autre branche)"
    ])
    make_code_content(prs,"Fusionner (merge) et supprimer",[
        "# 1. Fusionner une branche dans main", "# D'abord, assure-toi d'être sur main", "git checkout main",
        "# Puis fusionne", "git merge nouvelle-fonctionnalite",
        "# Tout le contenu de la branche arrive dans main",
        "",
        "# 2. Supprimer une branche (après merge réussi)", "git branch -d nouvelle-fonctionnalite",
        "",
        "# 3. Gérer les conflits (quand Git ne sait pas quoi faire)", "# Exemple: fichier modifié sur les 2 branches aux mêmes endroits", "# Git met des marqueurs :", "# <<<<<<< HEAD", "# code de main", "# =======", "# code de la branche", "# >>>>>>> feature",
        "# Tu dois ÉDITER le fichier manuellement pour choisir",
        "git add fichier-resolu.txt", "git commit -m \"Résolution conflit merge\""
    ])
    make_content(prs,"GitHub - Pourquoi ?",[
        ("GitHub = un site web qui héberge des dépôts Git.",15),
        "","Avantages :","  ☁️ SAUVEGARDE en ligne de ton code","  🌍 PARTAGE avec le monde","  👥 COLLABORATION avec d'autres développeurs","  🔄 SYNCHRONISATION entre plusieurs ordis","  🎯 PORTFOLIO pour montrer tes projets",
        "","Alternatives : GitLab, Bitbucket (mais GitHub est le plus populaire)",
        "","Concepts importants :","  • remote (distant) = le dépôt sur GitHub","  • origin = nom par défaut du remote","  • push = envoyer vers GitHub","  • pull = récupérer depuis GitHub","  • clone = copier un dépôt GitHub en local"
    ])
    make_code_content(prs,"Travailler avec GitHub",[
        "# 1. Créer un dépôt sur GitHub (via le site web)",
        "#    Ne PAS cocher \"Initialize with README\"",
        "",
        "# 2. Connecter le dépôt local au distant", "git remote add origin https://github.com/tonpseudo/mon-projet.git",
        "# origin = nom du remote (convention)",
        "",
        "# 3. Pousser (push) le code", "git push -u origin main",
        "# -u = set upstream (souviens-toi du remote par défaut)",
        "",
        "# 4. Cloner un dépôt existant", "git clone https://github.com/autre/util/projet.git",
        "# Crée un dossier \"projet\" avec tout l'historique",
        "",
        "# 5. Récupérer les changements", "git pull origin main  # fetch + merge en un coup",
        "# ou: git fetch + git merge (plus de contrôle)",
        "",
        "# 6. Voir les remotes", "git remote -v  # liste les dépôts distants"
    ])
    make_content(prs,"Pull Requests (PR) - Le travail d'équipe",[
        ("Une Pull Request = une demande pour fusionner ta branche dans une autre.",15),
        "","Flux de travail typique :","  1. Fork (copie) du dépôt principal","  2. Clone sur ton ordi","  3. Crée une branche (git checkout -b ma-fonctionnalite)","  4. Travaille et commit","  5. Push la branche vers GitHub","  6. Crée une Pull Request sur GitHub","  7. Discussion + revue de code","  8. Merge !",
        "","💡 Les PR permettent :","  • La REVUE DE CODE par les autres","  • Les TESTS AUTOMATIQUES avant fusion","  • La DISCUSSION sur les changements",
        "","C'est comme ça que les projets Open Source fonctionnent !"
    ])
    make_ex(prs,"Exercice : Branches et GitHub",[
        "1. Dans \"apprentissage_git\", crée une branche \"calculatrice\"\n2. Bascule dessus, ajoute un fichier calc.py avec +, -, *, /\n3. Commit (\"Ajout calculatrice\")\n4. Revient sur main\n5. Crée une branche \"interface\"\n6. Bascule dessus, ajoute un input() dans main.py\n7. Commit (\"Ajout interface utilisateur\")\n8. Revient sur main, merge \"calculatrice\" → tout le fichier calc.py arrive\n9. Merge \"interface\" → gère conflit si besoin\n10. Fait git log --oneline --graph pour admirer le résultat\n11. BONUS : crée un compte GitHub, initie un dépôt, push ton projet"
    ])
    make_summary(prs,["Branche = version parallèle (git branch, checkout, merge)","Conflit = Git ne sait pas quoi prendre → résolution manuelle","GitHub = hébergement distant + collaboration","push = envoyer, pull = recevoir, clone = copier","Pull Request = demande de fusion avec revue de code","Le travail en équipe devient FACILE avec Git+GitHub"])
    save(prs,"39_Git_branches_github.pptx")
    return prs

def p40_mini_projet_pendu():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Mini-Projet : Le Pendu","Jeu complet avec Git et GitHub","Partie 4 - POO & Git")
    make_obj(prs,["Créer un jeu complet","Appliquer Git tout au long du projet","Organiser le code en modules","Publier sur GitHub"])
    make_content(prs,"Le projet : Jeu du Pendu",[
        ("Règles : trouver un mot en devinant les lettres, une par une.",15),
        "","  • Un mot secret est choisi aléatoirement","  • Le joueur propose des lettres","  • Si la lettre est dans le mot → elle est révélée","  • Sinon → on dessine une partie du pendu","  • 6 erreurs max → PERDU","  • Toutes les lettres trouvées → GAGNÉ",
        "","Structure du projet :","pendu/","├── mots.py        # Liste de mots","├── pendu.py       # Logique du jeu","├── affichage.py   # Interface console + dessin","├── scores.py      # Sauvegarde des scores","└── main.py        # Point d'entrée",
        "","On va utiliser Git pour versionner chaque étape !"
    ])
    make_code_content(prs,"mots.py - Gestion des mots",[
        "# mots.py", "import random",
        "MOTS = [", '    "python", "ordinateur", "programmation", "developpeur",', '    "algorithme", "variable", "fonction", "boucle",', '    "condition", "dictionnaire", "liste", "module",', '    "internet", "navigation", "clavier", "ecran",', '    "souris", "fichier", "reseau", "application"',
        "]",
        "", "# Sélection par difficulté", "def choisir_mot(difficulte=\"facile\"):", "    if difficulte == \"facile\":", "        return random.choice([m for m in MOTS if len(m) <= 6])", "    elif difficulte == \"moyen\":", "        return random.choice([m for m in MOTS if 7 <= len(m) <= 10])", "    else:  # difficile", "        return random.choice([m for m in MOTS if len(m) > 10])"
    ])
    make_code_content(prs,"pendu.py - Logique du jeu",[
        "# pendu.py", "class Pendu:", "    MAX_ERREURS = 6",
        "    def __init__(self, mot_secret):", "        self.mot = mot_secret.lower()", "        self.trouve = ['_'] * len(self.mot)", "        self.erreurs = 0", "        self.lettres_proposees = set()",
        "    def proposer(self, lettre):", "        lettre = lettre.lower()", "        if lettre in self.lettres_proposees:", '            return "déjà proposée"', "        self.lettres_proposees.add(lettre)",
        "        if lettre in self.mot:", "            for i, c in enumerate(self.mot):", "                if c == lettre:", "                    self.trouve[i] = c", '            return "trouvée"', "        else:", "            self.erreurs += 1", '            return "ratée"',
        "    @property", "    def est_gagne(self):", "        return '_' not in self.trouve",
        "    @property", "    def est_perdu(self):", "        return self.erreurs >= self.MAX_ERREURS"
    ])
    make_code_content(prs,"affichage.py - Le dessin du pendu",[
        "# affichage.py", "DESSINS = [", '  """\n    +---+\n        |\n        |\n        |\n       ===\n    """,',
        '  """\n    +---+\n    O   |\n        |\n        |\n       ===\n    """,',
        '  """\n    +---+\n    O   |\n    |   |\n        |\n       ===\n    """,',
        '  """\n    +---+\n    O   |\n   /|   |\n        |\n       ===\n    """,',
        '  """\n    +---+\n    O   |\n   /|\\  |\n        |\n       ===\n    """,',
        '  """\n    +---+\n    O   |\n   /|\\  |\n   /    |\n       ===\n    """,',
        '  """\n    +---+\n    O   |\n   /|\\  |\n   / \\  |\n       ===\n    """',
        "]",
        "def afficher_partie(pendu):", "    print(DESSINS[pendu.erreurs])", '    print("Mot: " + " ".join(pendu.trouve))', '    print(f"Erreurs: {pendu.erreurs}/{Pendu.MAX_ERREURS}")', '    print(f"Lettres: {",".join(sorted(pendu.lettres_proposees))}")'
    ])
    make_code_content(prs,"Scores et main.py",[
        "# scores.py", "import json, os",
        "FICHIER = \"scores_pendu.json\"",
        "def charger_scores():", "    if not os.path.exists(FICHIER): return []", "    with open(FICHIER) as f: return json.load(f)",
        "def sauvegarder_score(pseudo, erreurs, mot):", "    scores = charger_scores()", "    scores.append({\"pseudo\":pseudo,\"erreurs\":erreurs,\"mot\":mot})", "    with open(FICHIER,\"w\") as f:", "        json.dump(scores, f, indent=2)",
        "def meilleur_score():", "    scores = charger_scores()", "    if not scores: return None", "    return min(scores, key=lambda s: s[\"erreurs\"])",
        "",
        "# main.py", "from pendu import Pendu", "from mots import choisir_mot", "from affichage import afficher_partie", "from scores import sauvegarder_score, meilleur_score",
        "def jouer():", "    mot = choisir_mot()", "    jeu = Pendu(mot)", "    while not jeu.est_gagne and not jeu.est_perdu:", "        afficher_partie(jeu)", "        lettre = input(\"Lettre: \")", "        jeu.proposer(lettre)",
        "    if jeu.est_gagne:", '        print(f"GAGNÉ ! {mot}")', "    else:", '        print(f"PERDU ! Le mot était {mot}")'
    ])
    make_ex(prs,"À toi de jouer avec Git !",[
        "Étapes Git (fais un commit à chaque étape) :\n\n1. git init + git add . + commit \"Initialisation\"\n2. Crée mots.py → commit\n3. Crée pendu.py → commit\n4. Crée affichage.py → git add/commit\n5. Crée scores.py → commit\n6. Crée main.py → commit \"Jeu fonctionnel\"\n7. git checkout -b améliorations\n8. Ajoute difficulté, catégories, ou couleurs → commit\n9. git checkout main + git merge améliorations\n10. BONUS : crée un repo GitHub et push",
        ["# Améliorations possibles", "# Couleurs avec ANSI:", 'def colorer(texte, couleur):', '    codes = {"rouge":"31","vert":"32","jaune":"33","bleu":"34"}', '    return f"\\033[{codes[couleur]}m{texte}\\033[0m"']
    ])
    make_summary(prs,["Pendu = projet parfait pour consolider Python + Git","Architecture modulaire : mots, logique, affichage, scores","Git commit après chaque module fonctionnel","Branche pour les améliorations","GitHub pour partager et sauvegarder","La pratique CONSTANTE est la clé de l'apprentissage"])
    save(prs,"40_Mini_Projet_Pendu.pptx")
    return prs

def p41_environnements_virtuels():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Environnements Virtuels et pip","Isoler ses projets Python","Partie 5 - Projets Avancés")
    make_obj(prs,["Comprendre les environnements virtuels","Créer et activer un venv","Utiliser pip correctement","Gérer les dépendances avec requirements.txt"])
    make_content(prs,"Pourquoi des environnements virtuels ?",[
        ("Problème : chaque projet Python peut avoir besoin de versions différentes des mêmes bibliothèques.",15),
        "","  Projet A : utilise Django 4.2 (a besoin de Python 3.10+)", "  Projet B : utilise Django 3.2 (marche avec Python 3.8)",
        "","Sans venv : impossible d'avoir les DEUX versions installées en même temps","(ou alors c'est le chaos total)",
        "","Solution : ENVIRONNEMENT VIRTUEL = une copie isolée de Python","  • Chaque projet a SON PROPRE Python et SES PROPRES paquets","  • Pas de conflit entre les projets",
        "","  Système : Python 3.14 + paquets système","    ├── Projet A : venv avec Django 4.2 + pandas 2.0","    └── Projet B : venv avec Django 3.2 + numpy 1.21",
        "","C'est la norme dans l'industrie !"
    ])
    make_code_content(prs,"Créer et utiliser un venv",[
        "# 1. Créer un environnement virtuel", "python -m venv mon_env",
        "# Crée un dossier mon_env/ avec :", "#   bin/   → python, pip activé", "#   lib/   → paquets installés", "#   pyvenv.cfg → config",
        "",
        "# 2. Activer l'environnement",
        "# Linux/macOS:", "source mon_env/bin/activate",
        "# Windows:", "mon_env\\Scripts\\activate",
        "",
        "# Après activation :", "# (mon_env) $  → le nom apparait dans le terminal",
        "which python  # → .../mon_env/bin/python", "python --version  # version du venv",
        "",
        "# 3. Désactiver", "deactivate"
    ])
    make_code_content(prs,"pip - Le gestionnaire de paquets",[
        "# pip = Python Install Packages",
        "",
        "# Installer un paquet", "pip install requests",
        "# Version spécifique", "pip install django==4.2",
        "# Version minimale", "pip install \"flask>=2.0\"",
        "",
        "# Désinstaller", "pip uninstall requests",
        "",
        "# Lister les paquets installés", "pip list",
        "# Afficher les infos d'un paquet", "pip show requests",
        "",
        "# Mettre à jour un paquet", "pip install --upgrade pip",
        "",
        "# Rechercher un paquet sur PyPI", "pip search mot-cle  # (déprécié, utiliser le site web)",
        "",
        "# Installer depuis un fichier", "pip install -r requirements.txt"
    ])
    make_code_content(prs,"requirements.txt - Les dépendances",[
        "# requirements.txt = liste des paquets nécessaires au projet",
        "",
        "# Créer requirements.txt (à partir de l'environnement actif)", "pip freeze > requirements.txt",
        "# Résultat :", "# requests==2.31.0", "# flask==3.0.0", "# numpy==1.26.0",
        "",
        "# Installer les dépendances d'un projet", "pip install -r requirements.txt",
        "",
        "# workflow typique :", "# 1. git clone mon-projet", "# 2. python -m venv venv", "# 3. source venv/bin/activate", "# 4. pip install -r requirements.txt", "# 5. python main.py",
        "",
        "# ATTENTION : NE PAS commiter le dossier venv !", "# Ajouter venv/ dans .gitignore",
        "",
        "# Différence entre paquets :", "# - Standard : random, math, os, sys (toujours disponibles)", "# - Externes : requests, flask, numpy (à installer via pip)"
    ])
    make_content(prs,"Bonnes pratiques",[
        ("Règles à suivre :",15),
        "","1. Toujours créer un venv pour chaque projet","   python -m venv .venv  (point + nom = convention moderne)",
        "","2. Activer le venv avant de travailler sur le projet","   source .venv/bin/activate",
        "","3. Toujours utiliser pip install dans le venv, PAS en système","   (évite de casser Python système)",
        "","4. requirements.txt versionné (mais PAS le dossier venv)",
        "","5. Fichier .gitignore :", "   .venv/", "   __pycache__/", "   *.pyc",
        "","6. Vérifie : python -m pip install (plus sûr que pip install)",
        "","7. python -m venv .venv --prompt mon_projet  (nom personnalisé)"
    ])
    make_ex(prs,"Exercice : Mise en place propre",[
        "1. Crée un dossier \"mon_app_flask\"\n2. Crée un venv : python -m venv .venv\n3. Active le venv\n4. Installe flask : pip install flask\n5. Crée un fichier app.py avec :\n     from flask import Flask\n     app = Flask(__name__)\n     @app.route('/')\n     def hello(): return \"Hello Venv!\"\n     app.run(debug=True)\n6. Lance l'app : python app.py\n7. Va sur http://localhost:5000\n8. Fait pip freeze > requirements.txt\n9. Désactive le venv, réactive-le, pip install -r requirements.txt\n10. Ajoute .venv/ dans .gitignore et initialise Git"
    ])
    make_summary(prs,["Environnement virtuel = Python isolé par projet","python -m venv .venv → source .venv/bin/activate → deactivate","pip = installation de paquets externes","pip freeze > requirements.txt → partage les dépendances","NE JAMAIS versionner le dossier venv","Un venv par projet = bonne pratique INDISPENSABLE"])
    save(prs,"41_Environnements_virtuels.pptx")
    return prs

def p42_apis_http():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"APIs HTTP avec requests","Communiquer avec le web","Partie 5 - Projets Avancés")
    make_obj(prs,["Comprendre ce qu'est une API","Connaître les méthodes HTTP","Utiliser le module requests","Travailler avec JSON et APIs"])
    make_content(prs,"Qu'est-ce qu'une API ?",[
        ("API = Application Programming Interface",15),
        "","Interface de communication entre programmes.",
        "","Métaphore du restaurant :","  • Toi = ton programme","  • Le menu = l'API (les actions possibles)","  • Le serveur = la requête API","  • La cuisine = le serveur distant","  • Ton plat = la réponse JSON",
        "","Concrètement : ton programme envoie une REQUÊTE à un SERVEUR distant","et reçoit une RÉPONSE (souvent en JSON)",
        "","Exemples d'API publiques :","  • météo → api.openweathermap.org","  • GitHub → api.github.com","  • Météo France, PokéAPI, etc.",
        "","/!\\ Besoin d'un module externe : pip install requests"
    ])
    make_content(prs,"Méthodes HTTP (CRUD)",[
        ("HTTP = le protocole du web. Chaque requête a une MÉTHODE :",15),
        "","  GET    → LIRE des données (comme SELECT en base)","  POST   → CRÉER des données (comme INSERT)","  PUT    → METTRE À JOUR des données (comme UPDATE)","  DELETE → SUPPRIMER des données (comme DELETE)",
        "","🔑 GET est la plus utilisée (et la seule dont on a besoin pour débuter)",
        "","URL d'API typique :","  https://api.github.com/users/octocat","  │          │              │             │","  protocole  domaine       chemin      paramètre",
        "","Paramètres dans l'URL :","  https://api.exemple.com/data?page=2&limit=10","  ? sépare les paramètres, & les sépare entre eux"
    ])
    make_code_content(prs,"Première API avec requests",[
        "import requests",
        "", "# GET simple", "reponse = requests.get(\"https://api.github.com\")",
        "", "# Vérifier le statut", "print(reponse.status_code)  # 200 = OK", "# 200 = succès, 404 = pas trouvé, 500 = erreur serveur", "# 301/302 = redirection, 401 = non autorisé, 403 = interdit",
        "",
        "# Voir la réponse en texte brut", "print(reponse.text[:200])",
        "",
        "# Convertir en JSON (dictionnaire Python)", "data = reponse.json()", "print(data[\"current_user_url\"])",
        "",
        "# Voir les en-têtes", "print(reponse.headers[\"Content-Type\"])",
        "",
        "# Vérifier le succès", "if reponse.ok:", '    print("Tout va bien")', "else:", '    print(f"Erreur {reponse.status_code}")'
    ])
    make_code_content(prs,"API GitHub en pratique",[
        "import requests",
        "", "# Récupérer les infos d'un utilisateur GitHub", "user = requests.get(\"https://api.github.com/users/python\").json()",
        "print(f\"Login: {user['login']}\")", "print(f\"Nom: {user['name']}\")", "print(f\"Repos: {user['public_repos']}\")", "print(f\"Followers: {user['followers']}\")",
        "",
        "# Récupérer ses dépôts", "repos = requests.get(\"https://api.github.com/users/python/repos\").json()",
        "for repo in repos[:5]:", "    print(f\"- {repo['name']} : {repo['description']}\")",
        "",
        "# Avec paramètres", "params = {\"per_page\": 3, \"sort\": \"updated\"}", "repos = requests.get(\"https://api.github.com/users/python/repos\", params=params).json()",
        "for r in repos:", "    print(r[\"full_name\"])"
    ])
    make_content(prs,"APIs avec paramètres et authentification",[
        ("APIs avec clé (API Key) :",15),
        "","  # Beaucoup d'APIs nécessitent une clé (gratuite)", "  # Exemple avec OpenWeatherMap",
        "  API_KEY = \"ta_clé_ici\"", "  ville = \"Paris\"", "  url = \"https://api.openweathermap.org/data/2.5/weather\"",
        "  params = {\"q\": ville, \"appid\": API_KEY, \"units\": \"metric\"}", "  reponse = requests.get(url, params=params)",
        "  if reponse.status_code == 200:", "      data = reponse.json()", "      print(f\"Température: {data['main']['temp']}°C\")",
        "","","  💡 Jamais de clé API en dur dans le code !", "  Utiliser :", "    • Fichier .env (python-dotenv)", "    • Variables d'environnement : os.environ[\"API_KEY\"]", "    • .gitignore pour ne pas exposer la clé"
    ])
    make_code_content(prs,"POST, en-têtes et gestion d'erreurs",[
        "import requests",
        "", "# POST : envoyer des données", "data = {\"title\": \"Mon article\", \"body\": \"Contenu\"}", 'reponse = requests.post("https://jsonplaceholder.typicode.com/posts", json=data)',
        "print(reponse.status_code)  # 201 = Created", "print(reponse.json())  # l'objet créé avec son ID",
        "",
        "# En-têtes personnalisés", "headers = {\"Authorization\": \"Bearer token123\", \"User-Agent\": \"MonApp/1.0\"}", "reponse = requests.get(\"https://api.github.com/user\", headers=headers)",
        "",
        "# Gestion d'erreurs robuste", "try:", "    reponse = requests.get(\"https://api.github.com/users/inconnu\", timeout=5)", "    reponse.raise_for_status()  # lève exception si pas 2xx", "    print(reponse.json())", "except requests.exceptions.HTTPError as e:", "    print(f\"HTTP Error: {e}\")", "except requests.exceptions.ConnectionError:", "    print(\"Pas de connexion Internet\")", "except requests.exceptions.Timeout:", "    print(\"Timeout (serveur trop lent)\")", "except requests.exceptions.RequestException as e:", "    print(f\"Erreur: {e}\")"
    ])
    make_ex(prs,"Exercice : Explorer les APIs",[
        "1. Récupère les infos de l'utilisateur GitHub \"octocat\"\n2. Affiche : nom, bio, localisation, nb de followers\n3. Récupère ses 3 derniers dépôts\n4. Pour chaque dépôt affiche : nom, étoiles, langage principal\n\nBONUS :\n5. Utilise l'API https://api.chucknorris.io/ pour une blague aléatoire\n6. Utilise https://api.coindesk.com/v1/bpi/currentprice.json pour le prix du Bitcoin\n7. Crée une fonction qui prend un username GitHub et affiche son profil\n\nEXTRA BONUS :\n8. Utilise https://restcountries.com/v3.1/name/france pour les infos pays\n9. Gère les erreurs (si l'utilisateur n'existe pas)",
        ["import requests", "def profil_github(username):", '    url = f"https://api.github.com/users/{username}"', "    r = requests.get(url)", "    if r.status_code == 404:", '        return "Utilisateur inconnu"', "    u = r.json()", '    return f"{u[\"name\"]} - {u[\"bio\"]}"']
    ])
    make_summary(prs,["API = interface entre programmes via HTTP","GET = lire, POST = créer, PUT = maj, DELETE = supprimer","module requests : get(), post(), json(), status_code","Toujours gérer les erreurs (timeout, 404, connexion)","API Key = fichier .env + .gitignore","Le monde du développement est rempli d'APIs !"])
    save(prs,"42_APIs_HTTP.pptx")
    return prs

if __name__ == "__main__":
    funcs = [p35_exceptions_avancees, p36_mini_projet_jeu_cartes,
             p37_git_intro, p38_git_commandes, p39_git_branches_github,
             p40_mini_projet_pendu, p41_environnements_virtuels, p42_apis_http]
    for f in funcs:
        print(f"\n Génération : {f.__name__}")
        try:
            f()
        except Exception as e:
            print(f"  ❌ Erreur dans {f.__name__}: {e}")
    print(f"\n{'='*50}\n✅ {len(funcs)} présentations générées !\n📁 Dossier : {OUT}\n{'='*50}")
