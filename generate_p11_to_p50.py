#!/usr/bin/env python3
"""Generate presentations 11 to 50 - Python learning curriculum continued."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os, math

# ═══════ DESIGN SYSTEM (same as part 1) ═══════
class Theme:
    DARK=RGBColor(0x0A,0x0A,0x2E); PRIMARY=RGBColor(0x16,0x21,0x3E)
    SECONDARY=RGBColor(0x1A,0x1A,0x4E); ACCENT_OR=RGBColor(0xFF,0x6B,0x35)
    ACCENT_CY=RGBColor(0x00,0xB4,0xD8); ACCENT_PU=RGBColor(0x7B,0x2F,0xF7)
    ACCENT_GR=RGBColor(0x06,0xD6,0xA0); WHITE=RGBColor(0xFF,0xFF,0xFF)
    LIGHT=RGBColor(0xF8,0xF9,0xFA); GRAY=RGBColor(0x6C,0x75,0x7D)
    LIGHT_GRAY=RGBColor(0xDE,0xE2,0xE6); CODE_BG=RGBColor(0x1E,0x1E,0x2E)
    CODE_TEXT=RGBColor(0xCD,0xD6,0xF4)

OUTPUT_DIR = "/home/akaletekoffilevis/Bureau/Coach/presentations"
os.makedirs(OUTPUT_DIR, exist_ok=True)

def add_slide_bg(slide, color):
    bg=slide.background; f=bg.fill; f.solid(); f.fore_color.rgb=color

def add_shape(slide,l,t,w,h,fc,bc=None):
    s=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=fc
    if bc: s.line.color.rgb=bc; s.line.width=Pt(1)
    else: s.line.fill.background()
    return s

def add_rounded_rect(slide,l,t,w,h,fc):
    s=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=fc; s.line.fill.background()
    return s

def add_textbox(slide,l,t,w,h,text,fs=14,color=None,bold=False,align=PP_ALIGN.LEFT,fn='Calibri'):
    tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(fs); p.font.color.rgb=color or Theme.WHITE; p.font.bold=bold; p.font.name=fn; p.alignment=align
    return tb

def add_rich_textbox(slide,l,t,w,h):
    tb=slide.shapes.add_textbox(l,t,w,h); tf=tb.text_frame; tf.word_wrap=True; return tb

def add_para(tf,text,size=14,color=None,bold=False,name='Calibri',align=PP_ALIGN.LEFT,sb=0,sa=0,level=0):
    if len(tf.paragraphs)==1 and tf.paragraphs[0].text=='': p=tf.paragraphs[0]
    else: p=tf.add_paragraph()
    p.text=text; p.font.size=Pt(size); p.font.color.rgb=color or Theme.WHITE; p.font.bold=bold; p.font.name=name; p.alignment=align; p.space_before=Pt(sb); p.space_after=Pt(sa); p.level=level
    return p

def add_code_block(slide,l,t,w,h,lines,fs=11):
    add_rounded_rect(slide,l,t,w,h,Theme.CODE_BG)
    tb=slide.shapes.add_textbox(l+Inches(0.2),t+Inches(0.15),w-Inches(0.4),h-Inches(0.3))
    tf=tb.text_frame; tf.word_wrap=True; first=True
    for line in lines:
        if first: p=tf.paragraphs[0]; first=False
        else: p=tf.add_paragraph()
        p.text=line; p.font.name='Consolas'; p.font.size=Pt(fs); p.font.color.rgb=Theme.CODE_TEXT; p.space_before=Pt(0); p.space_after=Pt(0)

def make_title_slide(prs,title,subtitle,part=""):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.DARK)
    add_shape(s,Inches(0),Inches(0),prs.slide_width,Inches(0.06),Theme.ACCENT_OR)
    add_shape(s,Inches(0),Inches(0),Inches(0.06),prs.slide_height,Theme.ACCENT_CY)
    if part: add_textbox(s,Inches(1.5),Inches(2.0),Inches(8),Inches(0.5),part,16,Theme.ACCENT_OR,True)
    add_textbox(s,Inches(1.5),Inches(2.6),Inches(8),Inches(1.5),title,40,Theme.WHITE,True)
    add_textbox(s,Inches(1.5),Inches(4.2),Inches(8),Inches(1),subtitle,20,Theme.GRAY)
    add_shape(s,Inches(1.5),Inches(5.2),Inches(3),Inches(0.04),Theme.ACCENT_OR)
    add_textbox(s,Inches(1.5),Inches(6.5),Inches(8),Inches(0.5),"Cours de Programmation Python",12,Theme.GRAY)

def make_section_slide(prs,num,title,desc=""):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.PRIMARY)
    add_textbox(s,Inches(0.8),Inches(1.0),Inches(2),Inches(2),f"{num:02d}",72,Theme.ACCENT_OR,True)
    add_textbox(s,Inches(0.8),Inches(3.0),Inches(8.5),Inches(1.2),title,36,Theme.WHITE,True)
    add_shape(s,Inches(0.8),Inches(4.3),Inches(4),Inches(0.04),Theme.ACCENT_CY)
    if desc: add_textbox(s,Inches(0.8),Inches(4.6),Inches(8.5),Inches(1.5),desc,18,Theme.LIGHT_GRAY)
    add_shape(s,Inches(0),Inches(7.1),prs.slide_width,Inches(0.06),Theme.ACCENT_GR)

def make_content_slide(prs,title,bullets,notes=""):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.LIGHT)
    add_shape(s,Inches(0),Inches(0),prs.slide_width,Inches(0.9),Theme.SECONDARY)
    add_shape(s,Inches(0),Inches(0.9),prs.slide_width,Inches(0.04),Theme.ACCENT_OR)
    add_textbox(s,Inches(0.6),Inches(0.12),Inches(9),Inches(0.7),title,26,Theme.WHITE,True)
    tb=add_rich_textbox(s,Inches(0.6),Inches(1.3),Inches(8.8),Inches(5.5)); tf=tb.text_frame; first=True
    for b in bullets:
        if first: p=tf.paragraphs[0]; first=False
        else: p=tf.add_paragraph()
        if isinstance(b,tuple): text,lvl=b; p.level=lvl
        else: text=b; p.level=0
        p.text=text; p.font.name='Calibri'; p.font.color.rgb=Theme.DARK; p.space_before=Pt(2); p.space_after=Pt(2)
        if p.level==0: p.font.size=Pt(18); p.font.bold=True
        else: p.font.size=Pt(15); p.font.bold=False
    if notes: add_textbox(s,Inches(0.6),Inches(6.8),Inches(8.8),Inches(0.5),f" {notes}",11,Theme.GRAY)
    add_shape(s,Inches(0),Inches(7.35),prs.slide_width,Inches(0.04),Theme.ACCENT_CY)
    return s

def make_code_slide(prs,title,code_lines,explanation=""):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.LIGHT)
    add_shape(s,Inches(0),Inches(0),prs.slide_width,Inches(0.9),Theme.SECONDARY)
    add_shape(s,Inches(0),Inches(0.9),prs.slide_width,Inches(0.04),Theme.ACCENT_GR)
    add_textbox(s,Inches(0.6),Inches(0.12),Inches(9),Inches(0.7),f"  {title}",24,Theme.WHITE,True)
    add_code_block(s,Inches(0.6),Inches(1.2),Inches(8.8),Inches(4.5),code_lines,13)
    if explanation:
        tb=add_rich_textbox(s,Inches(0.6),Inches(5.9),Inches(8.8),Inches(1.2)); tf=tb.text_frame
        p=tf.paragraphs[0]; p.text=f" {explanation}"; p.font.size=Pt(13); p.font.color.rgb=Theme.GRAY; p.font.name='Calibri'
    add_shape(s,Inches(0),Inches(7.35),prs.slide_width,Inches(0.04),Theme.ACCENT_GR)
    return s

def make_exercise_slide(prs,title,ex_text,solution=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.LIGHT)
    add_shape(s,Inches(0),Inches(0),prs.slide_width,Inches(0.9),Theme.ACCENT_PU)
    add_shape(s,Inches(0),Inches(0.9),prs.slide_width,Inches(0.04),Theme.ACCENT_OR)
    add_textbox(s,Inches(0.6),Inches(0.12),Inches(9),Inches(0.7),f"  {title}",24,Theme.WHITE,True)
    tb=add_rich_textbox(s,Inches(0.6),Inches(1.3),Inches(8.8),Inches(2.5)); tf=tb.text_frame
    p=tf.paragraphs[0]; p.text=ex_text; p.font.size=Pt(16); p.font.color.rgb=Theme.DARK; p.font.name='Calibri'
    if solution:
        add_code_block(s,Inches(0.6),Inches(4.0),Inches(8.8),Inches(3.2),solution,11)
        add_textbox(s,Inches(0.6),Inches(3.7),Inches(3),Inches(0.3),"Solution :",12,Theme.ACCENT_PU,True)
    add_shape(s,Inches(0),Inches(7.35),prs.slide_width,Inches(0.04),Theme.ACCENT_PU)

def make_objectives_slide(prs,objectives):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.LIGHT)
    add_shape(s,Inches(0),Inches(0),prs.slide_width,Inches(0.9),Theme.SECONDARY)
    add_shape(s,Inches(0),Inches(0.9),prs.slide_width,Inches(0.04),Theme.ACCENT_OR)
    add_textbox(s,Inches(0.6),Inches(0.12),Inches(9),Inches(0.7),"  Objectifs",24,Theme.WHITE,True)
    tb=add_rich_textbox(s,Inches(0.6),Inches(1.3),Inches(8.8),Inches(5.5)); tf=tb.text_frame; first=True
    for o in objectives:
        if first: p=tf.paragraphs[0]; first=False
        else: p=tf.add_paragraph()
        p.text=f"✅  {o}"; p.font.size=Pt(17); p.font.color.rgb=Theme.DARK; p.font.name='Calibri'; p.space_before=Pt(8); p.space_after=Pt(4)
    add_shape(s,Inches(0),Inches(7.35),prs.slide_width,Inches(0.04),Theme.ACCENT_CY)

def make_summary_slide(prs,points):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.DARK)
    add_shape(s,Inches(0),Inches(0),prs.slide_width,Inches(0.04),Theme.ACCENT_OR)
    add_textbox(s,Inches(0.8),Inches(0.5),Inches(9),Inches(0.7),"  Ce qu'il faut retenir",28,Theme.WHITE,True)
    add_shape(s,Inches(0.8),Inches(1.2),Inches(3),Inches(0.04),Theme.ACCENT_CY)
    tb=add_rich_textbox(s,Inches(0.8),Inches(1.6),Inches(8.4),Inches(5.0)); tf=tb.text_frame; first=True
    for pt in points:
        if first: p=tf.paragraphs[0]; first=False
        else: p=tf.add_paragraph()
        p.text=f"✦  {pt}"; p.font.size=Pt(16); p.font.color.rgb=Theme.LIGHT_GRAY; p.font.name='Calibri'; p.space_before=Pt(6); p.space_after=Pt(4)
    add_shape(s,Inches(0),Inches(7.35),prs.slide_width,Inches(0.04),Theme.ACCENT_GR)

def make_two_col_slide(prs,title,lt,li,rt,ri):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.LIGHT)
    add_shape(s,Inches(0),Inches(0),prs.slide_width,Inches(0.9),Theme.SECONDARY)
    add_shape(s,Inches(0),Inches(0.9),prs.slide_width,Inches(0.04),Theme.ACCENT_OR)
    add_textbox(s,Inches(0.6),Inches(0.12),Inches(9),Inches(0.7),title,24,Theme.WHITE,True)
    add_shape(s,Inches(0.4),Inches(1.2),Inches(4.3),Inches(0.04),Theme.ACCENT_CY)
    add_textbox(s,Inches(0.4),Inches(1.3),Inches(4.3),Inches(0.4),lt,16,Theme.ACCENT_CY,True)
    tb=add_rich_textbox(s,Inches(0.4),Inches(1.8),Inches(4.3),Inches(5.0)); tf=tb.text_frame; first=True
    for it in li:
        if first: p=tf.paragraphs[0]; first=False
        else: p=tf.add_paragraph()
        p.text=it; p.font.size=Pt(14); p.font.color.rgb=Theme.DARK; p.font.name='Calibri'; p.space_before=Pt(3); p.space_after=Pt(2)
    add_shape(s,Inches(5.3),Inches(1.2),Inches(4.3),Inches(0.04),Theme.ACCENT_OR)
    add_textbox(s,Inches(5.3),Inches(1.3),Inches(4.3),Inches(0.4),rt,16,Theme.ACCENT_OR,True)
    tb=add_rich_textbox(s,Inches(5.3),Inches(1.8),Inches(4.3),Inches(5.0)); tf=tb.text_frame; first=True
    for it in ri:
        if first: p=tf.paragraphs[0]; first=False
        else: p=tf.add_paragraph()
        p.text=it; p.font.size=Pt(14); p.font.color.rgb=Theme.DARK; p.font.name='Calibri'; p.space_before=Pt(3); p.space_after=Pt(2)
    add_shape(s,Inches(0),Inches(7.35),prs.slide_width,Inches(0.04),Theme.ACCENT_GR)

def make_definition_slide(prs,term,definition,example=""):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_slide_bg(s,Theme.LIGHT)
    add_shape(s,Inches(0),Inches(0),prs.slide_width,Inches(0.9),Theme.SECONDARY)
    add_shape(s,Inches(0),Inches(0.9),prs.slide_width,Inches(0.04),Theme.ACCENT_OR)
    add_textbox(s,Inches(0.6),Inches(0.12),Inches(9),Inches(0.7),"  Définition",24,Theme.WHITE,True)
    add_textbox(s,Inches(0.6),Inches(1.5),Inches(8.8),Inches(0.8),term,30,Theme.ACCENT_OR,True)
    add_shape(s,Inches(0.6),Inches(2.3),Inches(2),Inches(0.03),Theme.ACCENT_CY)
    tb=add_rich_textbox(s,Inches(0.6),Inches(2.6),Inches(8.8),Inches(2.0)); tf=tb.text_frame
    p=tf.paragraphs[0]; p.text=definition; p.font.size=Pt(17); p.font.color.rgb=Theme.DARK; p.font.name='Calibri'
    if example:
        add_rounded_rect(s,Inches(0.6),Inches(4.8),Inches(8.8),Inches(2.2),RGBColor(0xE8,0xF4,0xFD))
        add_textbox(s,Inches(0.8),Inches(4.9),Inches(1.5),Inches(0.3),"Exemple :",13,Theme.ACCENT_CY,True)
        add_textbox(s,Inches(0.8),Inches(5.3),Inches(8.4),Inches(1.5),example,15,Theme.DARK)
    add_shape(s,Inches(0),Inches(7.35),prs.slide_width,Inches(0.04),Theme.ACCENT_CY)

def save_presentation(prs,filename):
    path=os.path.join(OUTPUT_DIR,filename); prs.save(path); print(f"  ✅ {filename}")

# ═══════════════════════════════════════════════════════════════
# PART 2 (continued): COMPLETE PYTHON BASICS
# ═══════════════════════════════════════════════════════════════

def p11_nombres_operations():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Nombres et Opérations Mathématiques", "Calculer avec Python", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Comprendre les types numériques (int, float)", "Maîtriser les opérateurs arithmétiques", "Connaître l'ordre des opérations", "Découvrir le module math"])
    make_content_slide(prs, "Les nombres en Python", [
        ("Deux types principaux :", 0), "int (integer) : nombres entiers → 42, -3, 0", "float (floating point) : nombres décimaux → 3.14, -0.5, 2.0",
        "", ("Différence clé :", 0), "Les entiers sont EXACTS, les flottants sont APPROCHÉS", "Exemple : 0.1 + 0.2 = 0.30000000000000004 (problème de précision)",
        "", ("Opérations de base :", 0), "a = 10; b = 3", "Addition : a + b = 13", "Soustraction : a - b = 7", "Multiplication : a * b = 30", "Division : a / b = 3.333... (résultat TOUJOURS en float)",
    ])
    make_code_slide(prs, "Les opérateurs arithmétiques", [
        "a = 10; b = 3", "", "# Opérations de base", "print(a + b)   # 13  (addition)", "print(a - b)   # 7   (soustraction)", "print(a * b)   # 30  (multiplication)", "print(a / b)   # 3.333... (division → float)",
        "", "# Opérations spéciales", "print(a // b)  # 3   (division ENTIÈRE)", "print(a % b)   # 1   (MODULO : reste de la division)", "print(a ** b)  # 1000 (PUISSANCE : 10³ = 10×10×10)",
        "", "# Raccourcis", "x = 5", "x += 3   # x = x + 3 → 8", "x -= 2   # x = x - 2 → 6", "x *= 4   # x = x * 4 → 24", "x /= 2   # x = x / 2 → 12.0",
    ], "Le modulo % est super utile pour savoir si un nombre est pair (n % 2 == 0)")
    make_content_slide(prs, "Ordre des opérations (PEMDAS)", [
        ("Comme en mathématiques, Python suit un ordre précis :", 0), "",
        "1. Parenthèses ( )", "2. Exposants **", "3. Multiplication *, Division /, //, Modulo %", "4. Addition +, Soustraction -",
        "", ("Exemple :", 0), "resultat = 2 + 3 * 4 ** 2", "          = 2 + 3 * 16    (d'abord 4² = 16)", "          = 2 + 48        (ensuite 3 × 16 = 48)", "          = 50            (enfin 2 + 48 = 50)",
        "", ("Avec parenthèses :", 0), "resultat = (2 + 3) * 4 ** 2", "          = 5 * 16 = 80",
        "", "En cas de doute, ajoute des parenthèses ! (et rend le code plus lisible)"
    ])
    make_code_slide(prs, "Conversion entre types (casting)", [
        "# Convertir en entier : int()", "print(int(3.14))     # 3  (tronque la partie décimale)", "print(int(\"42\"))    # 42 (string → int)", "print(int(True))     # 1",
        "", "# Convertir en décimal : float()", "print(float(3))      # 3.0", "print(float(\"3.14\")) # 3.14",
        "", "# Convertir en chaîne : str()", "print(str(42))       # \"42\"", "print(str(3.14))     # \"3.14\"", "print(str(True))     # \"True\"",
        "", "# Attention aux erreurs !", "int(\"Bonjour\")  # ValueError !"
    ], "Le casting est TRÈS important quand on utilise input() qui retourne toujours une chaîne")
    make_code_slide(prs, "Le module math", [
        "import math  # On importe le module math", "", "# Constantes", "print(math.pi)        # 3.141592653589793", "print(math.e)         # 2.718281828459045",
        "", "# Arrondis", "print(math.ceil(3.2)) # 4  (arrondit à l'entier supérieur)", "print(math.floor(3.8))# 3  (arrondit à l'entier inférieur)", "print(round(3.5))     # 4  (arrondi standard)",
        "", "# Fonctions mathématiques", "print(math.sqrt(16))  # 4.0  (racine carrée)", "print(math.pow(2, 3)) # 8.0", "print(math.fabs(-5))  # 5.0  (valeur absolue)", "print(math.factorial(5)) # 120",
    ], "import math donne accès à des dizaines de fonctions mathématiques avancées")
    make_exercise_slide(prs, "Exercice : Calculs et conversions",
        "1. Calcule et affiche l'aire d'un cercle de rayon 5 (π × r²)\n2. Calcule et affiche : (15 + 3) × 4 ÷ 2 - 1\n3. Demande un nombre à l'utilisateur et affiche son double\n4. Prends deux nombres, affiche : somme, différence, produit, quotient\n5. Vérifie si un nombre est pair ou impair (utilise le %)\n\nBONUS : Affiche les 10 premiers nombres avec leur carré et leur cube",
        ["import math", "rayon = 5", "aire = math.pi * rayon ** 2", "print(f\"Aire du cercle : {aire:.2f}\")", "n = int(input(\"Nombre : \"))", "print(f\"Le double est {n * 2}\")", "print(f\"Pair\" if n % 2 == 0 else \"Impair\")"]
    )
    make_summary_slide(prs, ["int = entiers, float = décimaux", "/ donne float, // division entière, % modulo, ** puissance", "Ordre : Parenthèses → Exposants → ×/÷ → +-", "int(), float(), str() convertissent entre types", "import math pour des fonctions mathématiques avancées"])
    save_presentation(prs, "11_Nombres_operations.pptx")
    return prs

def p12_input_output():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Input / Output", "Interagir avec l'utilisateur", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Maîtriser input() pour récupérer des saisies", "Formater l'affichage avec print()", "Comprendre la conversion des types", "Créer des programmes interactifs"])
    make_content_slide(prs, "input() - Récupérer une saisie", [
        ("La fonction input() permet de demander quelque chose à l'utilisateur :", 0), "",
        "prenom = input(\"Quel est ton prénom ? \")", "print(f\"Bonjour {prenom} !\")",
        "", ("Ce qui se passe :", 0), "1. Le programme affiche le message (optionnel)", "2. Il attend que l'utilisateur tape quelque chose", "3. Il appuie sur Entrée", "4. input() retourne ce qui a été tapé (en chaîne)",
        "", ("CRUCIAL : input() retourne TOUJOURS une chaîne (str)", 0), "Même si tu tapes un nombre, c'est une chaîne !", "Il faut CONVERTIR avec int() ou float() pour faire des calculs"
    ])
    make_code_slide(prs, "Input avec conversion", [
        "# Problème : input retourne toujours du texte", "age = input(\"Quel est ton âge ? \")", "print(age + 1)   # ERREUR : str + int impossible",
        "", "# Solution : convertir avec int() ou float()", "age = int(input(\"Quel est ton âge ? \"))", "print(f\"L'an prochain, tu auras {age + 1} ans\")",
        "", "# Pour les nombres décimaux", "note = float(input(\"Ta note sur 20 : \"))", "print(f\"Tu as {note * 5} sur 100\")",
        "", "# Gérer les erreurs de saisie", "try:", "    age = int(input(\"Âge : \"))", "except ValueError:", "    print(\"Ce n'est pas un nombre valide !\")"
    ], "Toujours penser à convertir ! input() → str")
    make_code_slide(prs, "print() avancé", [
        "# print peut accepter plusieurs arguments", "", "# Séparateur personnalisé (sep)", "print(\"a\", \"b\", \"c\")                    # a b c", "print(\"a\", \"b\", \"c\", sep=\"-\")          # a-b-c", "print(\"a\", \"b\", \"c\", sep=\", \")          # a, b, c",
        "", "# Fin de ligne personnalisée (end)", "print(\"Bonjour\", end=\" \")", "print(\"Lucas\")                       # Bonjour Lucas (sur la même ligne)",
        "", "# Saut de ligne dans le texte", "print(\"Ligne 1\\nLigne 2\\nLigne 3\")",
    ], "sep et end sont très utiles pour formater joliment l'affichage")
    make_code_slide(prs, "Les f-strings avancées", [
        "# Formatage avancé avec les f-strings", "prix = 49.5", "taux = 0.076",
        "", "# Formatage des nombres décimaux", "print(f\"Prix : {prix:.2f} €\")        # Prix : 49.50 €", "print(f\"Taux : {taux:.1%}\")          # Taux : 7.6%",
        "", "# Alignement", "print(f\"|{'gauche':<10}|\")  # |gauche    |", "print(f\"|{'droite':>10}|\")  # |    droite|", "print(f\"|{'centre':^10}|\")  # |  centre  |",
        "", "# Grands nombres", "print(f\"{1000000:,}\")        # 1,000,000", "print(f\"{1000000:_}\")        # 1_000_000",
        "", "# Expressions dans les f-strings", "x, y = 5, 3", "print(f\"{x} + {y} = {x + y}\")  # 5 + 3 = 8"
    ], "Les f-strings sont extrêmement puissantes pour formater l'affichage proprement")
    make_exercise_slide(prs, "Exercice : Questionnaire interactif",
        "Crée un programme qui :\n1. Demande le prénom et l'âge de l'utilisateur\n2. Demande sa couleur préférée\n3. Demande un nombre\n4. Affiche un récapitulatif formaté :\n   \"Bonjour [Prénom] ! Tu as [âge] ans.\"\n5. Calcule l'année de naissance\n\nBONUS : Ajoute un cadre décoratif autour du message",
        ["prenom = input(\"Prénom : \")", "age = int(input(\"Âge : \"))", "couleur = input(\"Couleur préférée : \")", "nb = float(input(\"Un nombre : \"))", "annee = 2025 - age", "print(\"=\" * 40)", "print(f\"Bonjour {prenom} !\")", "print(f\"Tu as {age} ans (né(e) en {annee})\")", "print(f\"Couleur : {couleur}\")", "print(f\"Double de {nb} = {nb * 2}\")", "print(\"=\" * 40)"]
    )
    make_summary_slide(prs, ["input() récupère la saisie utilisateur (toujours en str)", "int(input()) ou float(input()) pour les nombres", "print() accepte sep, end pour personnaliser l'affichage", "Les f-strings (f\"...\") formatent joliment", "Les programmes deviennent INTERACTIFS grâce à input() !"])
    save_presentation(prs, "12_Input_Output.pptx")
    return prs

def p13_conditions():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Les Conditions (if/elif/else)", "Prendre des décisions dans le code", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Comprendre les conditions en programmation", "Maîtriser if, elif, else", "Utiliser les opérateurs de comparaison", "Combiner des conditions avec and, or, not"])
    make_definition_slide(prs, "Une condition", "Une condition permet d'exécuter du code SEULEMENT si quelque chose est vrai. C'est comme un embranchement : selon la réponse, on prend un chemin différent.", "SI il pleut ALORS je prends un parapluie\nSINON je ne prends rien")
    make_code_slide(prs, "Syntaxe de base : if", [
        "# Structure de base", "age = 16", "", "if age >= 18:", "    print(\"Tu es majeur\")", "else:", "    print(\"Tu es mineur\")",
        "", "# Règles IMPORTANTES :", "# 1. Le : (deux-points) est OBLIGATOIRE après if et else", "# 2. Le bloc indenté (4 espaces) est le corps du if/else", "# 3. Python est très strict sur l'indentation !"
    ], "L'indentation (4 espaces) définit les blocs en Python")
    make_code_slide(prs, "elif : plusieurs cas possibles", [
        "# Avec elif pour plusieurs conditions", "note = 15", "", "if note >= 18:", "    print(\"Excellent !\")", "elif note >= 15:", "    print(\"Très bien\")", "elif note >= 12:", "    print(\"Bien\")", "elif note >= 10:", "    print(\"Passable\")", "else:", "    print(\"Insuffisant\")",
        "", "# Note : elif peut être utilisé autant de fois qu'on veut", "# Le else est TOUJOURS optionnel à la fin"
    ], "Les conditions sont vérifiées dans L'ORDRE. Dès qu'une est vraie, on sort !")
    make_content_slide(prs, "Les opérateurs de comparaison", [
        ("a == b : a est égal à b", 1), ("a != b : a est différent de b", 1), ("a < b : a est strictement inférieur à b", 1), ("a <= b : a est inférieur ou égal à b", 1), ("a > b : a est strictement supérieur à b", 1), ("a >= b : a est supérieur ou égal à b", 1),
        "", ("ATTENTION : ne pas confondre :", 0), "=  → affectation (met une valeur dans une variable)", "== → comparaison (vérifie l'égalité)",
        "", ("Exemple :", 0), "x = 5       # affectation : x prend la valeur 5", "x == 5      # comparaison : est-ce que x vaut 5 ? (True)"
    ])
    make_code_slide(prs, "Conditions combinées : and, or, not", [
        "# and : les DEUX conditions doivent être vraies", "age = 20", "if age >= 18 and age <= 25:", "    print(\"Tu es dans la tranche 18-25 ans\")",
        "", "# or : AU MOINS UNE condition doit être vraie", "jour = \"Samedi\"", "if jour == \"Samedi\" or jour == \"Dimanche\":", "    print(\"C'est le week-end !\")",
        "", "# not : INVERSE la condition", "est_pluvieux = False", "if not est_pluvieux:", "    print(\"Pas besoin de parapluie\")",
        "", "# Combinaison complexe", "age = 17; avec_parent = True", "if (age >= 18) or (age >= 16 and avec_parent):", "    print(\"Tu peux entrer\")"
    ], "Utilise des parenthèses pour grouper les conditions")
    make_content_slide(prs, "Valeurs truthy/falsy", [
        ("En Python, TOUTE valeur peut être évaluée comme True ou False", 0), "",
        ("Valeurs FALSY (évaluées à False) :", 0), "False, 0, 0.0, \"\" (chaîne vide), [] (liste vide)", "{} (dictionnaire vide), None, () (tuple vide)",
        "", ("Tout le reste est TRUTHY (évalué à True)", 0), "",
        ("Exemple :", 0), "nom = input(\"Ton nom : \")", "if nom:   # équivalent à : if nom != \"\"", "    print(f\"Bonjour {nom}\")", "else:", "    print(\"Tu n'as pas entré de nom !\")",
    ], "Cette syntaxe est très Pythonique et concise !")
    make_code_slide(prs, "Opérateur ternaire", [
        "# Opérateur ternaire : condition en une ligne", "# Syntaxe : valeur_si_vrai if condition else valeur_si_faux",
        "", "age = 16", "# Version classique", "if age >= 18:", "    statut = \"Majeur\"", "else:", "    statut = \"Mineur\"",
        "", "# Version ternaire (équivalente, plus courte)", "statut = \"Majeur\" if age >= 18 else \"Mineur\"", "print(statut)  # Mineur",
        "", "# Utile pour des assignations simples", "x = 10", "resultat = \"Pair\" if x % 2 == 0 else \"Impair\"", "print(f\"{x} est {resultat}\")  # 10 est Pair"
    ], "Le ternaire est pratique pour les conditions SIMPLES")
    make_exercise_slide(prs, "Exercice : Jeu de combat",
        "Crée un programme de combat simplifié :\n1. Le joueur a 100 points de vie\n2. Demande : \"Quelle attaque ? (1=épée 15dégâts, 2=magie 25, 3=arc 10)\"\n3. Affiche les dégâts et la vie restante\n4. Si vie ≤ 0 → \"Game Over\"\n5. Si vie > 50 → \"En forme !\"\n6. Si vie entre 20 et 50 → \"Blessé\"\n7. Si vie < 20 → \"En danger !\"\n\nBONUS : Ajoute un bouclier qui réduit les dégâts",
        ["vie = 100", "attaque = int(input(\"Attaque (1/2/3) : \"))", "degats = 15 if attaque == 1 else 25 if attaque == 2 else 10", "vie -= degats", "print(f\"Dégâts : {degats}, Vie : {vie}\")", "if vie <= 0: print(\" Game Over\")", "elif vie > 50: print(\" En forme !\")", "elif vie >= 20: print(\" Blessé\")", "else: print(\" En danger !\")"]
    )
    make_summary_slide(prs, ["if/elif/else selon des conditions", "Les deux-points : et l'indentation (4 espaces) sont OBLIGATOIRES", "Opérateurs : ==, !=, <, <=, >, >=, and, or, not", "Valeurs falsy : 0, \"\", [], None, False...", "Les conditions rendent le code INTELLIGENT !"])
    save_presentation(prs, "13_Conditions.pptx")
    return prs

def p14_exercices_conditions():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Exercices : Conditions", "Mets en pratique les conditions !", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Consolider les notions de conditions", "Pratiquer if/elif/else", "Utiliser des conditions combinées", "Créer des programmes intéressants"])
    make_exercise_slide(prs, "Exercice 1 : Nombre pair ou impair",
        "Écris un programme qui :\n1. Demande un nombre à l'utilisateur\n2. Affiche si le nombre est pair ou impair\n3. Affiche si le nombre est positif, négatif ou nul\n4. Affiche si le nombre est un multiple de 3\n\nAstuce : n % 2 == 0 → pair, n % 3 == 0 → multiple de 3",
        ["n = int(input(\"Nombre : \"))", "if n % 2 == 0: print(f\"{n} est pair\")", "else: print(f\"{n} est impair\")", "if n > 0: print(\"Positif\")", "elif n < 0: print(\"Négatif\")", "else: print(\"Nul\")", "if n % 3 == 0: print(f\"{n} est multiple de 3\")"]
    )
    make_exercise_slide(prs, "Exercice 2 : Calculatrice conditionnelle",
        "Écris une mini-calculatrice :\n1. Demande deux nombres (a et b)\n2. Demande l'opération (+, -, *, /)\n3. Affiche le résultat\n4. Gère l'erreur : division par zéro → message d'erreur\n5. Gère l'opération inconnue → message d'erreur\n\nBONUS : Ajoute la puissance (^) et le modulo (%)",
        ["a = float(input(\"a : \"))", "b = float(input(\"b : \"))", "op = input(\"Opération (+, -, *, /) : \")", "if op == '+': print(f\"{a}+{b}={a+b}\")", "elif op == '-': print(f\"{a}-{b}={a-b}\")", "elif op == '*': print(f\"{a}×{b}={a*b}\")", "elif op == '/':", "    print(f\"{a}÷{b}={a/b}\") if b!=0 else print(\"Division par zéro !\")", "else: print(\"Opération invalide\")"]
    )
    make_exercise_slide(prs, "Exercice 3 : Convertisseur de notes",
        "Crée un convertisseur de notes :\n1. Demande une note sur 20\n2. Affiche la mention correspondante :\n   18-20 → Excellent | 15-17 → Très Bien\n   12-14 → Bien | 10-11 → Passable\n   8-9 → Insuffisant | 0-7 → Médiocre\n3. Vérifie que la note est entre 0 et 20\n4. Affiche aussi la note en lettre (A, B, C, D, E, F)",
        ["note = float(input(\"Note sur 20 : \"))", "if note < 0 or note > 20: print(\"Invalide\")", "elif note >= 18: print(f\"{note}/20 → Excellent (A)\")", "elif note >= 15: print(f\"Très Bien (B)\")", "elif note >= 12: print(f\"Bien (C)\")", "elif note >= 10: print(f\"Passable (D)\")", "elif note >= 8: print(f\"Insuffisant (E)\")", "else: print(f\"Médiocre (F)\")"]
    )
    make_exercise_slide(prs, "Exercice 4 : Année bissextile",
        "Détermine si une année est bissextile :\nRègles :\n- Divisible par 4 → bissextile\n- SAUF si divisible par 100 → pas bissextile\n- SAUF si divisible par 400 → bissextile à nouveau\n\nExemples : 2024 → oui, 2100 → non, 2000 → oui",
        ["annee = int(input(\"Année : \"))", "if annee % 400 == 0 or (annee % 4 == 0 and annee % 100 != 0):", "    print(f\"{annee} est bissextile\")", "else:", "    print(f\"{annee} n'est pas bissextile\")"]
    )
    make_exercise_slide(prs, "Exercice 5 : Pierre-Feuille-Ciseaux",
        "Crée un jeu :\n1. Le joueur choisit : 1=Pierre, 2=Feuille, 3=Ciseaux\n2. L'ordinateur choisit aléatoirement (random.randint)\n3. Compare les choix :\n   - Pierre bat Ciseaux\n   - Ciseaux bat Feuille\n   - Feuille bat Pierre\n4. Affiche le gagnant et les choix",
        ["import random", "choix = int(input(\"1=Pierre 2=Feuille 3=Ciseaux : \"))", "ordi = random.randint(1, 3)", "noms = [\"\", \"Pierre\", \"Feuille\", \"Ciseaux\"]", "print(f\"Ordi : {noms[ordi]}\")", "if choix == ordi: print(\"Égalité !\")", "elif (choix==1 and ordi==3) or (choix==2 and ordi==1) or (choix==3 and ordi==2):", "    print(\" Tu as gagné !\")", "else: print(\" L'ordi gagne...\")"]
    )
    make_summary_slide(prs, ["Les conditions permettent de prendre des décisions", "if/elif/else, : et indentation obligatoires", "and, or, not combinent plusieurs conditions", "Les conditions sont la base de TOUT programme interactif", "Le prochain chapitre : les boucles !"])
    save_presentation(prs, "14_Exercices_conditions.pptx")
    return prs

def p15_boucle_while():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "La Boucle while", "Répéter tant qu'une condition est vraie", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Comprendre le concept de boucle", "Maîtriser la boucle while", "Éviter les boucles infinies", "Utiliser break et continue"])
    make_content_slide(prs, "Pourquoi des boucles ?", [
        ("Sans boucle, on répète le code :", 0), "print(1); print(2); print(3); print(4); print(5)... → bof",
        "", ("Avec une boucle :", 0), "compteur = 1", "while compteur <= 5:", "    print(compteur)", "    compteur += 1",
        "", ("La boucle while répète un bloc TANT QUE la condition est vraie.", 0),
        "", ("Clé : il faut que la condition devienne fausse à un moment", 0), "Sinon → boucle INFINIE (le programme ne s'arrête jamais)"
    ])
    make_code_slide(prs, "Syntaxe et utilisation de while", [
        "# Syntaxe :", "while condition:", "    # bloc à répéter",
        "", "# Exemple 1 : Compter de 1 à 5", "i = 1", "while i <= 5:", "    print(i, end=\" \")", "    i += 1", "# Résultat : 1 2 3 4 5",
        "", "# Exemple 2 : Compter de 10 à 0", "i = 10", "while i >= 0:", "    print(i, end=\" \")", "    i -= 1", "# Résultat : 10 9 8 7 6 5 4 3 2 1 0"
    ], "La variable de condition DOIT changer dans la boucle !")
    make_code_slide(prs, "Applications concrètes de while", [
        "# 1. Saisie jusqu'à réponse correcte", "mdp = \"python\"; saisie = \"\"", "while saisie != mdp:", "    saisie = input(\"Mot de passe : \")", "print(\"Accès autorisé !\")",
        "", "# 2. Validation de saisie", "nombre = -1", "while nombre < 1 or nombre > 10:", "    nombre = int(input(\"Entre 1-10 : \"))",
        "", "# 3. Menu interactif", "choix = 0", "while choix != 4:", "    print(\"1.Jouer 2.Scores 3.Options 4.Quitter\")", "    choix = int(input(\"Choix : \"))"
    ], "while est PARFAIT quand on ne sait pas à l'avance combien de fois répéter")
    make_content_slide(prs, "Boucle infinie - Attention !", [
        ("BOUCLE INFINIE = programme qui ne s'arrête JAMAIS", 0), "",
        ("Ce code ne s'arrêtera jamais :", 0), "while True:", "    print(\"Je tourne en boucle !\")",
        "", ("Pourquoi ? La condition est TOUJOURS True", 0),
        "", ("Cas qui plante :", 0), "i = 1", "while i > 0:", "    i += 1  # i augmente → reste > 0 → INFINI !",
        "", ("Comment arrêter ?", 0), "Ctrl+C dans le terminal (arrêt d'urgence)"
    ])
    make_code_slide(prs, "break et continue", [
        "# break : SORT immédiatement de la boucle", "while True:", "    reponse = input(\"Tape 'quit' pour quitter : \")", "    if reponse == 'quit':", "        print(\"Au revoir !\")", "        break",
        "", "# continue : PASSE à l'itération suivante", "i = 0", "while i < 10:", "    i += 1", "    if i % 2 == 0:", "        continue  # saute les pairs", "    print(i, end=\" \")", "# Résultat : 1 3 5 7 9"
    ], "break = sortir de la boucle, continue = sauter au prochain tour")
    make_exercise_slide(prs, "Exercice : Devine le nombre",
        "Crée un jeu \"Devine le nombre\" :\n1. Python choisit un nombre aléatoire entre 1 et 100\n2. Le joueur propose des nombres (while True)\n3. Indique \"Trop grand\" ou \"Trop petit\"\n4. Quand il trouve : affiche le nombre d'essais\n5. Propose de rejouer\n\nBONUS : Limite le nombre d'essais à 7",
        ["import random", "secret = random.randint(1, 100); essais = 0", "while True:", "    prop = int(input(\"Devine (1-100) : \"))", "    essais += 1", "    if prop < secret: print(\"Trop petit\")", "    elif prop > secret: print(\"Trop grand\")", "    else:", "        print(f\" Trouvé en {essais} essais !\")", "        if input(\"Rejouer ? (o/n) : \") != 'o': break", "        secret = random.randint(1, 100); essais = 0"]
    )
    make_summary_slide(prs, ["while répète TANT QUE la condition est vraie", "La variable de condition DOIT changer (sinon : boucle infinie)", "Ctrl+C pour arrêter d'urgence", "break = sortir de la boucle, continue = sauter au tour suivant", "while est parfait pour validation de saisie, menus, jeux"])
    save_presentation(prs, "15_Boucle_while.pptx")
    return prs

def p16_boucle_for():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "La Boucle for", "Parcourir des séquences facilement", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Comprendre la boucle for", "Maîtriser range()", "Parcourir des chaînes, listes", "Utiliser la boucle for avec break/continue"])
    make_content_slide(prs, "for : la boucle pour", [
        ("La boucle for parcourt les éléments d'une SEQUENCE un par un.", 0), "",
        "for variable in sequence:", "    # bloc exécuté pour chaque élément",
        "", ("Contrairement à while, for sait combien de fois répéter :", 0), "C'est le nombre d'éléments dans la séquence !",
        "", ("Exemple avec une chaîne :", 0), "for lettre in \"Python\":", "    print(lettre)  # P y t h o n",
        "", ("Exemple avec une liste :", 0), "for fruit in [\"pomme\", \"banane\", \"cerise\"]:", "    print(fruit)"
    ])
    make_code_slide(prs, "range() - générer des séquences de nombres", [
        "# range(n) : de 0 à n-1", "for i in range(5):", "    print(i, end=\" \")  # 0 1 2 3 4",
        "", "# range(debut, fin) : de debut à fin-1", "for i in range(3, 8):", "    print(i, end=\" \")  # 3 4 5 6 7",
        "", "# range(debut, fin, pas) : avec un pas", "for i in range(0, 10, 2):", "    print(i, end=\" \")  # 0 2 4 6 8",
        "", "# range décroissant", "for i in range(10, 0, -1):", "    print(i, end=\" \")  # 10 9 8 7 6 5 4 3 2 1",
        "", "# Répéter N fois (sans utiliser la variable)", "for _ in range(3):", "    print(\"Python c'est génial !\")"
    ], "_ (underscore) est une convention pour une variable qu'on n'utilise pas")
    make_code_slide(prs, "Parcourir différents types de séquences", [
        "# Parcourir une chaîne", "for lettre in \"Python\":", "    print(lettre.upper(), end=\" \")  # P Y T H O N",
        "", "# Parcourir une liste", "amis = [\"Alice\", \"Bob\", \"Charlie\"]", "for ami in amis:", "    print(f\"Salut {ami} !\")",
        "", "# enumerate() : obtenir index ET valeur", "for index, ami in enumerate(amis):", "    print(f\"{index+1}. {ami}\")",
        "", "# zip() : parcourir plusieurs listes", "noms = [\"Alice\", \"Bob\"]; ages = [20, 25]", "for nom, age in zip(noms, ages):", "    print(f\"{nom} a {age} ans\")"
    ], "enumerate() et zip() sont très pratiques et Pythoniques !")
    make_code_slide(prs, "for avec break et continue", [
        "# break : arrête la boucle", "for i in range(1, 20):", "    if i == 10: break", "    print(i, end=\" \")  # 1 2 3 4 5 6 7 8 9",
        "", "# continue : saute une itération", "for i in range(1, 11):", "    if i % 3 == 0: continue", "    print(i, end=\" \")  # 1 2 4 5 7 8 10",
        "", "# else dans une boucle for (exécuté si pas de break)", "for i in range(5):", "    if i == 10: break", "else:", "    print(\"Aucun break, fin normale\")"
    ], "Le else d'une boucle s'exécute seulement si pas de break")
    make_exercise_slide(prs, "Exercice : Tables de multiplication",
        "Crée un programme qui :\n1. Demande un nombre à l'utilisateur\n2. Affiche la table de multiplication de 1 à 10\n3. Utilise une boucle for\n4. Formate joliment\n\nBONUS : Affiche les tables de 1 à 10 complètes (tableau)\navec deux boucles for imbriquées",
        ["n = int(input(\"Table de multiplication : \"))", "for i in range(1, 11):", "    print(f\"{n} × {i:>2} = {n * i:>3}\")", "", "# Tableau complet", "for i in range(1, 11):", "    for j in range(1, 11):", "        print(f\"{i*j:>4}\", end=\"\")", "    print()"]
    )
    make_summary_slide(prs, ["for parcourt les éléments d'une séquence", "range() génère des séquences de nombres", "enumerate() donne (index, valeur), zip() parcourt plusieurs listes", "break = arrêt, continue = saut, else = si pas de break", "for est plus sûr que while (pas de risque de boucle infinie)"])
    save_presentation(prs, "16_Boucle_for.pptx")
    return prs

def p17_listes():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Les Listes (list)", "Stocker plusieurs valeurs dans une seule variable", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Comprendre ce qu'est une liste", "Créer et manipuler des listes", "Utiliser les méthodes des listes", "Parcourir et modifier des listes"])
    make_content_slide(prs, "Qu'est-ce qu'une liste ?", [
        ("Une liste est une collection ORDONNÉE et MODIFIABLE d'éléments.", 0), "",
        ("Création :", 0), "liste_vide = []", "nombres = [1, 2, 3, 4, 5]", "mixte = [42, \"Python\", True, 3.14]",
        "", ("Caractéristiques :", 0), "✅ Ordonnée : les éléments ont un ordre défini", "✅ Modifiable : on peut ajouter, supprimer, modifier", "✅ Hétérogène : différents types dans la même liste", "✅ Indexée : on accède par position (index)", "✅ Accepte les doublons"
    ])
    make_code_slide(prs, "Accéder aux éléments (indexation)", [
        "fruits = [\"pomme\", \"banane\", \"cerise\", \"datte\", \"figue\"]", "", "# Index positif (0 = premier)", "print(fruits[0])    # pomme", "print(fruits[2])    # cerise", "print(fruits[-1])   # figue (dernier)", "print(fruits[-2])   # datte",
        "", "# Slicing (tranches)", "print(fruits[1:3])   # ['banane', 'cerise']", "print(fruits[:3])    # ['pomme', 'banane', 'cerise']", "print(fruits[2:])    # ['cerise', 'datte', 'figue']", "print(fruits[::-1])  # inversé",
        "", "# Modifier un élément", "fruits[1] = \"mangue\"", "print(fruits)  # ['pomme', 'mangue', 'cerise', 'datte', 'figue']"
    ], "L'indexation et le slicing sont les mêmes que pour les chaînes !")
    make_code_slide(prs, "Méthodes principales des listes", [
        "fruits = [\"pomme\", \"banane\"]", "", "# AJOUTER", "fruits.append(\"cerise\")    # à la fin", "fruits.insert(1, \"mangue\") # à l'index 1", "fruits.extend([\"kiwi\", \"datte\"])  # plusieurs à la fin",
        "", "# SUPPRIMER", "fruits.remove(\"banane\")    # supprime le premier trouvé", "dernier = fruits.pop()     # supprime et retourne le dernier", "del fruits[1]              # supprime l'index 1",
        "", "# TRIER", "nombres = [3, 1, 4, 1, 5, 9]", "nombres.sort()            # [1, 1, 3, 4, 5, 9] (sur place)", "nombres.sort(reverse=True) # [9, 5, 4, 3, 1, 1]", "liste_triee = sorted(nombres)  # nouvelle liste"
    ], "sort() modifie sur place, sorted() retourne une nouvelle liste")
    make_code_slide(prs, "Parcourir et chercher dans une liste", [
        "nombres = [10, 20, 30, 40, 50]", "", "# 1. Parcourir les éléments", "for n in nombres:", "    print(n, end=\" \")  # 10 20 30 40 50",
        "", "# 2. enumerate()", "for i, n in enumerate(nombres):", "    print(f\"Index {i} : {n}\")",
        "", "# 3. Vérifier la présence", "if 30 in nombres:", "    print(\"30 est dans la liste !\")",
        "", "# 4. Trouver l'index", "print(nombres.index(30))  # 2",
        "", "# 5. Longueur, min, max, somme", "print(len(nombres))     # 5", "print(min(nombres))     # 10", "print(max(nombres))     # 50", "print(sum(nombres))     # 150"
    ], "len(), min(), max(), sum() sont des fonctions intégrées très utiles")
    make_content_slide(prs, "Copie de listes - Attention !", [
        ("PIÈGE CLASSIQUE : copier une liste", 0), "",
        ("Mauvaise façon :", 0), "liste1 = [1, 2, 3]", "liste2 = liste1          # NON ! Ce n'est PAS une copie !", "liste2.append(4)", "print(liste1)  # [1, 2, 3, 4]  ← modifié aussi !",
        "", ("Pourquoi ? liste1 et liste2 pointent vers LA MÊME liste.", 0),
        "", ("Bonnes façons de copier :", 0), "liste2 = liste1.copy()", "liste2 = list(liste1)", "liste2 = liste1[:]",
        "", ("Pour les listes imbriquées :", 0), "import copy", "liste2 = copy.deepcopy(liste1)"
    ])
    make_exercise_slide(prs, "Exercice : Gestion de liste de notes",
        "Crée un programme qui gère une liste de notes :\n1. Crée une liste avec quelques notes (sur 20)\n2. Affiche : nombre, note min, note max, moyenne\n3. Ajoute une nouvelle note saisie par l'utilisateur\n4. Supprime la première note\n5. Affiche toutes les notes triées du meilleur au pire\n6. Affiche les notes au-dessus de 10\n\nBONUS : Calcule la médiane",
        ["notes = [12, 15, 8, 19, 10, 14, 6]", "print(f\"Notes : {notes}\")", "print(f\"Nb : {len(notes)}\")", "print(f\"Min: {min(notes)}, Max: {max(notes)}\")", "print(f\"Moyenne: {sum(notes)/len(notes):.2f}\")", "notes.append(int(input(\"Nouvelle note : \")))", "notes.pop(0)", "print(f\"Triées: {sorted(notes, reverse=True)}\")", "print(f\"+10 : {[n for n in notes if n > 10]}\")"]
    )
    make_summary_slide(prs, ["Une liste est une collection ORDONNÉE et MODIFIABLE", "Accès par index (0, 1, 2...) ou slicing ([:])", "Méthodes : append, insert, remove, pop, sort, clear", "len(), min(), max(), sum() sont des fonctions intégrées", "Attention : l'assignation (list2 = list1) ne copie PAS !"])
    save_presentation(prs, "17_Listes.pptx")
    return prs

def p18_dictionnaires():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Les Dictionnaires (dict)", "Stocker des données avec des clés", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Comprendre le concept clé-valeur", "Créer et manipuler des dictionnaires", "Utiliser les méthodes des dict", "Savoir quand utiliser une liste vs un dict"])
    make_content_slide(prs, "Qu'est-ce qu'un dictionnaire ?", [
        ("Un dictionnaire associe des CLÉS à des VALEURS.", 0), "",
        "C'est comme un vrai dictionnaire : mot (clé) → définition (valeur)",
        "",         ("Création :", 0), "vide = {}", "personne = {\"nom\": \"Utilisateur\", \"age\": 25, \"ville\": \"Paris\"}",
        "", ("Caractéristiques :", 0), "✅ Les clés sont UNIQUES (pas de doublon)", "✅ Accès rapide par la clé", "✅ Les valeurs peuvent être de n'importe quel type"
    ])
    make_code_slide(prs, "Créer et accéder à un dictionnaire", [
        "# Création", "etudiant = {", "    \"nom\": \"Utilisateur\",", "    \"age\": 25,", "    \"classe\": \"3ème\",", "    \"notes\": [12, 15, 8, 19]", "}",
        "", "# Accès par la clé", "print(etudiant[\"nom\"])     # Lucas", "print(etudiant[\"notes\"])   # [12, 15, 8, 19]",
        "", "# Accès avec get() (plus sûr)", "print(etudiant.get(\"nom\"))           # Lucas", "print(etudiant.get(\"pays\"))          # None (pas d'erreur !)", "print(etudiant.get(\"pays\", \"Inconnu\"))  # Inconnu",
        "", "# Modifier une valeur", "etudiant[\"age\"] = 15", "# Ajouter une nouvelle clé", "etudiant[\"email\"] = \"lucas@email.com\"",
        "", "# Vérifier si une clé existe", "if \"nom\" in etudiant:", "    print(\"La clé 'nom' existe !\")"
    ], "get() est plus sûr que [] car il ne lève pas d'erreur si la clé manque")
    make_code_slide(prs, "Méthodes principales des dictionnaires", [
        "personne = {\"nom\": \"Utilisateur\", \"age\": 25, \"ville\": \"Paris\"}",
        "", "# keys(), values(), items()", "print(personne.keys())    # dict_keys(['nom', 'age', 'ville'])", "for cle in personne.keys(): print(cle)",
        "", "print(personne.values())  # dict_values(['Lucas', 14, 'Paris'])",
        "", "for cle, valeur in personne.items():", "    print(f\"{cle}: {valeur}\")",
        "", "# Supprimer", "del personne[\"ville\"]     # supprime la clé", "age = personne.pop(\"age\")  # supprime et retourne",
        "", "# Mettre à jour (fusion)", "infos = {\"taille\": 170}", "personne.update(infos)"
    ], "items() est très pratique pour parcourir clés ET valeurs en même temps")
    make_content_slide(prs, "Liste vs Dictionnaire - Quand utiliser quoi ?", [
        ("Utilise une LISTE quand :", 0), "Les éléments sont ORDONNÉS (position importante)", "Tu veux accéder par position (index)", "Ex : liste de notes, noms d'amis",
        "", ("Utilise un DICTIONNAIRE quand :", 0), "Tu veux accéder par NOM (pas par position)", "Tu as des paires clé-valeur", "Ex : profil d'un utilisateur, configuration",
        "", ("On combine souvent les deux !", 0), "        liste_de_dicts = [", "    {\"nom\": \"Alice\", \"age\": 20},", "    {\"nom\": \"Bob\", \"age\": 25}", "]"
    ])
    make_code_slide(prs, "Applications pratiques", [
        "# 1. Compteur de mots", "phrase = \"le chat mange le poisson et le chat dort\"", "mots = phrase.split()", "compteur = {}", "for mot in mots:", "    if mot in compteur: compteur[mot] += 1", "    else: compteur[mot] = 1", "print(compteur)",
        "", "# 2. Annuaire téléphonique", "annuaire = {\"Alice\": \"06.12.34.56.78\", \"Bob\": \"07.98.76.54.32\"}", "nom = input(\"Nom à chercher : \")", "tel = annuaire.get(nom)", "print(f\"{nom} : {tel}\") if tel else print(\"Inconnu\")"
    ], "Les dictionnaires sont PARFAITS pour les compteurs, annuaires, configurations...")
    make_exercise_slide(prs, "Exercice : Gestion de contacts",
        "Crée un programme de gestion de contacts :\n1. Crée un dictionnaire avec 3 contacts (nom, âge, ville)\n2. Affiche tous les contacts formatés\n3. Ajoute un nouveau contact saisi par l'utilisateur\n4. Permet de chercher un contact par son nom\n5. Affiche les noms des contacts triés alphabétiquement\n\nBONUS : Ajoute un menu interactif (1=Ajouter, 2=Chercher, 3=Afficher, 4=Quitter)",
        ["contacts = {\"Alice\": {\"age\":15, \"ville\":\"Lyon\"}, \"Bob\": {\"age\":14, \"ville\":\"Paris\"}}", "for nom, infos in contacts.items():", "    print(f\"{nom}: {infos['age']} ans, {infos['ville']}\")", "nom = input(\"Nouveau contact - Nom : \")", "age = int(input(\"Âge : \"))", "ville = input(\"Ville : \")", "contacts[nom] = {\"age\": age, \"ville\": ville}", "rech = input(\"Chercher : \")", "print(contacts.get(rech, \"Inconnu\"))"]
    )
    make_summary_slide(prs, ["Un dictionnaire stocke des paires clé-valeur", "Accès par clé : dict[\"clé\"] ou dict.get(\"clé\")", "Méthodes : keys(), values(), items(), update(), pop()", "Les clés sont UNIQUES, les valeurs peuvent être de tout type", "Les dictionnaires sont partout : configs, APIs, JSON..."])
    save_presentation(prs, "18_Dictionnaires.pptx")
    return prs

def p19_exercices_listes_boucles():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Exercices : Listes et Boucles", "Pratique avec les listes et les boucles", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Consolider les listes", "Maîtriser les boucles for et while", "Combiner listes, boucles et conditions", "Résoudre des problèmes concrets"])
    make_exercise_slide(prs, "Exercice 1 : Somme et moyenne",
        "1. Crée une liste de 10 nombres\n2. Calcule la somme avec une boucle (SANS sum())\n3. Calcule la moyenne\n4. Trouve le plus grand nombre (SANS max())\n5. Trouve le plus petit nombre (SANS min())\n6. Compte le nombre de nombres pairs\n\nBONUS : Calcule l'écart-type",
        ["nombres = [12, 7, 19, 5, 3, 15, 8, 10, 6, 1]", "somme = 0", "for n in nombres: somme += n", "moyenne = somme / len(nombres)", "print(f\"Somme: {somme}, Moyenne: {moyenne:.2f}\")", "maxi = min(nombres)  # start bas", "for n in nombres:", "    if n > maxi: maxi = n", "print(f\"Max: {maxi}\")", "pairs = sum(1 for n in nombres if n % 2 == 0)", "print(f\"Pairs: {pairs}\")"]
    )
    make_exercise_slide(prs, "Exercice 2 : Filtrer une liste",
        "Crée un programme qui :\n1. Crée une liste de mots\n2. Affiche les mots de plus de 5 lettres\n3. Affiche les mots contenant la lettre 'a'\n4. Convertit tous les mots en majuscules\n5. Trie les mots par longueur\n6. Trouve le mot le plus long\n\nBONUS : Affiche un classement par ordre alphabétique inversé",
        ["mots = [\"python\", \"java\", \"javascript\", \"c\", \"ruby\", \"rust\", \"php\"]", "print(\"+ de 5 lettres:\")", "for m in mots:", "    if len(m) > 5: print(f\"  - {m}\")", "print(\"Contiennent 'a':\")", "for m in mots:", "    if 'a' in m: print(f\"  - {m}\")", "mots_maj = [m.upper() for m in mots]", "print(f\"Maj: {mots_maj}\")", "trie_long = sorted(mots, key=len)", "print(f\"Par longueur: {trie_long}\")", "plus_long = max(mots, key=len)", "print(f\"Le + long: {plus_long}\")"]
    )
    make_exercise_slide(prs, "Exercice 3 : Simulateur de dé",
        "Crée un simulateur de lancer de dé :\n1. Lance un dé à 6 faces 100 fois (randint)\n2. Stocke chaque résultat dans une liste\n3. Calcule la fréquence de chaque face\n4. Affiche un histogramme en ASCII\n5. Calcule la moyenne (doit être ~3.5)\n6. Affiche la face la plus fréquente",
        ["import random", "resultats = [random.randint(1,6) for _ in range(100)]", "for face in range(1,7):", "    count = resultats.count(face)", "    print(f\"Face {face}: {'█'*count} ({count})\")", "print(f\"Moyenne: {sum(resultats)/len(resultats):.2f}\")", "plus_freq = max(range(1,7), key=lambda f: resultats.count(f))", "print(f\"+ fréquente: {plus_freq}\")"]
    )
    make_exercise_slide(prs, "Exercice 4 : Palindrome",
        "Un palindrome est un mot qui se lit dans les deux sens.\nExemples : radar, kayak, été, ressasser\n\nCrée un programme qui :\n1. Demande un mot à l'utilisateur\n2. Vérifie si c'est un palindrome\n3. Affiche le résultat\n\nBONUS : ignore la casse et les accents\nBONUS 2 : fonctionne pour les phrases (sans espaces)",
        ["mot = input(\"Mot à vérifier : \").lower()", "if mot == mot[::-1]:", "    print(f\"{mot} est un palindrome !\")", "else:", "    print(f\"{mot} n'est pas un palindrome\")", "# Méthode boucle :", "est_pal = True", "for i in range(len(mot)//2):", "    if mot[i] != mot[-(i+1)]: est_pal = False; break", "print(f\"Vérification: {est_pal}\")"]
    )
    make_exercise_slide(prs, "Exercice 5 : Générateur de mots de passe",
        "Crée un générateur de mot de passe sécurisé :\n1. Demande la longueur souhaitée\n2. Génère un mot de passe avec : minuscules, majuscules, chiffres, symboles\n3. Assure au moins 1 de chaque type\n4. Affiche le mot de passe et sa force (faible/moyen/fort)",
        ["import random, string", "long = int(input(\"Longueur : \"))", "chars = string.ascii_letters + string.digits + \"!@#$%^&*\"", "mdp = \"\".join(random.choice(chars) for _ in range(long))", "print(f\"Mot de passe : {mdp}\")", "force = \"Faible\"", "if any(c.islower() for c in mdp) and any(c.isupper() for c in mdp): force = \"Moyen\"", "if any(c.isdigit() for c in mdp) and any(c in \"!@#$%^&*\" for c in mdp): force = \"Fort\"", "print(f\"Force : {force}\")"]
    )
    make_summary_slide(prs, ["Les boucles et listes sont des outils complémentaires très puissants", "On peut filtrer, transformer, analyser des listes avec des boucles", "Pratique régulièrement pour être à l'aise !", "Ces exercices sont typiques de la programmation", "Bravo si tu as tout fait !"])
    save_presentation(prs, "19_Exercices_listes_boucles.pptx")
    return prs

def p20_mini_projet_calculatrice():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Mini-Projet : Calculatrice Interactive", "Mets tout en pratique !", "Partie 2 - Les Bases de Python")
    make_objectives_slide(prs, ["Appliquer variables, conditions, boucles", "Créer un programme complet et utile", "Structurer son code proprement", "Tester et améliorer son programme"])
    make_content_slide(prs, "Objectif du projet", [
        ("Créer une calculatrice en ligne de commande qui permet de :", 0), "",
        "1. Choisir une opération dans un menu", "2. Saisir les nombres", "3. Afficher le résultat", "4. Revenir au menu jusqu'à 'Quitter'",
        "", ("Fonctionnalités :", 0), "✅ Addition, Soustraction, Multiplication", "✅ Division avec gestion de /0", "✅ Puissance, Racine carrée, Modulo", "✅ Historique des calculs", "✅ Gestion des erreurs de saisie",
        "", "Tout ce qu'on a appris en action !"
    ])
    make_code_slide(prs, "Structure du programme", [
        "import math", "historique = []",
        "", "def afficher_menu():",
        "    print(\"=\"*30)", "    print(\"    CALCULATRICE\")", "    print(\"=\"*30)", "    print(\"1.Addition 2.Soustraction\")", "    print(\"3.Multiplication 4.Division\")", "    print(\"5.Puissance 6.Racine\")", "    print(\"7.Modulo 8.Historique 9.Quitter\")",
        "", "while True:", "    afficher_menu()", "    choix = input(\"Choix (1-9) : \")",
        "", "    if choix == '9': break",
        "    if choix == '8':", "        for c in historique: print(f\"  {c}\")", "        continue",
        "", "    if choix == '6':", "        a = float(input(\"Nombre : \"))", "        if a < 0: print(\"Racine négative !\"); continue", "        r = math.sqrt(a); print(f\"√{a} = {r}\")",
        "    else:", "        # operations... (voir slide suivant)"
    ], "On structure le code avec des fonctions pour le rendre clair")
    make_code_slide(prs, "Le corps principal", [
        "def addition(a,b): return a+b", "def soustraction(a,b): return a-b", "def multiplication(a,b): return a*b", "def division(a,b):", "    return \"Erreur: /0!\" if b==0 else a/b",
        "", "operations = {", "    '1': ('+', addition),", "    '2': ('-', soustraction),", "    '3': ('×', multiplication),", "    '4': ('÷', division),", "    '5': ('^', lambda x,y: x**y),", "    '7': ('%', lambda x,y: x%y)", "}",
        "", "a = float(input(\"a : \"))", "b = float(input(\"b : \"))", "symbole, fonction = operations[choix]", "resultat = fonction(a, b)", "print(f\"{a} {symbole} {b} = {resultat}\")",
        "", "if not isinstance(resultat, str):", "    historique.append(f\"{a} {symbole} {b} = {resultat}\")", "    if len(historique) > 10: historique.pop(0)"
    ], "Le dictionnaire operations permet de choisir la fonction selon le choix")
    make_exercise_slide(prs, "Améliorations à ajouter",
        "Ajoute ces fonctionnalités :\n\n1. 🔄 Logarithme, sinus, cosinus\n2. 📊 Moyenne des résultats de l'historique\n3. 💾 Sauvegarde de l'historique dans un fichier\n4. 🎨 Couleurs dans le terminal\n5. 🔣 Calcul en une ligne : \"2 + 3 * 4\"\n6. 📋 Permet de copier le résultat\n\nTeste et vérifie :\n- Division par zéro\n- Racine d'un nombre négatif\n- Saisie texte au lieu de nombre\n- Choix invalide dans le menu"
    )
    make_summary_slide(prs, ["Ce projet combine : variables, conditions, boucles, listes, fonctions", "On a appris à structurer un programme complet", "L'historique montre comment stocker et afficher des données", "La gestion d'erreurs rend le programme robuste", "Continue à ajouter des fonctionnalités pour t'améliorer !"])
    save_presentation(prs, "20_Mini_Projet_Calculatrice.pptx")
    return prs

# ═══════════════════════════════════════════════════════════════
# PARTIE 3 : FONCTIONS, FICHIERS & MODULES
# ═══════════════════════════════════════════════════════════════

def p21_fonctions_definition():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Les Fonctions - Définition", "Écrire du code réutilisable", "Partie 3 - Fonctions & Modules")
    make_objectives_slide(prs, ["Comprendre l'utilité des fonctions", "Savoir définir une fonction avec def", "Comprendre le return", "Appeler une fonction"])
    make_content_slide(prs, "Pourquoi des fonctions ?", [
        ("Problème SANS fonction :", 0), "print(\"Bonjour Lucas\")", "print(\"Bonjour Marie\")", "# Si on change le message, tout modifier !",
        "", ("Avec UNE fonction :", 0), "def saluer(prenom):", "    print(f\"Bonjour {prenom}\")", "saluer(\"Lucas\"); saluer(\"Marie\")",
        "", ("Avantages :", 0), "✅ RÉUTILISATION : appelable autant de fois qu'on veut", "✅ MAINTENANCE : changer à UN SEUL endroit", "✅ LISIBILITÉ : code plus clair et organisé", "✅ TEST : chaque fonction indépendante"
    ])
    make_code_slide(prs, "Définir et appeler une fonction", [
        "# Définition", "def ma_fonction():", "    print(\"Dans la fonction !\")", "", "# Appel (parenthèses OBLIGATOIRES)", "ma_fonction()",
        "", "# Fonction avec paramètre", "def dire_bonjour(prenom):", "    print(f\"Bonjour {prenom} !\")", "dire_bonjour(\"Alice\")",
        "", "# Fonction avec plusieurs paramètres", "def presentation(nom, age):", "    print(f\"Je m'appelle {nom}, j'ai {age} ans\")", "presentation(\"Utilisateur\", 25)",
        "", "# Fonction avec RETURN", "def carre(x):", "    return x * x", "resultat = carre(5); print(resultat)  # 25"
    ], "def = définition, () = paramètres, : = bloc, return = valeur de retour")
    make_content_slide(prs, "return vs print", [
        ("DIFFÉRENCE FONDAMENTALE :", 0), "",
        ("print() = AFFICHE dans la console (visible à l'écran)", 0), ("return = RENVOIE une valeur au programme (utilisable dans le code)", 0), "",
        ("Exemple :", 0), "def addition_print(a, b):", "    print(a + b)   # affiche le résultat",
        "def addition_return(a, b):", "    return a + b   # renvoie le résultat", "",
        "addition_print(3, 4)  # → affiche 7", "x = addition_return(3, 4)  # x = 7, rien n'est affiché", "print(x * 2)  # → 14 (on peut UTILISER le résultat !)",
        "", "En général : les fonctions calculent (return) et le programme principal affiche (print)"
    ])
    make_code_slide(prs, "return : points importants", [
        "        # return peut renvoyer plusieurs valeurs", "def get_infos():", "    return \"Utilisateur\", 25  # retourne un tuple", "nom, age = get_infos()",
        "", "# return arrête IMMÉDIATEMENT la fonction", "def verifier_age(age):", "    if age < 18:", "        return \"Mineur\"  # on sort ici", "    return \"Majeur\"",
        "", "# Sans return, la fonction retourne None", "def sans_return():", "    x = 5", "resultat = sans_return()", "print(resultat)  # None", "print(type(resultat))  # <class 'NoneType'>"
    ], "None = le type spécial qui signifie \"rien\"")
    make_content_slide(prs, "Documenter une fonction (docstring)", [
        ("La docstring explique ce que fait la fonction :", 0), "",
        "def calculer_imc(poids, taille):", '    """Calcule l\'Indice de Masse Corporelle.', "", '    Args:', '        poids (float): Poids en kg', '        taille (float): Taille en mètres', "", '    Returns:', '        float: L\'IMC calculé', '    """', "    return poids / (taille ** 2)",
        "", ("On consulte avec :", 0), "help(calculer_imc)", "print(calculer_imc.__doc__)",
        "", "Les docstrings (triples guillemets) sont la convention officielle (PEP 257)"
    ])
    make_exercise_slide(prs, "Exercice : Bibliothèque de fonctions",
        "Crée ces fonctions et teste-les :\n1. saluer(nom) → affiche \"Bonjour [nom] !\"\n2. carre(x) → retourne x²\n3. est_pair(n) → retourne True si n est pair\n4. max3(a, b, c) → retourne le plus grand des 3\n5. factorielle(n) → retourne n! (produit de 1 à n)\n\nChaque fonction doit avoir une docstring !",
        ["def saluer(nom):", '    """Dit bonjour."""', "    print(f\"Bonjour {nom} !\")", "def carre(x): return x * x", "def est_pair(n): return n % 2 == 0", "def max3(a, b, c): return max(a,b,c)", "def factorielle(n):", "    r = 1", "    for i in range(2, n+1): r *= i", "    return r"]
    )
    make_summary_slide(prs, ["Une fonction est un bloc de code réutilisable avec def", "Les paramètres sont les entrées, return est la sortie", "return ≠ print : return donne une valeur, print affiche", "Sans return, la fonction retourne None", "Les fonctions rendent le code : modulaire, réutilisable, lisible"])
    save_presentation(prs, "21_Fonctions_definition.pptx")
    return prs

def p22_fonctions_parametres():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Fonctions - Paramètres Avancés", "Paramètres par défaut, *args, **kwargs", "Partie 3 - Fonctions & Modules")
    make_objectives_slide(prs, ["Maîtriser les paramètres par défaut", "Arguments positionnels et nommés", "Utiliser *args pour nombre variable d'arguments", "Utiliser **kwargs pour arguments nommés variables"])
    make_content_slide(prs, "Paramètres par défaut", [
        ("On peut donner des valeurs par défaut :", 0), "",
        "def saluer(prenom, message=\"Bonjour\"):", "    print(f\"{message}, {prenom} !\")",
        "saluer(\"Lucas\")           # Bonjour, Lucas !", "saluer(\"Lucas\", \"Salut\")  # Salut, Lucas !",
        "", ("Règle : paramètres avec défaut APRÈS ceux sans défaut", 0), "❌ def exemple(a=1, b):  # ERREUR !", "✅ def exemple(b, a=1):  # OK",
        "", ("Attention aux MUTABLES par défaut !", 0), "def ajouter(element, liste=[]):  # PIÈGE !", "    liste.append(element); return liste", "print(ajouter(1))  # [1]", "print(ajouter(2))  # [1, 2] → PARTAGÉ !",
        "", "Solution : def ajouter(element, liste=None):", "    if liste is None: liste = []"
    ])
    make_code_slide(prs, "Arguments positionnels vs nommés", [
        "# Arguments POSITIONNELS (dans l'ordre)", "def profil(nom, age, ville):", "    print(f\"{nom}, {age} ans, {ville}\")", "profil(\"Utilisateur\", 25, \"Paris\")",
        "", "# Arguments NOMMÉS (keyword arguments)", "profil(age=25, nom=\"Utilisateur\", ville=\"Paris\")  # ordre libre",
        "", "# Mélange : positionnels D'ABORD", "profil(\"Utilisateur\", ville=\"Paris\", age=25)  # OK", "profil(nom=\"Utilisateur\", 25, \"Paris\")  # ERREUR !",
        "", "# Utilité : rend le code plus lisible", "def configurer(hote, port=80, ssl=True, timeout=30):", "    ...", "configurer(\"192.168.1.1\", port=8080, ssl=True)"
    ], "Les arguments nommés rendent le code PLUS LISIBLE")
    make_code_slide(prs, "*args - Nombre variable d'arguments", [
        "# *args = nombre VARIABLE d'arguments positionnels", "def somme(*args):", "    print(f\"Reçus : {args}\")  # args est un tuple", "    total = 0", "    for n in args: total += n", "    return total",
        "print(somme(1, 2))           # 3", "print(somme(1, 2, 3, 4, 5))   # 15", "print(somme())                  # 0",
        "", "# *args avec d'autres paramètres", "def afficher(titre, *elements):", "    print(f\"=== {titre} ===\")", "    for e in elements: print(f\"  - {e}\")", "afficher(\"Courses\", \"Pain\", \"Lait\", \"Œufs\")",
        "", "# Déballer une liste en arguments", "nombres = [10, 20, 30]", "print(somme(*nombres))  # 60"
    ], "*args récupère les arguments supplémentaires dans un TUPLE")
    make_code_slide(prs, "**kwargs - Arguments nommés variables", [
        "# **kwargs = nombre VARIABLE d'arguments NOMMÉS", "def afficher_profil(**kwargs):", "    print(f\"Reçus : {kwargs}\")  # kwargs est un dict", "    for k, v in kwargs.items(): print(f\"  {k}: {v}\")",
        "afficher_profil(nom=\"Utilisateur\", age=25, ville=\"Paris\")",
        "", "# *args ET **kwargs combinés", "def super_fonction(a, b, *args, **kwargs):", "    print(f\"a={a}, b={b}, args={args}, kwargs={kwargs}\")",
        "super_fonction(1, 2, 3, 4, x=10, y=20)", "# a=1, b=2, args=(3, 4), kwargs={'x':10, 'y':20}",
        "", "# Ordre correct : positionnels → *args → défauts → **kwargs"
    ], "**kwargs est utilisé dans Django, Flask, décorateurs...")
    make_exercise_slide(prs, "Exercice : Fonctions avancées",
        "1. Crée multiplier(*nombres) qui multiplie tous les nombres\n2. Crée creer_utilisateur(nom, **infos) qui affiche le profil\n3. Crée formater_texte(texte, couleur=\"blanc\", taille=12, gras=False)\n4. Crée une fonction qui calcule la moyenne d'une liste de notes\n5. Crée wrapper_temps(fonction) qui mesure le temps d'exécution\n\nBONUS : Démontre l'ordre correct des paramètres",
        ["def multiplier(*nombres):", "    r = 1", "    for n in nombres: r *= n", "    return r",
         "def creer_utilisateur(nom, **infos):", "    print(f\"Profil: {nom}\")", "    for k,v in infos.items(): print(f\"  {k}: {v}\")",
         "creer_utilisateur(\"Lucas\", age=14, ville=\"Paris\")"]
    )
    make_summary_slide(prs, ["Paramètres par défaut : def f(x, y=10)", "Arguments positionnels (ordre) vs nommés (nom=)", "*args → tuple d'arguments positionnels variables", "**kwargs → dict d'arguments nommés variables", "Ordre : positionnels → *args → défauts → **kwargs"])
    save_presentation(prs, "22_Fonctions_parametres.pptx")
    return prs

def p23_portee_variables():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Portée des Variables (Scope)", "Où les variables sont-elles accessibles ?", "Partie 3 - Fonctions & Modules")
    make_objectives_slide(prs, ["Comprendre la notion de portée", "Distinguer variable locale et globale", "Connaître l'ordre LEGB", "Utiliser global et nonlocal"])
    make_content_slide(prs, "C'est quoi la portée (scope) ?", [
        ("La portée d'une variable = où elle est accessible dans le code.", 0), "",
        ("Règle d'or :", 0), "Une variable définie DANS une fonction n'existe que DANS cette fonction.", "Une variable définie en DEHORS de toute fonction est globale.",
        "", ("Exemple :", 0), "x = 10  # variable GLOBALE", "def ma_fonction():", "    y = 5  # variable LOCALE", "    print(x)  # OK (lecture globale)", "    print(y)  # OK",
        "", "ma_fonction()", "print(x)  # OK (x est globale)", "print(y)  # ERREUR ! y n'existe pas ici",
        "", "Chaque fonction a son propre ESPACE (environnement) pour les variables."
    ])
    make_code_slide(prs, "Variable locale vs globale", [
        "x = 100  # Variable GLOBALE", "", "def modifier_x():", "    x = 50  # Variable LOCALE (masque la globale)", "    print(f\"Dans fonction: x = {x}\")",
        "modifier_x()           # Dans fonction: x = 50", "print(f\"En dehors: x = {x}\")  # En dehors: x = 100  (inchangée !)",
        "", "# Pour modifier une variable globale :", "y = 10", "def modifier_globale():", "    global y  # déclaration", "    y = 20",
        "modifier_globale()", "print(y)  # 20 (la globale a changé)"
    ], "global à utiliser avec PRUDENCE. Préfère return.")
    make_content_slide(prs, "L'ordre LEGB (Python Scope)", [
        ("Quand Python cherche une variable, il suit cet ordre :", 0), "",
        ( "L → Local (dans la fonction actuelle)", 1), ( "E → Enclosing (dans les fonctions englobantes)", 1), ( "G → Global (au niveau du module/fichier)", 1), ( "B → Built-in (fonctions intégrées de Python)", 1), "",
        ("Exemple :", 0), "x = \"global\"", "def externe():", "    x = \"enclosing\"", "    def interne():", "        x = \"local\"", "        print(x)  # → \"local\" (LEGB: trouve LOCAL d'abord)", "    interne()", "externe()",
        "", "Python remonte du plus spécifique (local) au plus général (built-in)"
    ])
    make_code_slide(prs, "nonlocal - modifier une variable englobante", [
        "def externe():", "    x = 10  # enclosing", "    def interne():", "        nonlocal x  # on veut modifier x de l'englobante", "        x = 20", "    interne()", "    print(f\"Externe: {x}\")  # Externe: 20",
        "externe()",
        "", "# Sans nonlocal, ça crée une variable locale", "def externe2():", "    x = 10", "    def interne2():", "        x = 99  # nouvelle variable locale", "    interne2()", "    print(x)  # 10 (inchangé)"
    ], "nonlocal = \"modifier la variable de la fonction parente\"")
    make_content_slide(prs, "Bonnes pratiques sur le scope", [
        ("À FAIRE :", 0),
        "✅ Utiliser des paramètres pour passer des données aux fonctions",
        "✅ Retourner les valeurs modifiées avec return",
        "✅ Préférer des fonctions pures (qui dépendent seulement de leurs entrées)",
        "", ("À ÉVITER :", 0),
        "❌ Trop de variables globales → code difficile à suivre",
        "❌ Modifier des globales dans des fonctions (sauf si nécessaire)",
        "", ("Exemple :", 0), "# ❌ Pas bien", "total = 0", "def ajouter(x):", "    global total; total += x",
        "", "# ✅ Bien", "def ajouter(total, x):", "    return total + x"
    ])
    make_exercise_slide(prs, "Exercice : Comprendre le scope",
        "Sans exécuter le code, trouve ce qui va s'afficher :\n\na = 5\ndef f1():\n    print(a)\nf1()  # ?\n\nb = 10\ndef f2():\n    b = 20\n    print(b)\nf2()\nprint(b)  # ?\n\nc = 1\ndef f3():\n    global c\n    c += 1\n    print(c)\nf3()\nprint(c)  # ?\n\nVérifie en exécutant !",
        ["# Réponses:", "# f1() → 5 (lit la globale a)", "# f2() → 20 (locale), print(b) → 10 (globale inchangée)", "# f3() → 2 (globale modifiée), print(c) → 2"]
    )
    make_summary_slide(prs, ["Portée = où une variable est accessible", "Variable locale = définie dans une fonction", "Variable globale = en dehors de toute fonction", "LEGB : Local, Enclosing, Global, Built-in", "global = modifier une globale, nonlocal = modifier une englobante", "Préfère les fonctions pures (paramètres + return)"])
    save_presentation(prs, "23_Portee_variables.pptx")
    return prs

def p24_tuples_sets():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Tuples et Sets", "Collections immutables et ensembles", "Partie 3 - Fonctions & Modules")
    make_objectives_slide(prs, ["Comprendre les tuples (immuables)", "Comprendre les sets (ensembles sans doublons)", "Savoir quand utiliser chaque type", "Connaître les opérations sur les sets"])
    make_two_col_slide(prs, "Tuple vs List vs Set",
        "List", ["✅ Modifiable", "✅ Ordonné", "✅ Doublons", "✅ Indexé", "Syntaxe: [...]"],
        "Tuple vs Set",
        ["❌/✅/❌ Modifiable", "✅/✅/❌ Ordonné", "✅/✅/❌ Doublons", "✅/✅/❌ Indexé", "(...) / {...}"]
    )
    make_code_slide(prs, "Les Tuples - Collections immutables", [
        "# Un tuple est une liste... qui ne peut pas être modifiée !",
        "", "# Création", "vide = ()", "coords = (10, 20)", "trois = 1, 2, 3  # parenthèses optionnelles",
        "", "# Tuple avec 1 élément (attention à la virgule !)", "simple = (5,)   # tuple avec 1 élément", "pas_tuple = (5)  # simple entier !",
        "", "# Accès (comme les listes)", "print(coords[0])   # 10", "print(coords[-1])  # 20",
        "", "# Tentative de modification → ERREUR", "coords[0] = 5  # TypeError !"
    ], "Les tuples sont utilisés pour les données qui ne doivent PAS changer")
    make_code_slide(prs, "Pourquoi utiliser des tuples ?", [
        "# 1. RETOUR MULTIPLE de fonctions", "def min_max(liste):", "    return min(liste), max(liste)  # retourne un tuple",
        "mini, maxi = min_max([3, 1, 7, 2, 9])", "print(mini, maxi)  # 1 9",
        "", "# 2. CLÉS DE DICTIONNAIRE (tuples oui, listes non)", "positions = {(0,0): \"Origine\", (1,0): \"Droite\"}", "print(positions[(0,0)])  # Origine",
        "", "# 3. DONNÉES IMMUTABLES", "TAILLE = (1920, 1080)", "COULEUR = (255, 200, 100)"
    ], "Les tuples sont plus LÉGERS et plus RAPIDES que les listes")
    make_code_slide(prs, "Les Sets - Ensembles sans doublons", [
        "# Un set = collection NON ordonnée, SANS doublons",
        "", "# Création", "set_vide = set()  # {} créerait un dict vide !", "fruits = {\"pomme\", \"banane\", \"cerise\", \"pomme\"}", "print(fruits)  # {'cerise', 'pomme', 'banane'} - pomme qu'une fois !",
        "", "# Créer un set à partir d'une liste", "nombres = [1, 2, 2, 3, 3, 3, 4, 5, 5]", "set_nombres = set(nombres)", "print(set_nombres)  # {1, 2, 3, 4, 5}"
    ], "Les sets sont parfaits pour SUPPRIMER LES DOUBLONS")
    make_code_slide(prs, "Opérations sur les sets", [
        "a = {1, 2, 3, 4, 5}; b = {4, 5, 6, 7, 8}",
        "", "# UNION : tous les éléments", "print(a | b)   # {1, 2, 3, 4, 5, 6, 7, 8}", "print(a.union(b))",
        "", "# INTERSECTION : éléments COMMUNS", "print(a & b)   # {4, 5}", "print(a.intersection(b))",
        "", "# DIFFÉRENCE : dans a mais PAS dans b", "print(a - b)   # {1, 2, 3}",
        "", "# DIFFÉRENCE SYMÉTRIQUE : dans a ou b, pas les deux", "print(a ^ b)   # {1, 2, 3, 6, 7, 8}",
        "", "# Vérifications", "print(3 in a)  # True", "print({1,2}.issubset(a))  # True"
    ], "Les opérations sur les sets sont TRÈS RAPIDES (O(1) pour la recherche)")
    make_exercise_slide(prs, "Exercice : Tuples et Sets",
        "1. Crée un tuple avec 5 nombres, essaie de le modifier → constate l'erreur\n2. Crée une fonction qui retourne (somme, moyenne, min, max) d'une liste\n3. Crée deux listes avec doublons, trouve : a) communs b) uniques c) éléments distincts\n4. Vérifie si \"python\" est dans un set de langages\n5. (BONUS) Trouve les caractères uniques dans une phrase",
        ["# Caractères uniques", "phrase = \"le chat et le chien jouent\"", "caracteres = set(phrase.replace(\" \", \"\"))", "print(f\"Uniques: {sorted(caracteres)}\")", "print(f\"Nombre: {len(caracteres)}\")"]
    )
    make_summary_slide(prs, ["Tuple = liste immuable (ne peut pas être modifiée)", "Tuple : (), return multiple, clés de dict, données fixes", "Set = ensemble NON ordonné SANS doublons", "Opérations : | & - ^", "Choisis le bon type selon ton besoin !"])
    save_presentation(prs, "24_Tuples_Sets.pptx")
    return prs

def p25_gestion_erreurs():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Gestion des Erreurs", "try/except pour un code robuste", "Partie 3 - Fonctions & Modules")
    make_objectives_slide(prs, ["Comprendre les exceptions", "Utiliser try/except", "Connaître les types d'exceptions courants", "Utiliser finally et else"])
    make_content_slide(prs, "Pourquoi gérer les erreurs ?", [
        ("Un programme peut planter pour plusieurs raisons :", 0),
        "- L'utilisateur entre \"abc\" quand on attend un nombre",
        "- Un fichier n'existe pas", "- Division par zéro", "- Connexion Internet perdue",
        "", ("SANS gestion d'erreur → le programme CRASH (s'arrête brutalement)", 0),
        ("AVEC gestion d'erreur → on peut continuer ou montrer un message utile", 0), "",
        "try:", "    age = int(input(\"Âge : \"))", "except ValueError:", "    print(\"Ce n'est pas un nombre valide !\")", "    age = 0"
    ])
    make_code_slide(prs, "Syntaxe complète de try/except", [
        "try:", "    # Code qui pourrait planter", "    fichier = open(\"inexistant.txt\", \"r\")", "    nombre = int(input(\"Nombre : \"))", "    resultat = 10 / nombre",
        "", "except FileNotFoundError:", "    print(\"Le fichier n'existe pas\")",
        "except ValueError:", "    print(\"Entrez un nombre valide\")",
        "except ZeroDivisionError:", "    print(\"Division par zéro !\")",
        "except Exception as e:", "    print(f\"Erreur inattendue: {e}\")",
        "", "else:", "    print(f\"Résultat: {resultat}\")  # si PAS d'erreur",
        "", "finally:", "    print(\"Toujours exécuté\")  # nettoyage"
    ], "finally = fermer fichiers, libérer ressources...")
    make_content_slide(prs, "Exceptions courantes", [
        ("ValueError : valeur incorrecte", 0), "int(\"abc\")",
        ("TypeError : opération sur un type inapproprié", 0), "\"texte\" + 5",
        ("ZeroDivisionError : division par zéro", 0), "10 / 0",
        ("FileNotFoundError : fichier introuvable", 0), "open(\"inconnu.txt\")",
        ("IndexError : index hors limites", 0), "liste = [1, 2, 3]; liste[10]",
        ("KeyError : clé inexistante", 0), "d = {\"a\": 1}; d[\"b\"]",
        ("ImportError : module introuvable", 0), "import module_inexistant",
    ])
    make_code_slide(prs, "Astuces et bonnes pratiques", [
        "# 1. Attraper avec message", "try:", "    resultat = 10 / 0", "except ZeroDivisionError as e:", "    print(f\"Erreur: {e}\")  # division by zero",
        "", "# 2. Grouper plusieurs exceptions", "try:", "    valeur = int(input(\"Nombre: \"))", "except (ValueError, TypeError):", "    print(\"Entrée invalide\")",
        "", "# 3. Ne JAMAIS faire : except nu", "try:", "    ...", "except:  # masque tout !", "    pass",
        "", "# 4. raise : déclencher une exception", "def diviser(a, b):", "    if b == 0:", "        raise ValueError(\"Division par zéro interdite\")", "    return a / b"
    ], "raise permet de créer TES PROPRES erreurs pour valider des entrées")
    make_code_slide(prs, "Exemple complet : saisie robuste", [
        "def demander_entier(message):", '    """Demander un entier de manière robuste."""', "    while True:", "        try:", "            return int(input(message))", "        except ValueError:", "            print(\"Nombre invalide !\")",
        "", "def demander_entre_bornes(mini, maxi):", "    while True:", "        try:", "            v = int(input(f\"Entre {mini}-{maxi}: \"))", "            if mini <= v <= maxi: return v", "        except ValueError:", "            print(\"Invalide\")",
        "", "# Utilisation", "age = demander_entre_bornes(0, 120)", "note = demander_entre_bornes(0, 20)"
    ], "Cette technique est très utilisée dans les vrais programmes")
    make_exercise_slide(prs, "Exercice : Saisie sécurisée",
        "1. Crée une fonction lire_entier(message) qui :\n   - Affiche le message\n   - Essaie de convertir en int\n   - Si erreur → \"Erreur !\" et redemande\n   - Continue jusqu'à obtenir un entier valide\n\n2. Crée un mini programme de calcul :\n   - Demande 2 nombres avec lire_entier()\n   - Demande l'opération (+, -, *, /)\n   - Gère la division par zéro\n   - Affiche le résultat",
        ["def lire_entier(message):", "    while True:", "        try:", "            return int(input(message))", "        except ValueError:", "            print(\"Entrez un nombre valide\")",
         "a = lire_entier(\"a : \"); b = lire_entier(\"b : \")", "op = input(\"Opération (+, -, *, /) : \")", "if op == '+': print(f\"{a}+{b}={a+b}\")", "elif op == '-': print(f\"{a}-{b}={a-b}\")", "elif op == '*': print(f\"{a}*{b}={a*b}\")", "elif op == '/':", "    try: print(f\"{a}/{b}={a/b}\")", "    except ZeroDivisionError: print(\"/0!\")"]
    )
    make_summary_slide(prs, ["Les exceptions sont des erreurs qui peuvent arriver pendant l'exécution", "try = bloc à risque, except = gestion, finally = toujours", "else s'exécute si PAS d'erreur dans try", "Attrape des exceptions SPÉCIFIQUES", "raise pour déclencher volontairement", "La gestion d'erreur rend les programmes ROBUSTES"])
    save_presentation(prs, "25_Gestion_erreurs.pptx")
    return prs

def p26_fichiers():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Lecture et Écriture de Fichiers", "Sauvegarder et lire des données", "Partie 3 - Fonctions & Modules")
    make_objectives_slide(prs, ["Comprendre comment ouvrir un fichier", "Lire un fichier", "Écrire dans un fichier", "Utiliser le context manager (with)"])
    make_content_slide(prs, "Pourquoi des fichiers ?", [
        ("Jusqu'à présent, les données sont PERDUES quand le programme s'arrête.", 0), "",
        ("Les fichiers permettent de :", 0), "✅ SAUVEGARDER des données entre les exécutions", "✅ LIRE des données existantes (configs)", "✅ EXPORTER des résultats", "✅ CRÉER des logs",
        "", ("Modes d'ouverture :", 0), ("\"r\" : read (lecture) - mode par défaut", 1), ("\"w\" : write (écriture) - ÉCRASE le fichier s'il existe", 1), ("\"a\" : append (ajout) - AJOUTE à la fin", 1), ("\"x\" : exclusive - crée nouveau fichier (erreur si existe)", 1),
    ])
    make_code_slide(prs, "Ouvrir et lire un fichier", [
        "# Méthode 1 : open() / close() (risque d'oubli)", "fichier = open(\"mon_fichier.txt\", \"r\", encoding=\"utf-8\")", "contenu = fichier.read()", "fichier.close()  # IMPORTANT !",
        "", "# Méthode 2 : with (RECOMMANDÉE) - fermeture auto", "with open(\"mon_fichier.txt\", \"r\", encoding=\"utf-8\") as f:", "    contenu = f.read()", "    print(contenu)", "# Fermé automatiquement après le bloc",
        "", "# Différentes façons de lire", "with open(\"fichier.txt\", \"r\") as f:", "    contenu = f.read()          # tout en une string", "    ligne = f.readline()        # une seule ligne", "    lignes = f.readlines()      # toutes les lignes dans une liste",
        "", "    # Parcourir ligne par ligne (recommandé)", "    for ligne in f:", "        print(ligne.strip())"
    ], "Toujours utiliser with open(...) as f:")
    make_code_slide(prs, "Écrire dans un fichier", [
        "# Écrire (mode 'w' = write, écrase le fichier)", "with open(\"sortie.txt\", \"w\", encoding=\"utf-8\") as f:", "    f.write(\"Première ligne\\n\")", "    f.write(\"Deuxième ligne\\n\")",
        "", "# Ajouter (mode 'a' = append)", "with open(\"log.txt\", \"a\", encoding=\"utf-8\") as f:", "    f.write(\"Nouvelle entrée\\n\")",
        "", "# Écrire une liste de lignes", "lignes = [\"Ligne 1\\n\", \"Ligne 2\\n\", \"Ligne 3\\n\"]", "with open(\"multi.txt\", \"w\") as f:", "    f.writelines(lignes)",
        "", "# Vérifier avant d'écraser", "import os", "if not os.path.exists(\"important.txt\"):", "    with open(\"important.txt\", \"w\") as f:", "        f.write(\"Données importantes\")"
    ], "encoding=\"utf-8\" pour gérer les accents")
    make_code_slide(prs, "Travailler avec des chemins", [
        "from pathlib import Path  # Solution moderne",
        "", "# pathlib (RECOMMANDÉ)", "dossier = Path(\"mes_donnees\")", "dossier.mkdir(exist_ok=True)  # crée le dossier",
        "fichier = dossier / \"data.txt\"  # utilise / pour les chemins !", "fichier.write_text(\"Hello depuis pathlib !\")", "print(fichier.read_text())",
        "", "# Lister les fichiers .txt", "for p in Path(\".\").glob(\"*.txt\"):", "    print(p.name)",
        "", "# os.path (ancienne méthode)", "import os", "print(os.getcwd())", "print(os.path.exists(\"test.txt\"))", "print(os.path.isfile(\"test.txt\"))"
    ], "pathlib est plus moderne et plus facile que os.path")
    make_exercise_slide(prs, "Exercice : Journal personnel",
        "Crée un journal personnel :\n1. Demande une entrée à l'utilisateur\n2. Sauvegarde chaque entrée avec date et heure\n3. Propose d'afficher toutes les entrées passées\n4. Propose de chercher un mot dans le journal\n5. Compte le nombre d'entrées total\n\nBONUS : Crée un fichier différent par mois (janvier.txt, février.txt...)",
        ["from datetime import datetime", "fichier = \"journal.txt\"", "action = input(\"1=Écrire 2=Lire 3=Chercher : \")", "if action == \"1\":", "    entree = input(\"Aujourd'hui... : \")", "    with open(fichier, \"a\", encoding=\"utf-8\") as f:", "        f.write(f\"[{datetime.now():%d/%m/%Y %H:%M}] {entree}\\n\")", "    print(\"Sauvegardé !\")", "elif action == \"2\":", "    with open(fichier, \"r\", encoding=\"utf-8\") as f:", "        print(f.read())"]
    )
    make_summary_slide(prs, ["Ouvrir un fichier : open() avec 'r' (lecture), 'w' (écriture), 'a' (ajout)", "Toujours utiliser with open(...) as f:", "Lire : read(), readline(), readlines(), ou for ligne in f:", "Écrire : write() et writelines()", "pathlib est la méthode moderne pour gérer les chemins", "Les fichiers RENDENT LES DONNÉES PERMANENTES"])
    save_presentation(prs, "26_Fichiers.pptx")
    return prs

def p27_modules_import():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    make_title_slide(prs, "Modules et Import", "Utiliser et créer des bibliothèques de code", "Partie 3 - Fonctions & Modules")

if __name__ == "__main__":
    funcs = [
        p11_nombres_operations, p12_input_output, p13_conditions, p14_exercices_conditions,
        p15_boucle_while, p16_boucle_for, p17_listes, p18_dictionnaires,
        p19_exercices_listes_boucles, p20_mini_projet_calculatrice,
        p21_fonctions_definition, p22_fonctions_parametres, p23_portee_variables,
        p24_tuples_sets, p25_gestion_erreurs, p26_fichiers, p27_modules_import,
    ]
    for f in funcs:
        print(f"\n Génération : {f.__name__}")
        try:
            f()
        except Exception as e:
            print(f"  ❌ Erreur dans {f.__name__}: {e}")
    print(f"\n{'='*50}\n✅ {len(funcs)} présentations générées !\n📁 Dossier : {OUTPUT_DIR}\n{'='*50}")
