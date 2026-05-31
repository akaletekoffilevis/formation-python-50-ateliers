#!/usr/bin/env python3
"""Generate presentations 27 to 34 - Part 3 continued."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

class Th:
    D=RGBColor(0x0A,0x0A,0x2E); P=RGBColor(0x16,0x21,0x3E)
    S=RGBColor(0x1A,0x1A,0x4E); O=RGBColor(0xFF,0x6B,0x35)
    C=RGBColor(0x00,0xB4,0xD8); PU=RGBColor(0x7B,0x2F,0xF7)
    G=RGBColor(0x06,0xD6,0xA0); W=RGBColor(0xFF,0xFF,0xFF)
    L=RGBColor(0xF8,0xF9,0xFA); GR=RGBColor(0x6C,0x75,0x7D)
    LG=RGBColor(0xDE,0xE2,0xE6); CB=RGBColor(0x1E,0x1E,0x2E)
    CT=RGBColor(0xCD,0xD6,0xF4); R=RGBColor(0xEF,0x44,0x44)

OUT = "/home/akaletekoffilevis/Bureau/Coach/presentations"
os.makedirs(OUT, exist_ok=True)

def bg(sl,c):
    f=sl.background.fill; f.solid(); f.fore_color.rgb=c

def sh(sl,l,t,w,h,fc,bc=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=fc
    s.line.fill.background()
    if bc: s.line.color.rgb=bc; s.line.width=Pt(1)
    return s

def rr(sl,l,t,w,h,fc):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=fc; s.line.fill.background()
    return s

def tb(sl,l,t,w,h,text,fs=14,c=None,b=False,a=PP_ALIGN.LEFT):
    tx=sl.shapes.add_textbox(l,t,w,h); tf=tx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(fs); p.font.color.rgb=c or Th.W
    p.font.bold=b; p.font.name='Calibri'; p.alignment=a
    return tx

def tb_rich(sl,l,t,w,h):
    tx=sl.shapes.add_textbox(l,t,w,h); tf=tx.text_frame; tf.word_wrap=True; return tx

def ap(tf,text,sz=14,c=None,b=False,al=PP_ALIGN.LEFT,sb=0,sa=0,name='Calibri'):
    if len(tf.paragraphs)==1 and tf.paragraphs[0].text=='': p=tf.paragraphs[0]
    else: p=tf.add_paragraph()
    p.text=text; p.font.size=Pt(max(sz,8)); p.font.color.rgb=c or Th.W; p.font.bold=b
    p.font.name=name; p.alignment=al; p.space_before=Pt(sb); p.space_after=Pt(sa)
    return p

def code(sl,l,t,w,h,lines,fs=11):
    rr(sl,l,t,w,h,Th.CB)
    tx=tb_rich(sl,l+Inches(0.3),t+Inches(0.15),w-Inches(0.6),h-Inches(0.3))
    for line in lines:
        ap(tx.text_frame,line,sz=fs,c=Th.CT,name='Consolas',al=PP_ALIGN.LEFT)

def make_title(prs,title,sub,part):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    rr(sl,Inches(1),Inches(1.5),Inches(11.333),Inches(4.5),Th.P)
    tb(sl,Inches(1.5),Inches(1.8),Inches(10.333),Inches(1),part,fs=12,c=Th.O,b=True)
    tb(sl,Inches(1.5),Inches(2.5),Inches(10.333),Inches(1.8),title,fs=36,b=True)
    tb(sl,Inches(1.5),Inches(4.2),Inches(10.333),Inches(1),sub,fs=18,c=Th.GR)

def make_obj(prs,items):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(0.08),Inches(7.5),Th.C)
    tb(sl,Inches(0.8),Inches(0.5),Inches(11),Inches(0.8),"Objectifs",fs=28,b=True,c=Th.C)
    for i,item in enumerate(items):
        rr(sl,Inches(1),Inches(1.6+i*1.1),Inches(11),Inches(0.8),Th.S)
        tb(sl,Inches(1.3),Inches(1.7+i*1.1),Inches(10.5),Inches(0.7),f"✦ {item}",fs=16)

def make_content(prs,title,lines):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.7),title,fs=26,b=True,c=Th.O)
    tx=tb_rich(sl,Inches(0.8),Inches(1.4),Inches(11.5),Inches(5.5))
    for line in lines:
        t,sz=line if isinstance(line,tuple) else (line,15)
        ap(tx.text_frame,t,sz=sz)

def make_code_content(prs,title,lines):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    tb(sl,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),title,fs=24,b=True,c=Th.O)
    code(sl,Inches(0.8),Inches(1.2),Inches(11.5),Inches(5.5),lines)

def make_code_tip(prs,title,lines,tip):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    tb(sl,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),title,fs=24,b=True,c=Th.O)
    code(sl,Inches(0.8),Inches(1.2),Inches(11.5),Inches(3.8),lines)
    rr(sl,Inches(0.8),Inches(5.3),Inches(11.5),Inches(1.2),Th.S)
    tb(sl,Inches(1.2),Inches(5.4),Inches(10.8),Inches(1),tip,fs=14,c=Th.C)

def make_ex(prs,title,desc,code_lines=None):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.G)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.6),title,fs=26,b=True,c=Th.G)
    rr(sl,Inches(0.8),Inches(1.2),Inches(11.5),Inches(3.5),Th.S)
    tx=tb_rich(sl,Inches(1.3),Inches(1.4),Inches(10.5),Inches(3.2))
    if isinstance(desc, list) and desc and isinstance(desc[-1], list):
        code_lines = desc[-1]
        desc = desc[:-1]
    txt = '\n'.join(desc) if isinstance(desc, list) else desc
    ap(tx.text_frame,txt,sz=15)
    if code_lines:
        code(sl,Inches(0.8),Inches(5),Inches(11.5),Inches(2.2),code_lines,fs=10)

def make_summary(prs,items):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(0.08),Inches(7.5),Th.O)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.7),"À retenir",fs=28,b=True,c=Th.O)
    for i,item in enumerate(items):
        rr(sl,Inches(1),Inches(1.5+i*0.9),Inches(11),Inches(0.65),Th.P)
        tb(sl,Inches(1.3),Inches(1.55+i*0.9),Inches(10.5),Inches(0.6),f"✓ {item}",fs=15,c=Th.LG)

def save(prs,name):
    p=os.path.join(OUT,name); prs.save(p); print(f"  ✅ {name}")

# ═══════════════════════════════════════════════════════════════
# PRESENTATIONS
# ═══════════════════════════════════════════════════════════════

def p27_modules_import():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Modules et Import","Utiliser et créer des bibliothèques","Partie 3 - Fonctions & Modules")
    make_obj(prs,["Comprendre ce qu'est un module","Importer avec import, from, as","Créer ses propres modules","Connaître le if __name__ == '__main__'"])
    make_content(prs,"Qu'est-ce qu'un module ?",[
        ("Un module = un fichier .py qui contient du code réutilisable.",15),
        "","Sans modules : tout dans un seul fichier → DÉSORDRE","Avec modules : code organisé en fichiers séparés → CLARTÉ","",
        ("Modules intégrés (built-in) :",0),"  • math → fonctions mathématiques","  • random → nombres aléatoires","  • datetime → dates et heures","  • os → système d'exploitation","  • sys → interpréteur Python","",
        ("Modules externes : installés avec pip",0),"","Modules personnels : tes propres fichiers .py"
    ])
    make_code_content(prs,"Syntaxes d'import",[
        "# 1. import simple", "import math", "print(math.sqrt(16))  # 4.0",
        "", "# 2. import avec alias", "import numpy as np", "print(np.array([1,2,3]))",
        "", "# 3. from ... import (sélectif)", "from math import pi, sqrt", "print(pi, sqrt(25))  # 3.14... 5.0",
        "", "# 4. from ... import * (⚠ DÉCONSEILLÉ)", "from math import *  # importe TOUT", "# Risque de conflit de noms !",
        "", "# 5. import d'une partie spécifique", "from datetime import datetime", "print(datetime.now())"
    ])
    make_code_content(prs,"Créer son propre module",[
        "# FICHIER: mes_outils.py", "def saluer(nom):", '    return f"Bonjour {nom}!"', "",
        "def additionner(a, b, c=0):", "    return a + b + c", "",
        "PI = 3.14159", "UTILISATEURS = []  # variable globale",
        "", "# FICHIER: main.py (autre fichier)", "import mes_outils",
        "print(mes_outils.saluer(\"Alice\"))", "print(mes_outils.additionner(3, 5))", "print(mes_outils.PI)",
        "", "# Ou avec from", "from mes_outils import saluer, PI", "print(saluer(\"Bob\"))", "print(PI)"
    ])
    make_code_content(prs,"Le __name__ == '__main__'",[
        "# FICHIER: salutations.py", "def bonjour(nom):", '    return f"Salut {nom}"', "",
        "def au_revoir(nom):", '    return f"À plus {nom}"', "",
        "# Test (s'exécute seulement si on lance CE fichier)", 'if __name__ == "__main__":', "    print(bonjour(\"Test\"))", "    print(au_revoir(\"Test\"))",
        "", "# Si on importe ce module AILLEURS, ces tests ne s'exécutent PAS", "",
        "# Pourquoi ?", "# __name__ vaut '__main__' quand on exécute le fichier directement", "# __name__ vaut le NOM du module quand on l'importe"
    ])
    make_content(prs,"Organisation des dossiers",[
        ("Bonne structure :",0),"","mon_projet/","├── main.py           # Point d'entrée","├── utils/             # Dossier = package","│   ├── __init__.py    # Rend le dossier importable","│   ├── maths.py       # Import: from utils.maths import ...","│   └── strings.py","├── data/              # Données","└── tests/             # Tests",
        "","Le fichier __init__.py peut être vide ou contenir du code d'initialisation","",
        ("Un dossier avec __init__.py = un PACKAGE (ensemble de modules)",0)
    ])
    make_code_content(prs,"Recharger un module sans redémarrer",[
        "# Pendant le développement", "import importlib", "import mes_outils",
        "", "# Modifier mes_outils.py...", "importlib.reload(mes_outils)  # recharge sans redémarrer",
        "", "# Lister tout ce que contient un module", "import math", "print(dir(math))  # liste toutes les fonctions",
        "", "# Voir le chemin d'un module", "print(math.__file__)  # chemin vers math.py"
    ])
    make_ex(prs,"Exercice : Crée ton module math",[
        "Crée un fichier mon_math.py avec :\n  1. Fonction factorielle(n) qui calcule n!\n  2. Fonction est_pair(n) qui retourne booléen\n  3. Fonction fibonacci(n) qui retourne les n premiers termes\n  4. Constante VITESSE_LUMIERE = 299792458\n  5. Bloc if __name__ == '__main__' pour tester\n\nPuis crée un test_math.py qui importe mon_math et :\n  - Calcule factorielle(5)\n  - Vérifie si 42 est pair\n  - Affiche les 10 premiers nombres de Fibonacci\n  - Affiche VITESSE_LUMIERE"
    ])
    make_summary(prs,["Un module = fichier .py réutilisable","import module, from module import func, import module as alias","if __name__ == '__main__' : test uniquement en exécution directe","Un dossier avec __init__.py = un package","Les modules rendent le code ORGANISÉ et RÉUTILISABLE"])
    save(prs,"27_Modules_Import.pptx")
    return prs

def p28_bibliotheques_utiles():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Bibliothèques Utiles","Découvrir les modules essentiels","Partie 3 - Fonctions & Modules")
    make_obj(prs,["Explorer les modules les plus courants","Apprendre random, datetime, math","Utiliser os, sys, json","Comprendre le choix du bon module"])
    make_content(prs,"Module random - Aléatoire",[
        ("Générer des valeurs aléatoires :",0),
        "","random.random() → float entre 0.0 et 1.0","random.randint(a,b) → entier entre a et b inclus","random.uniform(a,b) → float entre a et b","random.choice(liste) → élément aléatoire","random.sample(liste,k) → k éléments uniques","random.shuffle(liste) → mélange la liste","random.seed(n) → fixe le hasard (reproductible)",
        "","🔑 seed() permet de rejouer les mêmes tirages (tests, débogage)"
    ])
    make_code_content(prs,"random en action",[
        "import random",
        "", "# Jeu de dé", "de = random.randint(1, 6)", "print(f\"Tu as fait un {de}\")",
        "", "# Carte au hasard", "couleurs = [\"♥\",\"♦\",\"♣\",\"♠\"]", "valeurs = [\"A\",\"2\",\"3\",\"4\",\"5\",\"6\",\"7\",\"8\",\"9\",\"10\",\"J\",\"Q\",\"K\"]",
        "carte = random.choice(couleurs) + random.choice(valeurs)", "print(f\"Carte: {carte}\")",
        "", "# Mélanger une liste", "questions = [\"Q1\",\"Q2\",\"Q3\",\"Q4\"]", "random.shuffle(questions)", "print(questions)",
        "", "# Mot de passe aléatoire (4 chiffres)", "code = random.randint(1000, 9999)", "print(f\"Code: {code}\")"
    ])
    make_code_content(prs,"Module datetime - Dates et heures",[
        "from datetime import datetime, date, timedelta",
        "", "# Maintenant", "now = datetime.now()", "print(now)  # 2026-05-31 14:30:00...",
        "", "# Créer une date spécifique", "noel = date(2026, 12, 25)", "print(noel)",
        "", "# Formater une date", "print(now.strftime(\"%d/%m/%Y %H:%M\"))     # 31/05/2026 14:30", "print(now.strftime(\"%A %d %B %Y\"))     # Sunday 31 May 2026",
        "", "# Calculer un décalage", "demain = now + timedelta(days=1)", "hier = now - timedelta(days=1)", "print(f\"Demain: {demain}\")",
        "", "# Différence entre deux dates", "age = date.today() - date(2010, 6, 15)", "print(f\"{age.days} jours\")"
    ])
    make_content(prs,"Module math - Mathématiques",[
        ("Fonctions essentielles :",0),
        "","  math.sqrt(x)   → racine carrée","  math.ceil(x)    → arrondi supérieur","  math.floor(x)   → arrondi inférieur","  math.trunc(x)   → partie entière","  math.fabs(x)    → valeur absolue","  math.pow(x,y)   → x puissance y (équivalent à **)","  math.gcd(a,b)   → PGCD de a et b",
        "","Constantes :","  math.pi   → 3.141592...","  math.e    → 2.718281...","  math.inf  → infini","  math.nan  → Not a Number",
        "","Trigonométrie : sin(), cos(), tan(), radians(), degrees()"
    ])
    make_code_content(prs,"Module os et sys",[
        "import os, sys",
        "", "# os = interaction avec le système", "print(os.getcwd())          # dossier courant", "os.mkdir(\"nouveau\")         # crée un dossier (⚠ FileExistsError)", "os.makedirs(\"a/b/c\", exist_ok=True)  # crée une arborescence",
        "print(os.listdir(\".\"))      # liste les fichiers", "os.rename(\"old.txt\",\"new.txt\")  # renomme", "os.remove(\"fichier.txt\")    # supprime",
        "print(os.environ[\"HOME\"])  # variable d'environnement", "print(os.name)              # 'posix' sur Linux",
        "", "# os.path = opérations sur chemins", "print(os.path.join(\"a\",\"b.txt\"))  # a/b.txt", "print(os.path.exists(\"test.txt\"))  # booléen", "print(os.path.getsize(\"fichier.txt\"))  # taille en octets",
        "", "# sys = interpréteur Python", "print(sys.version)          # version Python", "print(sys.argv)             # arguments de la ligne de commande", "sys.exit(0)                 # quitte le programme"
    ])
    make_code_content(prs,"Module json - Données structurées",[
        "import json",
        "", "# JSON = format de stockage de données (très utilisé)", "",
        "# Python → JSON (sérialisation)", "data = {\"nom\": \"Alice\", \"age\": 25, \"notes\": [15, 18, 12]}", "json_str = json.dumps(data, indent=2, ensure_ascii=False)", "print(json_str)",
        "# Résultat:", '# {', '#   "nom": "Alice",', '#   "age": 25,', '#   "notes": [15, 18, 12]', '# }',
        "", "# JSON → Python (désérialisation)", "chargee = json.loads(json_str)", "print(chargee[\"nom\"])  # Alice",
        "", "# Lire/écrire dans un fichier", "with open(\"data.json\", \"w\", encoding=\"utf-8\") as f:", "    json.dump(data, f, indent=2)",
        "with open(\"data.json\", \"r\", encoding=\"utf-8\") as f:", "    data = json.load(f)"
    ])
    make_ex(prs,"Exercice : Générateur de citations",[
        "Crée un programme qui :\n1. Affiche une citation aléatoire parmi 10 dans une liste\n2. Affiche l'auteur (aléatoire aussi)\n3. Sauvegarde chaque citation affichée dans un fichier historique.json\n4. Au démarrage, demande si l'utilisateur veut :\n   a) Nouvelle citation aléatoire\n   b) Voir l'historique\n   c) Voir le nombre total de citations affichées\n5. Affiche la date et l'heure de chaque citation",
        ["import random, json", "from datetime import datetime", "CIT = [(\"La vie est belle\",\"Socrate\"),(\"Coder ou ne pas coder\",\"Moi\")]", "HIST = \"historique.json\"", "hist = json.load(open(HIST)) if os.path.exists(HIST) else []", "cit = random.choice(CIT)", "hist.append({\"c\":cit[0],\"a\":cit[1],\"d\":str(datetime.now())})", "json.dump(hist,open(HIST,\"w\"),indent=2)", "print(f'\"{cit[0]}\" - {cit[1]}')"]
    ])
    make_summary(prs,["random : aléatoire, choix, mélange","datetime : dates, heures, formatage, calculs","math : fonctions mathématiques et constantes","os : système de fichiers, chemins","sys : arguments, version, sortie","json : stockage structuré de données","Choisis TOUJOURS le module adapté à ton besoin"])
    save(prs,"28_Bibliotheques_utiles.pptx")
    return prs

def p29_mini_projet_gestionnaire_taches():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Mini-Projet : Gestionnaire de Tâches","Application console avec fichiers","Partie 3 - Fonctions & Modules")
    make_obj(prs,["Créer un vrai programme utile","Utiliser fichiers, modules, datetime","Apprendre à structurer un projet","Implémenter CRUD (Create, Read, Update, Delete)"])
    make_content(prs,"Le projet : Todo List",[
        ("On va créer un gestionnaire de tâches en ligne de commande :",15),
        "","Fonctionnalités :","  📝 AJOUTER une tâche (avec date limite)","  📋 LISTER les tâches (par statut)","  ✅ MARQUER comme terminée","  ❌ SUPPRIMER une tâche","  💾 Sauvegarder automatiquement dans un fichier JSON",
        "","Structure du projet :","taches/","├── main.py         # Point d'entrée (menu)","├── taches_manager.py  # Logique métier","├── taches_storage.py  # Sauvegarde/chargement","└── taches.json        # Données persistantes",
        "","C'est une ARCHITECTURE à 3 couches : affichage → logique → données"
    ])
    make_code_content(prs,"taches_storage.py - Stockage",[
        "# taches_storage.py", "import json, os",
        "FICHIER = \"taches.json\"",
        "", "def charger():", "    if not os.path.exists(FICHIER):", "        return []", "    with open(FICHIER, \"r\", encoding=\"utf-8\") as f:", "        return json.load(f)",
        "", "def sauvegarder(taches):", "    with open(FICHIER, \"w\", encoding=\"utf-8\") as f:", "        json.dump(taches, f, indent=2, ensure_ascii=False)",
        "", "def ajouter(taches, titre, date_limite=\"\"):", "    taches.append({", "        \"id\": len(taches) + 1,", '        "titre": titre,', '        "date_limite": date_limite,', '        "fait": False', "    })", "    sauvegarder(taches)"
    ])
    make_code_content(prs,"taches_manager.py - Logique",[
        "# taches_manager.py", "from datetime import datetime",
        "", "def afficher(taches, filtre=\"toutes\"):", '    """Affiche les tâches selon le filtre."""', "    for t in taches:", "        if filtre == \"faites\" and not t[\"fait\"]: continue", "        if filtre == \"en_cours\" and t[\"fait\"]: continue", "        statut = \"✓\" if t[\"fait\"] else \"◻\"", "        limite = f\" [{t['date_limite']}]\" if t['date_limite'] else \"\"", "        print(f\"{t['id']}. {statut} {t['titre']}{limite}\")",
        "", "def marquer_fait(taches, id_tache):", "    for t in taches:", "        if t[\"id\"] == id_tache:", "            t[\"fait\"] = not t[\"fait\"]", "            break", "    sauvegarder(taches)",
        "", "def supprimer(taches, id_tache):", "    taches[:] = [t for t in taches if t[\"id\"] != id_tache]", "    sauvegarder(taches)"
    ])
    make_code_content(prs,"main.py - Menu principal",[
        "# main.py (extrait)", "from taches_storage import charger, ajouter", "from taches_manager import afficher, marquer_fait, supprimer",
        "", "def menu():", "    taches = charger()", '    while (c := input(\"\\n[A]jouter [L]ister [✓]Faire [X]Sup [Q]uitter: \").upper()):',
        "        if c == \"A\":", '            t = input(\"Tâche: \")', '            d = input(\"Date limite (jj/mm): \")', "            ajouter(taches, t, d)",
        "        elif c == \"L\":", '            f = input(\"[T]outes [F]aites [E]n cours: \").upper()', "            afficher(taches, {\"T\":\"toutes\",\"F\":\"faites\",\"E\":\"en_cours\"}.get(f,\"toutes\"))",
        "        elif c == \"✓\":", "            afficher(taches, \"en_cours\")", "            id_ = int(input(\"ID à marquer: \"))", "            marquer_fait(taches, id_)",
        "        elif c == \"X\":", "            id_ = int(input(\"ID à supprimer: \"))", "            supprimer(taches, id_)",
        '        elif c == \"Q\": break',
        "", "if __name__ == '__main__':", "    menu()"
    ])
    make_content(prs,"Améliorations possibles",[
        ("Ton projet ne s'arrête pas là ! Ajoute :",0),
        "","🆕 Priorité (haute/moyenne/basse) avec des couleurs","🔍 Recherche par mot-clé","📊 Statistiques : combien faites/en cours","📅 Tâches en retard (date dépassée)","🔔 Rappel si une tâche est due aujourd'hui","📤 Export en CSV pour Excel",
        "","💡 Idée : transforme-le en app Web avec Flask (plus tard !)","","Conseil : développe PAR ÉTAPES, une fonctionnalité à la fois","        teste chaque ajout avant de passer à la suite"
    ])
    make_ex(prs,"À toi de jouer !",[
        "Objectif : Faire fonctionner le gestionnaire de tâches complet.\n\nÉtapes :\n1. Crée les 3 fichiers (storage, manager, main)\n2. Teste l'ajout et l'affichage\n3. Teste marquer comme fait\n4. Teste la suppression\n5. Ajoute au moins UNE amélioration personnelle\n\n⚠ N'oublie pas le if __name__ == '__main__' dans chaque fichier !\n\n📁 Tous les fichiers doivent être dans le même dossier."
    ])
    make_summary(prs,["CRUD = Create, Read, Update, Delete - les 4 opérations de base","Architecture 3 couches : présentation, logique, stockage","JSON pour la persistance des données","Modules séparés = code plus clair et maintenable","Un vrai projet se construit PAR ÉTAPES"])
    save(prs,"29_Mini_Projet_Gestionnaire_Taches.pptx")
    return prs

def p30_exercices_revision():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Exercices de Révision","Partie 3 - Fonctions, Fichiers, Modules","Partie 3 - Récapitulatif")
    make_obj(prs,["Réviser les fonctions et paramètres","Pratiquer fichiers et JSON","Renforcer les modules et imports","Combiner tous les concepts"])
    make_content(prs,"Exercice 1 : Analyseur de texte",[
        ("Crée un programme qui analyse un fichier texte :",15),
        "","1. Demande le nom du fichier à l'utilisateur","2. Affiche :","   • Nombre total de mots","   • Nombre de lignes","   • Nombre de caractères (sans espaces)","   • Le mot le plus long","   • La phrase la plus longue","3. Sauvegarde les résultats dans analyse_resultat.json",
        "","👣 Indices :",'  with open(nom, "r", encoding="utf-8") as f:', "  len(texte.split()) pour les mots", "  max(mots, key=len) pour le mot le plus long"
    ])
    make_content(prs,"Exercice 2 : Générateur de phrases",[
        ("Utilise random et des listes pour générer des phrases aléatoires :",15),
        "","sujets = [\"Le chat\", \"Un robot\", \"Ma soeur\", \"Le professeur\"]","verbes = [\"mange\", \"construit\", \"regarde\", \"programme\"]","complements = [\"une pizza\", \"un site web\", \"la télé\", \"un jeu vidéo\"]","lieux = [\"dans le jardin\", \"sur l'ordinateur\", \"à l'école\", \"dans la cuisine\"]",
        "","Génère 5 phrases aléatoires complètes :","\"Le chat mange une pizza dans le jardin.\"","\"Un robot programme un jeu vidéo sur l'ordinateur.\"",
        "","BONUS : Ajoute la date et sauvegarde dans phrases.txt","BONUS 2 : Assure-toi que les sujets/verbes s'accordent (singulier/pluriel)"
    ])
    make_content(prs,"Exercice 3 : Carnet de contacts",[
        ("Crée un carnet d'adresses avec JSON :",15),
        "","Structure : chaque contact a nom, téléphone, email","","Fonctions :","  • ajouter_contact(carnet)","  • rechercher_contact(carnet, nom)","  • lister_tous(carnet)","  • supprimer_contact(carnet, nom)","  • sauvegarder / charger depuis contacts.json",
        "","Le carnet est une liste de dictionnaires :","[{\"nom\":\"Alice\", \"tel\":\"06...\", \"email\":\"a@b.com\"}, ...]",
        "","💡 Organise en 2 fichiers : carnet.py (logique) et main_carnet.py (menu)"
    ])
    make_content(prs,"Exercice 4 : Statistiques",[
        ("Crée un module stats.py avec ces fonctions :",15),
        "","def moyenne(liste):","    return sum(liste) / len(liste)",
        "","def mediane(liste):","    liste_triee = sorted(liste)","    n = len(liste_triee)","    # Si n impair → milieu, si pair → moyenne des 2 du milieu",
        "","def ecart_type(liste):","    # sqrt(moyenne des (x - moyenne)²)",
        "","def nettoyer(liste):","    # Supprime les None, les strings, les négatifs",
        "        ","Puis crée un test_stats.py qui :","• Utilise les fonctions avec [15, 12, 18, 7, 16, None, 20, \"abc\"]","• Affiche moyenne, médiane, écart-type après nettoyage"
    ])
    make_ex(prs,"Exercice 5 : Mini-jeu de devinettes",[
        "Crée un jeu complet avec modules :\n\n1. Un module jeu.py avec :\n   - choisir_mot() → depuis une liste de 20 mots\n   - melanger(mot) → mélange les lettres du mot\n   - verifier(proposition, mot) → retourne booléen\n\n2. Un module scores.py avec :\n   - charger_scores() depuis scores.json\n   - sauvegarder_scores(scores)\n   - ajouter_score(pseudo, points)\n   - top_scores(n=5) → meilleurs scores\n\n3. Un main.py qui :\n   - Affiche le mot mélangé\n   - Donne 3 essais\n   - Compte les points\n   - Affiche le top 5 à la fin",
        ["# Mélange de lettres", "import random", "def melanger(mot):", "    l = list(mot)", "    random.shuffle(l)", "    return ''.join(l)"]
    ])
    make_summary(prs,["Fonctions : def, return, paramètres par défaut, *args/**kwargs","Portée : LEGB (Local, Enclosing, Global, Built-in)","Fichiers : open/with, read/write/append, JSON","Modules : import, if __name__, packages","Tuples (immuable), Sets (sans doublons, opérateurs mathématiques)","Organisation : séparer logique, stockage, présentation"])
    save(prs,"30_Exercices_Revision.pptx")
    return prs

def p31_poo_classes():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Programmation Orientée Objet","Classes et Objets en Python","Partie 4 - POO & Git")
    make_obj(prs,["Comprendre le concept d'objet","Créer une classe","Utiliser __init__ et self","Définir des méthodes et attributs"])
    make_content(prs,"Qu'est-ce qu'un objet ?",[
        ("Jusqu'à présent : programmation PROCÉDURALE (fonctions + données)",15),
        "","Données = variables classiques","Fonctions = opérations sur les données","",
        ("Avec la POO, on regroupe données ET fonctions dans un OBJET",0),
        "","🎯 Un objet = une « chose » qui a :","  • des ATTRIBUTS (caractéristiques, données)","  • des MÉTHODES (actions, fonctions)",
        "","Exemples concrets :","  • Chien : nom, âge, race / aboyer(), manger(), dormir()","  • Voiture : marque, couleur, vitesse / demarrer(), freiner(), klaxonner()","  • CompteBancaire : solde, titulaire / deposer(), retirer(), afficher()",
        ("Une CLASSE = le plan/blueprint. Un OBJET = une instance du plan.",0)
    ])
    make_code_content(prs,"Créer une classe - Syntaxe de base",[
        "# Définir une classe (convention : Nom en PascalCase)", "class Chien:", "    pass  # classe vide",
        "", "# Créer un objet (instance)", "mon_chien = Chien()", "print(type(mon_chien))  # <class '__main__.Chien'>",
        "", "# Ajouter des attributs (peu pratique)", "mon_chien.nom = \"Rex\"", "mon_chien.age = 3",
        "", "# MÉTHODE 1 : définir dans la classe", "class Chien:", "    def aboyer(self):", '        print("Woof!")',
        "", "rex = Chien()", "rex.aboyer()  # Woof!"
    ])
    make_code_content(prs,"Le constructeur __init__",[
        "# __init__ s'exécute AUTOMATIQUEMENT à la création", "class Chien:",
        "    def __init__(self, nom, age, race):", "        self.nom = nom     # attribut d'instance", "        self.age = age", "        self.race = race",
        "    def aboyer(self):", '        print(f"{self.nom} dit: Woof!")',
        "    def presentation(self):", '        print(f"Je suis {self.nom}, un {self.race} de {self.age} ans")',
        "", "# Création d'objets", "rex = Chien(\"Rex\", 3, \"Berger\")", "luna = Chien(\"Luna\", 1, \"Labrador\")",
        "rex.aboyer()  # Rex dit: Woof!", "luna.presentation()  # Je suis Luna, un Labrador de 1 ans", "",
        "self = l'objet lui-même (comme \"moi\" pour l'objet)"
    ])
    make_content(prs,"Attributs et Méthodes",[
        ("Types d'attributs :",0),
        "","  Attributs d'instance : propres à chaque objet (self.nom)","  Attributs de classe : partagés par tous les objets",
        "","  class Chien:", "      espece = \"Canis familiaris\"  # attribut de classe",
        "      def __init__(self, nom):", "          self.nom = nom  # attribut d'instance",
        "","  rex = Chien(\"Rex\"); luna = Chien(\"Luna\")", "  print(rex.espece)  # Canis familiaris (commun)", "  print(rex.nom)     # Rex (propre)",
        ("Types de méthodes :",0),
        "  Méthodes d'instance : prennent self, agissent sur l'objet","  Méthodes de classe (@classmethod) : prennent cls","  Méthodes statiques (@staticmethod) : ni self ni cls"
    ])
    make_code_content(prs,"Exemple complet : CompteBancaire",[
        "class CompteBancaire:",
        '    """Représente un compte en banque."""',
        "    frais_retrait = 0.50  # attribut de classe",
        "    nb_comptes = 0  # compteur global",
        "",
        "    def __init__(self, titulaire, solde=0):", "        self.titulaire = titulaire", "        self.solde = solde", "        self.operations = []  # historique", "        CompteBancaire.nb_comptes += 1",
        "    def deposer(self, montant):", "        self.solde += montant", '        self.operations.append(f"+{montant}")', '        print(f"Déposé {montant}€")',
        "    def retirer(self, montant):", "        total = montant + CompteBancaire.frais_retrait", "        if total > self.solde: return False", "        self.solde -= total", '        self.operations.append(f"-{montant} (frais 0.50)")', "        return True",
        "    def afficher(self):", '        print(f"{self.titulaire}: {self.solde}€")',
        "", "c = CompteBancaire(\"Alice\", 100)", "c.deposer(50); c.retirer(20); c.afficher()"
    ])
    make_code_content(prs,"Méthodes magiques __str__ et __repr__",[
        "class Livre:",
        "    def __init__(self, titre, auteur):", "        self.titre = titre", "        self.auteur = auteur",
        "", "    def __str__(self):", '        return f"{self.titre} par {self.auteur}"',
        "    def __repr__(self):", '        return f"Livre({self.titre!r}, {self.auteur!r})"',
        "", "l = Livre(\"1984\", \"Orwell\")", "print(str(l))    # 1984 par Orwell", "print(repr(l))   # Livre('1984', 'Orwell')",
        "", "# Autres méthodes magiques utiles:", "# __len__(self)  → len(objet)", "# __add__(self, other)  → objet1 + objet2", "# __eq__(self, other)  → objet1 == objet2", "# __lt__(self, other)  → objet1 < objet2"
    ])
    make_ex(prs,"Exercice : Classe Étudiant",[
        "Crée une classe Etudiant avec :\n  • Attributs : nom, age, note (liste de notes)\n  • Méthodes :\n    - ajouter_note(note) → ajoute une note (0-20)\n    - moyenne() → retourne la moyenne\n    - mention() → \"Excellent\"(≥17), \"Bien\"(≥15), \"Assez bien\"(≥13), \"Passable\"(≥11), \"Insuffisant\"\n    - __str__() → affiche \"nom: moy/20 (mention)\"\n\nPuis :\n  • Crée 3 étudiants\n  • Ajoute des notes à chacun\n  • Affiche le meilleur étudiant\n  • Affiche la moyenne de la classe"
    ])
    make_summary(prs,["Classe = plan, Objet = réalisation concrète (instance)","__init__ = constructeur, s'exécute à la création","self = l'objet lui-même","Attribut d'instance = self.x, Attribut de classe = Classe.x","Méthodes = fonctions attachées à un objet","__str__ = représentation lisible, __repr__ = représentation technique","La POO organise le code en entités du monde réel"])
    save(prs,"31_POO_Classes.pptx")
    return prs

def p32_poo_heritage():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"POO - Héritage et Polymorphisme","Réutiliser et étendre des classes","Partie 4 - POO & Git")
    make_obj(prs,["Comprendre l'héritage","Utiliser super()", "Appliquer le polymorphisme","Connaître l'encapsulation"])
    make_content(prs,"Principe de l'héritage",[
        ("L'héritage permet à une CLASSE FILLE d'hériter des attributs et méthodes d'une CLASSE MÈRE.",15),
        "","🔑 Pourquoi hériter ?","  • RÉUTILISER du code sans le copier","  • AJOUTER des fonctionnalités spécifiques","  • CRÉER une hiérarchie de classes",
        "","Exemple concret :","  Animal (classe mère) → attributs : nom, age / méthode : parler()","    ├── Chien → hérite + méthode aboyer()","    ├── Chat → hérite + méthode miauler()","    └── Oiseau → hérite + méthode voler()",
        "","Chaque classe fille est un ANIMAL mais avec ses PARTICULARITÉS","","Syntaxe : class Chien(Animal):  # Chien hérite de Animal"
    ])
    make_code_content(prs,"Héritage simple - Syntaxe",[
        "# Classe mère (parent)", "class Animal:", "    def __init__(self, nom):", "        self.nom = nom",
        "    def manger(self):", '        print(f"{self.nom} mange")',
        "    def dormir(self):", '        print(f"{self.nom} dort")',
        "", "# Classe fille qui hérite de Animal", "class Chien(Animal):", "    def aboyer(self):", '        print(f"{self.nom} aboie: Woof!")',
        "class Chat(Animal):", "    def miauler(self):", '        print(f"{self.nom} miaule: Miaou!")',
        "", "# Utilisation", "rex = Chien(\"Rex\")", "rex.manger()   # hérité → Rex mange", "rex.aboyer()   # propre → Rex aboie: Woof!",
        "minou = Chat(\"Minou\")", "minou.dormir()   # hérité", "minou.miauler() # propre"
    ])
    make_code_content(prs,"Le mot-clé super()",[
        "super() = appelle la méthode de la CLASSE MÈRE",
        "", "class Animal:", "    def __init__(self, nom):", "        self.nom = nom",
        "", "class Chien(Animal):", "    def __init__(self, nom, race):", "        super().__init__(nom)  # appelle Animal.__init__(self, nom)", "        self.race = race",
        "", "    def __str__(self):", '        return f"{self.nom} ({self.race})"',
        "", "# Explication :", "# Sans super(), on devrait réécrire: self.nom = nom", "# Avec super(), on RÉUTILISE le code du parent",
        "", "rex = Chien(\"Rex\", \"Berger\")", "print(rex)  # Rex (Berger)"
    ])
    make_content(prs,"Polymorphisme",[
        ("Poly-morphisme = plusieurs formes. Même méthode, comportements différents.",15),
        "","  class Animal:", "      def faire_son(self):", "          pass  # à redéfinir",
        "  class Chien(Animal):", "      def faire_son(self): return \"Woof\"",
        "  class Chat(Animal):", "      def faire_son(self): return \"Miaou\"",
        "  class Vache(Animal):", "      def faire_son(self): return \"Meuh\"",
        "","  animaux = [Chien(), Chat(), Vache()]", "  for a in animaux:", "      print(a.faire_son())  # Woof, Miaou, Meuh",
        "","On appelle la MÊME méthode faire_son() mais chaque animal répond DIFFÉREMMENT.","","C'est très puissant pour traiter des objets différents de façon uniforme !"
    ])
    make_content(prs,"Encapsulation (Portée des attributs)",[
        ("Contrôler l'accès aux attributs :",15),
        "","  Public    : self.nom       → accessible partout (par défaut)","  Protégé   : self._nom      → _ = \"usage interne\" (convention)","  Privé     : self.__nom     → __ = name mangling (Python le rend moins accessible)",
        "","  class CompteBancaire:", "      def __init__(self, solde):", "          self.__solde = solde  # attribut privé",
        "      def get_solde(self):", '          """Getter : accès contrôlé"""', "          return self.__solde",
        "      def deposer(self, montant):", "          if montant > 0:", "              self.__solde += montant",
        "","  c = CompteBancaire(100)", "  # print(c.__solde)  # ERREUR ! (attribut privé)", "  print(c.get_solde())  # 100 ✅"
    ])
    make_code_content(prs,"Exemple complet : Formes géométriques",[
        "import math",
        "class Forme:", "    def aire(self): return 0", "    def description(self): return \"Forme géométrique\"",
        "", "class Carre(Forme):", "    def __init__(self, cote): self.cote = cote", "    def aire(self): return self.cote ** 2", '    def description(self): return f"Carré de côté {self.cote}"',
        "", "class Cercle(Forme):", "    def __init__(self, rayon): self.rayon = rayon", "    def aire(self): return math.pi * self.rayon ** 2", '    def description(self): return f"Cercle de rayon {self.rayon}"',
        "", "class Triangle(Forme):", "    def __init__(self, base, hauteur): self.base = base; self.hauteur = hauteur", "    def aire(self): return 0.5 * self.base * self.hauteur",
        "", "# Polymorphisme en action", "formes = [Carre(4), Cercle(3), Triangle(5, 2)]", "for f in formes:", "    print(f\"{f.description()} → aire = {f.aire():.2f}\")"
    ])
    make_ex(prs,"Exercice : Système d'employés",[
        "Crée une classe Employe avec :\n  • Attributs : nom, salaire\n  • Méthode : travailler() → \"Employé travaille...\"\n\nCrée 3 classes filles :\n  1. Developpeur → travailler() → \"Développeur écrit du code\"\n  2. Designer → travailler() → \"Designer crée des maquettes\"\n  3. Manager → __init__ ajoute equipe=[], travailler() → \"Manager supervise\"\n\nAjoute :\n  • Méthode embaucher(employe) à Manager\n  • Un affichage liste des employés de tous types\n  • Calcule le salaire total de l'entreprise"
    ])
    make_summary(prs,["Héritage : class Fille(Parent) → réutilise le code","super() = appelle la méthode parent","Polymorphisme : même méthode, comportement différent selon la classe","Encapsulation : _protégé, __privé (name mangling)","Getter/Setter pour contrôler l'accès aux attributs","La POO permet de modéliser le monde réel proprement"])
    save(prs,"32_POO_Heritage.pptx")
    return prs

def p33_comprehensions():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Compréhensions et Expressions Avancées","Écrire du code concis et élégant","Partie 4 - POO & Git")
    make_obj(prs,["Comprendre les list comprehensions","Maîtriser les dict et set comprehensions","Utiliser les expressions conditionnelles","Améliorer la lisibilité du code"])
    make_content(prs,"Pourquoi des comprehensions ?",[
        ("Besoin courant : créer une liste à partir d'une autre liste.",15),
        "","SANS compréhension (classique) :","  carres = []", "  for x in range(10):", "      carres.append(x**2)",
        "","AVEC compréhension (1 ligne !) :","  carres = [x**2 for x in range(10)]",
        "","👉 Plus court, plus clair, plus rapide (Python optimise)","",
        ("Syntaxe générale :",0),"  [expression for élément in iterable [if condition]]",
        "","C'est comme un mini-langage dans Python pour transformer des séquences."
    ])
    make_code_content(prs,"List Comprehensions",[
        "# Structure: [expression for élément in série]",
        "", "# 1. Simple", "[x**2 for x in range(6)]  # [0, 1, 4, 9, 16, 25]",
        "", "# 2. Avec filtre (if)", "[x for x in range(20) if x % 2 == 0]  # pairs", "[x for x in range(20) if x % 3 == 0]  # multiples de 3",
        "", "# 3. Avec transformation", 'noms = ["alice", "bob", "charlie"]', '[n.capitalize() for n in noms]  # ["Alice", "Bob", "Charlie"]',
        "", "# 4. Double boucle", '[(x, y) for x in range(3) for y in range(3)]', "# [(0,0),(0,1),(0,2),(1,0),(1,1),(1,2),(2,0),(2,1),(2,2)]",
        "", "# 5. Flatten une liste 2D", "matrice = [[1,2],[3,4],[5,6]]", "[n for ligne in matrice for n in ligne]  # [1,2,3,4,5,6]"
    ])
    make_code_content(prs,"Dict et Set Comprehensions",[
        "# DICT COMPREHENSION: {clé: valeur for ...}",
        "", "# Carrés en dictionnaire", "{x: x**2 for x in range(6)}  # {0:0, 1:1, 2:4, 3:9, 4:16, 5:25}",
        "", "# Inverser clés/valeurs", "original = {\"a\": 1, \"b\": 2, \"c\": 3}", "inverse = {v: k for k, v in original.items()}", "print(inverse)  # {1:\"a\", 2:\"b\", 3:\"c\"}",
        "", "# Filtrer un dict", "notes = {\"Alice\":15, \"Bob\":8, \"Charlie\":18, \"David\":10}", "adm = {n: m for n, m in notes.items() if m >= 12}", "print(adm)  # {\"Alice\":15, \"Charlie\":18}",
        "", "# SET COMPREHENSION: {expression for ...}", 'phrase = "le chat et le chien"', "{c for c in phrase if c != \" \"}  # {'l','e','c','h','a','t','n','i'}",
        "", "# Doublons automatiquement supprimés !", "{x % 3 for x in range(10)}  # {0, 1, 2}"
    ])
    make_code_content(prs,"Expressions conditionnelles (ternaire)",[
        "# Syntaxe: valeur_si_vrai if condition else valeur_si_faux",
        "", "# Dans une compréhension", "notes = [12, 8, 15, 4, 18, 6]",
        'resultats = ["✅ Admis" if n >= 10 else "❌ Échec" for n in notes]',
        "# ['✅ Admis','❌ Échec','✅ Admis','❌ Échec','✅ Admis','❌ Échec']",
        "",
        "# Avec des opérations", "x = 10", 'parite = "pair" if x % 2 == 0 else "impair"',
        "", "# Combinaison filtrer + transformer + ternaire", "nombres = range(-5, 6)",
        '[f"{n} positif" if n > 0 else f"{n} nul" if n == 0 else f"{n} négatif" for n in nombres]'
    ])
    make_code_content(prs,"Astuces et performances",[
        "# 1. Comprehension vs map/filter (souvent plus lisible)",
        "# Comprehension", "[x**2 for x in range(10) if x > 5]",
        "# Equivalent map/filter", "list(map(lambda x: x**2, filter(lambda x: x > 5, range(10))))",
        "",
        "# 2. Tuple comprehension ?", "# C'est un GÉNÉRATEUR, pas un tuple !", "gen = (x**2 for x in range(5))", "print(gen)  # <generator object>", "print(list(gen))  # [0, 1, 4, 9, 16]",
        "",
        "# 3. Pour des listes TRÈS grandes, préférer les générateurs",
        "# Liste : prend toute la mémoire", "l = [x**2 for x in range(10**6)]  # ~8 Mo",
        "# Générateur : calcule à la demande", "g = (x**2 for x in range(10**6))  # ~0 Mo",
        "# next(g) donne le suivant, ou for val in g:",
        "",
        "# 4. NE PAS abuser : une compréhension > 2 lignes mérite une boucle"
    ])
    make_ex(prs,"Exercice : Comprehensions",[
        "1. Liste des carrés des nombres pairs entre 1 et 30",
        "2. Dict : mot → longueur pour une phrase donnée",
        "3. Liste des nombres premiers entre 1 et 50 (avec fonction est_premier)",
        "4. Depuis une liste de températures en °C, crée la liste en °F (F = C*9/5+32)",
        "5. Depuis une phrase, trouve les mots de plus de 5 lettres",
        "6. Set des voyelles présentes dans une phrase",
        "7. (BONUS) Matrice identité 4×4 en compréhension",
        "8. (BONUS) Depuis une liste de notes, crée {'admis':[...], 'echec':[...]}",
        "",
        ['# Aide pour 7:', '[[1 if i==j else 0 for j in range(4)] for i in range(4)]']
    ])
    make_summary(prs,["[expr for x in iterable] = list comprehension","{k:v for ...} = dict comprehension", "{x for ...} = set comprehension", "Condition : [x for x in s if cond], Ternaire : a if cond else b", "Générateur (...) = paresseux, économe en mémoire", "Plus CONCIS, plus CLAIR, plus RAPIDE"])
    save(prs,"33_Comprehensions.pptx")
    return prs

def p34_lambda_map_filter():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Fonctions Lambda et map/filter","Écrire des fonctions anonymes","Partie 4 - POO & Git")
    make_obj(prs,["Comprendre les fonctions lambda","Utiliser map() et filter()","Maîtriser la fonction sorted() avec key","Connaître reduce() et functools"])
    make_content(prs,"Fonctions Lambda",[
        ("Une lambda = une fonction ANONYME (sans nom), sur UNE SEULE LIGNE.",15),
        "","Syntaxe : lambda paramètres: expression",
        "","Équivalent à :","  def ma_fonction(x): return x * 2",
        "  lambda x: x * 2",
        "","Caractéristiques :","  • Pas de return (implicite)","  • Pas de nom","  • Pas de boucles ni if/elif complexes","  • UTILE comme argument à une autre fonction",
        "","Quand l'utiliser ?","  • Pour une opération simple et ponctuelle","  • Comme argument de map(), filter(), sorted()","  • Pour un tri personnalisé",
        ("⚠ Pas de return, pas d'annotations, pas de docstrings !",0)
    ])
    make_code_content(prs,"Lambda en pratique",[
        "# Lambda simple", "carre = lambda x: x ** 2", "print(carre(5))  # 25",
        "", "# Plusieurs paramètres", "addition = lambda a, b: a + b", "print(addition(3, 7))  # 10",
        "", "# Lambda avec condition (ternaire)", "max2 = lambda a, b: a if a > b else b",
        "", "# En pratique, on utilise lambda DIRECTEMENT sans la stocker", "# (c'est son intérêt !)",
        "", "# Exemple : tri personnalisé", "etudiants = [(\"Alice\",15),(\"Bob\",8),(\"Charlie\",18)]", "tries = sorted(etudiants, key=lambda e: e[1], reverse=True)", "print(tries)  # [('Charlie',18),('Alice',15),('Bob',8)]",
        "", "# Trier des dictionnaires", "livres = [{\"titre\":\"Z\",\"pages\":200},{\"titre\":\"A\",\"pages\":50}]", 'tries = sorted(livres, key=lambda l: l["pages"])'
    ])
    make_code_content(prs,"map() - Appliquer une fonction à chaque élément",[
        "# map(fonction, itérable) → retourne un itérateur",
        "", "# 1. Transformer des nombres", "notes = [8, 12, 15, 6, 18]",
        "# Sans map", "[n + 1 for n in notes]",
        "# Avec map", "list(map(lambda n: n + 1, notes))  # [9, 13, 16, 7, 19]",
        "",
        "# 2. Convertir des types", "chiffres = [\"1\", \"2\", \"3\", \"4\"]", "list(map(int, chiffres))  # [1, 2, 3, 4]",
        "",
        "# 3. Avec plusieurs listes", "a = [1, 2, 3]; b = [10, 20, 30]", "list(map(lambda x, y: x + y, a, b))  # [11, 22, 33]",
        "",
        "# 4. map + méthode de string", 'noms = ["alice", "bob", "charlie"]', 'list(map(str.capitalize, noms))  # ["Alice", "Bob", "Charlie"]'
    ])
    make_code_content(prs,"filter() - Garder les éléments qui satisfont une condition",[
        "# filter(fonction_booleenne, itérable)",
        "", "# 1. Garder les pairs", "nombres = range(1, 21)", "list(filter(lambda x: x % 2 == 0, nombres))  # [2,4,6,8,10,12,14,16,18,20]",
        "",
        "# 2. Filtrer les mots courts", 'mots = ["python", "js", "java", "c", "ruby", "rust"]', 'list(filter(lambda m: len(m) >= 4, mots))  # ["python", "java", "ruby", "rust"]',
        "",
        "# 3. Filtrer les None", "valeurs = [1, None, 3, None, 5, 0, [], \"\"]", "list(filter(None, valeurs))  # [1, 3, 5] (None, 0, [], \"\" sont falsy)",
        "",
        "# 4. Combiner map + filter", "notes = [12, 8, 15, 4, 18, 6]",
        "# Filtrer les admis puis ajouter bonus", "list(map(lambda n: n + 1, filter(lambda n: n >= 10, notes)))",
        "# Résultat: [13, 16, 19]",
        "",
        "# Equivalent comprehension (souvent plus lisible)", "[n + 1 for n in notes if n >= 10]"
    ])
    make_code_content(prs,"reduce() et autres applications",[
        "from functools import reduce",
        "", "# reduce(fonction, itérable) → accumule un résultat",
        "", "# 1. Somme (mieux avec sum())", "reduce(lambda a, b: a + b, [1, 2, 3, 4, 5])  # 15",
        "",
        "# 2. Produit", "reduce(lambda a, b: a * b, range(1, 6))  # 120 (5!)",
        "",
        "# 3. Maximum (mieux avec max())", "reduce(lambda a, b: a if a > b else b, [3, 7, 2, 9, 1])  # 9",
        "",
        "# 4. Concaténer des chaînes", 'reduce(lambda a, b: f"{a}-{b}", ["a","b","c","d"])  # "a-b-c-d"',
        "",
        "# ATTENTION :", "# reduce() est UTILE mais parfois moins lisible", "# Préférer des boucles ou sum()/max() quand c'est simple"
    ])
    make_ex(prs,"Exercice : Lambda et map/filter",[
        "1. Avec une liste de mots, crée une liste des longueurs (map)",
        "2. Filtre les mots qui contiennent la lettre 'a' (filter)",
        "3. Trier des tuples (nom, age) par âge décroissant (sorted + lambda key)",
        "4. Avec reduce, calcule le factoriel de 7",
        "5. Depuis une liste de prix HT, applique TVA 20% (map)",
        "6. Trouve les nombres divisibles par 3 ET 5 entre 1 et 100 (filter)",
        "7. (BONUS) Trier une phrase par longueur de mot",
        "8. (BONUS) Grouper des mots par première lettre avec un dict",
        "",
        ["# Aide 7:", 'phrase = "le python est un langage puissant"', "mots = phrase.split()", "trie = sorted(mots, key=lambda m: len(m))"]
    ])
    make_summary(prs,["lambda params: expr → fonction anonyme sur une ligne","map(f, iter) → applique f à chaque élément","filter(f, iter) → garde les éléments où f retourne True","reduce(f, iter) → accumule (functools)","sorted(iter, key=lambda ...) → tri personnalisé","Compréhension [x for x in ... if ...] souvent plus lisible"])
    save(prs,"34_Lambda_map_filter.pptx")
    return prs

if __name__ == "__main__":
    funcs = [p27_modules_import, p28_bibliotheques_utiles, p29_mini_projet_gestionnaire_taches,
             p30_exercices_revision, p31_poo_classes, p32_poo_heritage,
             p33_comprehensions, p34_lambda_map_filter]
    for f in funcs:
        print(f"\n Génération : {f.__name__}")
        try:
            f()
        except Exception as e:
            print(f"  ❌ Erreur dans {f.__name__}: {e}")
    print(f"\n{'='*50}\n✅ {len(funcs)} présentations générées !\n📁 Dossier : {OUT}\n{'='*50}")
