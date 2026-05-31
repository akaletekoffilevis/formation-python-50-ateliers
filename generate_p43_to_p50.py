#!/usr/bin/env python3
"""Generate presentations 43 to 50 - Part 5."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

class Th:
    D=RGBColor(0x0A,0x0A,0x2E); P=RGBColor(0x16,0x21,0x3E)
    S=RGBColor(0x1A,0x1A,0x4E); O=RGBColor(0xFF,0x6B,0x35)
    C=RGBColor(0x00,0xB4,0xD8); PU=RGBColor(0x7B,0x2F,0xF7)
    G=RGBColor(0x06,0xD6,0xA0); W=RGBColor(0xFF,0xFF,0xFF)
    L=RGBColor(0xF8,0xF9,0xFA); GR=RGBColor(0x6C,0x75,0x7D)
    LG=RGBColor(0xDE,0xE2,0xE6); CB=RGBColor(0x1E,0x1E,0x2E)
    CT=RGBColor(0xCD,0xD6,0xF4); R=RGBColor(0xEF,0x44,0x44)

OUT = "/home/akaletekoffilevis/Bureau/Coach/presentations"
os.makedirs(OUT, exist_ok=True)

def bg(sl,c):
    f=sl.background.fill; f.solid(); f.fore_color.rgb=c

def sh(sl,l,t,w,h,fc,bc=None):
    s=sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=fc
    s.line.fill.background()
    if bc: s.line.color.rgb=bc; s.line.width=Pt(1)
    return s

def rr(sl,l,t,w,h,fc):
    s=sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,l,t,w,h); s.fill.solid(); s.fill.fore_color.rgb=fc; s.line.fill.background()
    return s

def tb(sl,l,t,w,h,text,fs=14,c=None,b=False,a=PP_ALIGN.LEFT):
    tx=sl.shapes.add_textbox(l,t,w,h); tf=tx.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.text=text; p.font.size=Pt(fs); p.font.color.rgb=c or Th.W
    p.font.bold=b; p.font.name='Calibri'; p.alignment=a
    return tx

def tb_rich(sl,l,t,w,h):
    tx=sl.shapes.add_textbox(l,t,w,h); tf=tx.text_frame; tf.word_wrap=True; return tx

def ap(tf,text,sz=14,c=None,b=False,al=PP_ALIGN.LEFT,sb=0,sa=0,name='Calibri'):
    if len(tf.paragraphs)==1 and tf.paragraphs[0].text=='': p=tf.paragraphs[0]
    else: p=tf.add_paragraph()
    p.text=text; p.font.size=Pt(max(sz,8)); p.font.color.rgb=c or Th.W; p.font.bold=b
    p.font.name=name; p.alignment=al; p.space_before=Pt(sb); p.space_after=Pt(sa)
    return p

def code(sl,l,t,w,h,lines,fs=11):
    rr(sl,l,t,w,h,Th.CB)
    tx=tb_rich(sl,l+Inches(0.3),t+Inches(0.15),w-Inches(0.6),h-Inches(0.3))
    for line in lines:
        ap(tx.text_frame,line,sz=fs,c=Th.CT,name='Consolas',al=PP_ALIGN.LEFT)

def make_title(prs,title,sub,part):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    rr(sl,Inches(1),Inches(1.5),Inches(11.333),Inches(4.5),Th.P)
    tb(sl,Inches(1.5),Inches(1.8),Inches(10.333),Inches(1),part,fs=12,c=Th.O,b=True)
    tb(sl,Inches(1.5),Inches(2.5),Inches(10.333),Inches(1.8),title,fs=36,b=True)
    tb(sl,Inches(1.5),Inches(4.2),Inches(10.333),Inches(1),sub,fs=18,c=Th.GR)

def make_obj(prs,items):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(0.08),Inches(7.5),Th.C)
    tb(sl,Inches(0.8),Inches(0.5),Inches(11),Inches(0.8),"Objectifs",fs=28,b=True,c=Th.C)
    for i,item in enumerate(items):
        rr(sl,Inches(1),Inches(1.6+i*1.1),Inches(11),Inches(0.8),Th.S)
        tb(sl,Inches(1.3),Inches(1.7+i*1.1),Inches(10.5),Inches(0.7),f"✦ {item}",fs=16)

def make_content(prs,title,lines):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.7),title,fs=26,b=True,c=Th.O)
    tx=tb_rich(sl,Inches(0.8),Inches(1.4),Inches(11.5),Inches(5.5))
    for line in lines:
        t,sz=line if isinstance(line,tuple) else (line,15)
        ap(tx.text_frame,t,sz=sz)

def make_code_content(prs,title,lines):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.O)
    tb(sl,Inches(0.8),Inches(0.3),Inches(11),Inches(0.6),title,fs=24,b=True,c=Th.O)
    code(sl,Inches(0.8),Inches(1.2),Inches(11.5),Inches(5.5),lines)

def make_ex(prs,title,desc,code_lines=None):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.P)
    sh(sl,Inches(0),Inches(0),Inches(13.333),Inches(0.06),Th.G)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.6),title,fs=26,b=True,c=Th.G)
    rr(sl,Inches(0.8),Inches(1.2),Inches(11.5),Inches(3.5),Th.S)
    tx=tb_rich(sl,Inches(1.3),Inches(1.4),Inches(10.5),Inches(3.2))
    if isinstance(desc, list) and desc and isinstance(desc[-1], list):
        code_lines = desc[-1]
        desc = desc[:-1]
    txt = '\n'.join(desc) if isinstance(desc, list) else desc
    ap(tx.text_frame,txt,sz=15)
    if code_lines:
        code(sl,Inches(0.8),Inches(5),Inches(11.5),Inches(2.2),code_lines,fs=10)

def make_summary(prs,items):
    sl=prs.slides.add_slide(prs.slide_layouts[6]); bg(sl,Th.D)
    sh(sl,Inches(0),Inches(0),Inches(0.08),Inches(7.5),Th.O)
    tb(sl,Inches(0.8),Inches(0.4),Inches(11),Inches(0.7),"À retenir",fs=28,b=True,c=Th.O)
    for i,item in enumerate(items):
        rr(sl,Inches(1),Inches(1.5+i*0.9),Inches(11),Inches(0.65),Th.P)
        tb(sl,Inches(1.3),Inches(1.55+i*0.9),Inches(10.5),Inches(0.6),f"✓ {item}",fs=15,c=Th.LG)

def save(prs,name):
    p=os.path.join(OUT,name); prs.save(p); print(f"  ✅ {name}")

# ═══════════════════════════════════════════════════════════════
# PRESENTATIONS
# ═══════════════════════════════════════════════════════════════

def p43_mini_projet_meteo():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Mini-Projet : Application Météo","API + Python = données temps réel","Partie 5 - Projets Avancés")
    make_obj(prs,["Utiliser une vraie API Météo","structurer un projet avec API","Traiter des données JSON complexes","Gérer les clés API en sécurité"])
    make_content(prs,"Le projet : Console Météo",[
        ("On va créer une app qui affiche la météo d'une ville.",15),
        "","Fonctionnalités :","  • Entrer le nom d'une ville","  • Afficher température, humidité, description","  • Afficher une prévision simple","  • Sauvegarder les recherches récentes",
        "","API utilisée : OpenWeatherMap (gratuite)","  • S'inscrire sur https://openweathermap.org/api","  • Récupérer une clé API (API Key) gratuite","  • Limite : 60 requêtes/minute (gratuit)",
        "","Structure :","meteo/","├── config.py      # Clé API + constantes","├── meteo_api.py   # Appels API","├── affichage.py   # Interface console","├── historique.py  # Sauvegarde CSV", "└── main.py        # Menu principal"
    ])
    make_code_content(prs,"config.py et .env",[
        "# config.py - NE PAS commiter la vraie clé !",
        "# Méthode 1 : variable d'environnement (recommandée)", "import os", "CLE_API = os.environ.get(\"OPENWEATHER_API_KEY\", \"\")",
        "",
        "# Méthode 2 : fichier .env (pas versionné)", "# pip install python-dotenv", "from dotenv import load_dotenv", "load_dotenv()  # charge .env", "CLE_API = os.getenv(\"OPENWEATHER_API_KEY\")",
        "",
        "# .env (à créer, à ajouter dans .gitignore)", "# OPENWEATHER_API_KEY=votre_clé_ici",
        "",
        "# Constantes", "URL_BASE = \"https://api.openweathermap.org/data/2.5\"", "UNITE = \"metric\"  # Celsius", "LANGUE = \"fr\"",
        "",
        "# Pour créer une clé :", "# 1. https://home.openweathermap.org/users/sign_up", "# 2. https://home.openweathermap.org/api_keys", "# 3. Copier la clé API"
    ])
    make_code_content(prs,"meteo_api.py - Appels API",[
        "# meteo_api.py", "import requests", "from config import CLE_API, URL_BASE, UNITE, LANGUE",
        "def meteo_actuelle(ville):", '    """Retourne la météo actuelle pour une ville."""', "    params = {", '        "q": ville, "appid": CLE_API,', '        "units": UNITE, "lang": LANGUE', "    }", "    try:", "        r = requests.get(f\"{URL_BASE}/weather\", params=params, timeout=5)", "        r.raise_for_status()", "        return r.json()", "    except requests.exceptions.HTTPError:", "        if r.status_code == 404:", '            return {"erreur": f"Ville {ville} inconnue"}', '        return {"erreur": f"Erreur HTTP {r.status_code}"}', "    except Exception as e:", '        return {"erreur": f"Erreur: {e}"}',
        "def previsions(ville, jours=3):", '    """Retourne les prévisions sur plusieurs jours."""', "    params = {", '        "q": ville, "appid": CLE_API,', '        "units": UNITE, "lang": LANGUE, "cnt": jours * 8', "    }", "    r = requests.get(f\"{URL_BASE}/forecast\", params=params)", "    return r.json()"
    ])
    make_code_content(prs,"affichage.py et formatage",[
        "# affichage.py", "from datetime import datetime",
        "def afficher_actuelle(data):", '    """Affiche la météo actuelle."""', '    if "erreur" in data:', '        print(f"❌ {data[\"erreur\"]}")', "        return",
        '    print(f"\\n🌍 {data[\"name\"]}, {data[\"sys\"][\"country\"]}")', '    print(f"🌡️ {data[\"main\"][\"temp\"]:.1f}°C "', '          f"(ressenti {data[\"main\"][\"feels_like\"]:.1f}°C)")', '    print(f"💧 Humidité: {data[\"main\"][\"humidity\"]}%")', '    print(f"🌬️ Vent: {data[\"wind\"][\"speed\"]} m/s")', '    print(f"☁️ {data[\"weather\"][0][\"description\"].capitalize()}")',
        "def afficher_previsions(data):", '    """Affiche les prévisions."""', "    for prev in data[\"list\"][::8]:  # 1 par jour", "        dt = datetime.fromtimestamp(prev[\"dt\"])", "        t = prev[\"main\"][\"temp\"]", "        desc = prev[\"weather\"][0][\"description\"]", '        print(f"  {dt:%d/%m}: {t:.1f}°C - {desc}")'
    ])
    make_code_content(prs,"historique.py et main.py",[
        "# historique.py", "import csv, os", 'FICHIER = "recherches.csv"', "def ajouter(ville):", "    with open(FICHIER, \"a\", newline=\"\", encoding=\"utf-8\") as f:", "        csv.writer(f).writerow([ville])",
        "def lister():", "    if not os.path.exists(FICHIER): return []", "    with open(FICHIER, \"r\", encoding=\"utf-8\") as f:", "        return [row[0] for row in csv.reader(f)]",
        "",
        "# main.py", "from meteo_api import meteo_actuelle, previsions", "from affichage import afficher_actuelle, afficher_previsions", "from historique import ajouter, lister",
        "def main():", '    while (ville := input("\\nVille (ou Q): ").strip()):', '        if ville.upper() == "Q": break', "        data = meteo_actuelle(ville)", "        afficher_actuelle(data)", '        if "erreur" not in data:', "            ajouter(ville)"
    ])
    make_ex(prs,"Exercice : Améliorer la météo",[
        "À partir du code fourni :\n\n1. Ajoute un affichage des prévisions sur 3 jours\n2. Ajoute une option \"historique\" qui montre les 5 dernières villes\n3. Ajoute la pression atmosphérique et le lever/coucher du soleil\n4. Affiche un petit graphique ASCII des températures de la semaine\n   (ex: ███ 15°C, ██████ 20°C...)\n5. Gère les erreurs : pas de clé API, pas de connexion, ville inconnue\n\nBONUS : Ajoute la possibilité de sauvegarder des villes favorites"
    ])
    make_summary(prs,["API OpenWeatherMap : météo temps réel et prévisions","Ne JAMAIS hardcoder les clés API (fichier .env)","Structurer : config, API, affichage, historique","Gérer les erreurs : HTTP, connexion, timeout","Un projet API = lecture de docs + parsing JSON"])
    save(prs,"43_Mini_Projet_Meteo.pptx")
    return prs

def p44_sqlite():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Bases de Données avec SQLite","Stocker des données structurées","Partie 5 - Projets Avancés")
    make_obj(prs,["Comprendre les bases de données","Découvrir SQLite","Écrire des requêtes SQL","Intégrer SQLite en Python"])
    make_content(prs,"Qu'est-ce qu'une base de données ?",[
        ("Une BD = un fichier qui stocke des données structurées de façon OPTIMISÉE.",15),
        "","Fichier JSON vs Base de données :",
        "","  JSON : simple, lisible, mais LENT pour chercher parmi 1M d'entrées","  SQLite : rapide, puissant, avec un langage de requête (SQL)",
        "","SQLite = base de données FICHIER (pas de serveur)","  • Intégré dans Python (module sqlite3)","  • Pas besoin d'installer quoi que ce soit","  • Un seul fichier .db","  • Parfait pour les petites applis, prototypes, jeux",
        "","Autres BD :","  • MySQL / PostgreSQL : serveur, multi-utilisateurs, production","  • MongoDB : NoSQL (documents JSON)",
        "","On va apprendre les BASES du SQL avec SQLite."
    ])
    make_content(prs,"Concepts SQL de base",[
        ("Une BD contient des TABLES (comme des feuilles Excel).",15),
        "","  Table UTILISATEURS :","  ┌────┬──────────┬─────────────┐","  │ id │  nom     │  email       │","  ├────┼──────────┼─────────────┤","  │ 1  │ Alice    │ a@b.com      │","  │ 2  │ Bob      │ b@c.com      │","  │ 3  │ Charlie  │ c@d.com      │","  └────┴──────────┴─────────────┘",
        "","Chaque ligne = un enregistrement (record)","Chaque colonne = un champ (field)","id = clé primaire (identifiant unique)",
        "","Opérations CRUD en SQL :","  CREATE → INSERT INTO (ajouter)","  READ   → SELECT  (lire)","  UPDATE → UPDATE  (modifier)","  DELETE → DELETE  (supprimer)"
    ])
    make_code_content(prs,"sqlite3 - Premiers pas",[
        "import sqlite3",
        "", "# 1. Connexion (crée le fichier s'il n'existe pas)", "conn = sqlite3.connect(\"ma_base.db\")",
        "", "# 2. Créer un curseur (pour exécuter des requêtes)", "cur = conn.cursor()",
        "", "# 3. Créer une table", "cur.execute('''", "    CREATE TABLE IF NOT EXISTS utilisateurs (", "        id INTEGER PRIMARY KEY AUTOINCREMENT,", "        nom TEXT NOT NULL,", "        email TEXT UNIQUE,", "        age INTEGER", "    )", "''')",
        "", "# 4. Sauvegarder (COMMIT) et fermer", "conn.commit()", "conn.close()",
        "", "# AUTOINCREMENT = id auto-généré", "# NOT NULL = champ obligatoire", "# UNIQUE = pas de doublons"
    ])
    make_code_content(prs,"CRUD : INSERT, SELECT, UPDATE, DELETE",[
        "conn = sqlite3.connect(\"ma_base.db\")", "cur = conn.cursor()",
        "", "# INSERT - Ajouter des données", 'cur.execute("INSERT INTO utilisateurs (nom, email, age) VALUES (?, ?, ?)",', '           ("Alice", "alice@mail.com", 25))',
        "# Le ? évite les injections SQL (NE JAMAIS utiliser f-string)", 'cur.execute("INSERT INTO utilisateurs VALUES (NULL, ?, ?, ?)", ("Bob", "bob@mail.com", 15))',
        "",
        "# INSERT multiple", "users = [(\"Charlie\",\"c@mail.com\",13), (\"David\",\"d@mail.com\",14)]", "cur.executemany(\"INSERT INTO utilisateurs VALUES (NULL,?,?,?)\", users)",
        "",
        "# SELECT - Lire", "cur.execute(\"SELECT * FROM utilisateurs\")", "print(cur.fetchall())  # tous les résultats",
        "cur.execute(\"SELECT nom, age FROM utilisateurs WHERE age >= 18\")", "for row in cur.fetchall():", "    print(f\"{row[0]}: {row[1]} ans\")",
        "conn.commit(); conn.close()"
    ])
    make_code_content(prs,"Requêtes SELECT avancées",[
        "conn = sqlite3.connect(\"ma_base.db\")", "cur = conn.cursor()",
        "", "# Filtrer avec WHERE", 'cur.execute("SELECT * FROM utilisateurs WHERE age > 13")',
        "", "# Trier avec ORDER BY", 'cur.execute("SELECT * FROM utilisateurs ORDER BY age DESC")',
        "", "# Compter", 'cur.execute("SELECT COUNT(*) FROM utilisateurs")', 'cur.execute("SELECT AVG(age) FROM utilisateurs")',
        "", "# Limiter", 'cur.execute("SELECT * FROM utilisateurs LIMIT 2")',
        "", "# LIKE (recherche texte)", 'cur.execute("SELECT * FROM utilisateurs WHERE nom LIKE ?", ("%li%",))',
        "", "# UPDATE - Modifier", 'cur.execute("UPDATE utilisateurs SET age = ? WHERE nom = ?", (16, "Alice"))',
        "", "# DELETE - Supprimer", 'cur.execute("DELETE FROM utilisateurs WHERE nom = ?", ("David",))',
        "conn.commit(); conn.close()"
    ])
    make_code_content(prs,"SQLite avec gestionnaire de contexte",[
        "import sqlite3",
        "", "# Avec 'with', la fermeture est automatique", "with sqlite3.connect(\"ma_base.db\") as conn:", "    cur = conn.cursor()",
        "    # On peut aussi accéder aux lignes comme des dictionnaires", "    conn.row_factory = sqlite3.Row",
        '    cur.execute("SELECT * FROM utilisateurs")', "    for row in cur.fetchall():", '        print(dict(row))  # {"id":1, "nom":"Alice", ...}',
        "",
        "# Ou avec une classe", "class BaseDonnees:", "    def __init__(self, fichier):", "        self.conn = sqlite3.connect(fichier)", "        self.conn.row_factory = sqlite3.Row",
        "    def requete(self, sql, params=()):", "        return self.conn.execute(sql, params).fetchall()",
        "    def fermer(self):", "        self.conn.close()",
        "bd = BaseDonnees(\"ma_base.db\")", "users = bd.requete(\"SELECT * FROM utilisateurs\")", "for u in users: print(u[\"nom\"])", "bd.fermer()"
    ])
    make_ex(prs,"Exercice : Gestionnaire de livres",[
        "Crée un gestionnaire de bibliothèque avec SQLite :\n\n1. Crée une table 'livres' avec :\n   - id (INTEGER PRIMARY KEY AUTOINCREMENT)\n   - titre (TEXT NOT NULL)\n   - auteur (TEXT)\n   - annee (INTEGER)\n   - lu (INTEGER DEFAULT 0)  # 0=non, 1=oui\n\n2. Fonctions :\n   - ajouter_livre(titre, auteur, annee)\n   - lister_tous()\n   - chercher_par_titre(mot)\n   - marquer_lu(id)\n   - statistiques() → nb total, nb lus, %\n\n3. Menu console pour tout ça\n\nBONUS : Ajoute une table 'genres' et une table de liaison"
    ])
    make_summary(prs,["SQLite = base de données fichier (intégré dans Python)","SQL = Structured Query Language","CRUD : INSERT, SELECT, UPDATE, DELETE","Toujours utiliser ? pour les paramètres (pas f-string)","conn.commit() pour sauvegarder","Les BD sont INDISPENSABLES pour stocker des données structurées"])
    save(prs,"44_SQLite.pptx")
    return prs

def p45_flask_intro():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Introduction à Flask","Créer des applications web avec Python","Partie 5 - Projets Avancés")
    make_obj(prs,["Comprendre le web (côté serveur)","Installer Flask","Créer des routes et vues","Afficher des pages HTML"])
    make_content(prs,"Du code console au web",[
        ("Jusqu'à présent : programmes en LIGNE DE COMMANDE.",15),
        "","  Terminal : python main.py → l'utilisateur tape des réponses",
        "","Avec Flask : on crée un SITE WEB !","  Navigateur : http://localhost:5000 → clics, formulaires, pages",
        "","Flask = micro-framework web Python","  • Léger, simple à apprendre","  • Suffisant pour faire des vrais sites","  • Extensible avec des extensions",
        "","Comment ça marche :","  1. Tu définis des ROUTES (URLs) dans ton code","  2. Chaque route est liée à une FONCTION","  3. La fonction retourne du HTML à afficher",
        "","  / (accueil) → def accueil(): return \"<h1>Bonjour !</h1>\""
    ])
    make_content(prs,"Installation et première app",[
        ("Prérequis :",15),
        "","  pip install flask  # (dans un venv !)",
        "","Plus simple d'installer :","  # Créer un venv d'abord", "  python -m venv .venv", "  source .venv/bin/activate", "  pip install flask",
        "",
        "Première app (app.py) :",
        "","  from flask import Flask", "  app = Flask(__name__)", "",
        "  @app.route(\"/\")", "  def accueil():", '      return "Hello, Web !"', "",
        "  if __name__ == \"__main__\":", "      app.run(debug=True)",
        "",
        "Lancer : python app.py → http://localhost:5000"
    ])
    make_code_content(prs,"Routes et paramètres",[
        "from flask import Flask", "app = Flask(__name__)",
        "", "# Route simple", "@app.route(\"/\")", "def accueil():", '    return "<h1>Bienvenue !</h1>"',
        "", "# Plusieurs routes", "@app.route(\"/a-propos\")", "def a_propos():", '    return "<p>Apprentissage Python</p>"',
        "", "# Route avec paramètre dans l'URL", "@app.route(\"/utilisateur/<nom>\")", "def profil(nom):", '    return f"<h2>Profil de {nom}</h2>"',
        "", "# Route avec paramètre typé", "@app.route(\"/carrer/<int:nombre>\")", "def carre(nombre):", '    return f"<p>{nombre}² = {nombre**2}</p>"',
        "", "# Routes avec plusieurs paramètres", "@app.route(\"/addition/<int:a>/<int:b>\")", "def addition(a, b):", '    return f"{a} + {b} = {a+b}"'
    ])
    make_code_content(prs,"Afficher du vrai HTML",[
        "from flask import Flask, render_template", "app = Flask(__name__)",
        "", "# Les templates HTML vont dans dossier templates/",
        "@app.route(\"/\")", "def accueil():", '    return render_template("accueil.html")',
        "",
        "@app.route(\"/bonjour/<nom>\")", "def bonjour(nom):", '    return render_template("bonjour.html", nom=nom, age=25)',
        "",
        "# templates/accueil.html :", "# <!DOCTYPE html>", "# <html>", "# <head><title>Accueil</title></head>", '# <body><h1>Bienvenue !</h1></body>', "# </html>",
        "",
        "# templates/bonjour.html :", "# <h1>Bonjour {{ nom }} !</h1>", "# <p>Tu as {{ age }} ans.</p>",
        "",
        "# {{ variable }} = Jinja2 (moteur de templates)", "# On peut aussi faire : {% if %} {% for %}"
    ])
    make_content(prs,"Jinja2 - Templates dynamiques",[
        ("Jinja2 = moteur de templates intégré à Flask.",15),
        "","  <h1>Bonjour {{ nom }}!</h1>",
        "  {% if age >= 18 %}", "    <p>Tu es majeur.</p>", "  {% else %}", "    <p>Tu es mineur.</p>", "  {% endif %}",
        "",
        "  <ul>", "  {% for item in liste %}", '    <li>{{ item }}</li>', "  {% endfor %}", "  </ul>",
        "",
        "Exemple complet dans python :",
        "  @app.route(\"/liste\")", "  def liste():", '      fruits = ["Pomme", "Banane", "Cerise"]', "      return render_template(\"liste.html\", fruits=fruits)"
    ])
    make_code_content(prs,"Formulaires et méthodes GET/POST",[
        "from flask import Flask, request, render_template", "app = Flask(__name__)",
        "@app.route(\"/\", methods=[\"GET\", \"POST\"])", "def accueil():", "    if request.method == \"POST\":", '        nom = request.form["nom"]', '        return f"Bonjour {nom}!"', '    return """', '        <form method="POST">', '            <input name="nom" placeholder="Ton nom">', '            <button type="submit">Valider</button>', '        </form>', '    """',
        "", "# request.form = données POST du formulaire", "# request.args = paramètres GET (dans l'URL)",
        "",
        "# Exemple avec GET", "@app.route(\"/search\")", "def search():", '    q = request.args.get("q", "")', "    return f\"Recherche: {q}\"",
        "# http://localhost:5000/search?q=python"
    ])
    make_ex(prs,"Exercice : Premier site Flask",[
        "1. Crée un dossier \"mon_site\" avec venv + flask installé\n2. Crée app.py avec routes :\n   - / → page d'accueil (\"Bienvenue sur mon site !\")\n   - /salut/<nom> → \"Salut nom !\"\n   - /age/<int:age> → dit si majeur ou mineur\n   - /table/<int:n> → table de multiplication de n\n\n3. Crée le dossier templates/ avec :\n   - base.html (structure HTML commune)\n   - accueil.html qui hérite de base.html\n\n4. Lance : python app.py\n5. Va sur http://localhost:5000/table/7\n\nBONUS : Ajoute du CSS (fichier static/style.css)"
    ])
    make_summary(prs,["Flask = micro-framework web Python","Route = URL liée à une fonction (décorateur @app.route)","render_template() pour afficher du HTML","Jinja2 : {{ variable }}, {% if %}, {% for %}","request.form pour les formulaires POST","Flask ouvre la porte aux applications WEB"])
    save(prs,"45_Flask_intro.pptx")
    return prs

def p46_mini_projet_blog():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Mini-Projet : Blog avec Flask","Articles, base de données et templates","Partie 5 - Projets Avancés")
    make_obj(prs,["Créer une app web complète","Combiner Flask + SQLite","Implémenter CRUD pour des articles","Utiliser Bootstrap pour le style"])
    make_content(prs,"Le projet : Mini Blog",[
        ("On va créer un blog minimaliste mais fonctionnel.",15),
        "","Fonctionnalités :","  • LISTER les articles (page d'accueil)","  • VOIR un article en détail","  • AJOUTER un nouvel article (formulaire)","  • MODIFIER un article","  • SUPPRIMER un article",
        "","Pas d'authentification (pour l'instant)","Utilisation de Bootstrap (CSS) pour un rendu moderne",
        "","Structure :","blog/","├── .venv/","├── app.py           # Application Flask","├── templates/","│   ├── base.html     # Layout commun","│   ├── index.html    # Liste des articles","│   ├── article.html  # Un article","│   ├── edit.html     # Formulaire","└── static/","    └── style.css"
    ])
    make_code_content(prs,"app.py - La base de données",[
        "import sqlite3", "from flask import Flask, render_template, request, redirect, url_for",
        "app = Flask(__name__)",
        "def get_db():", "    conn = sqlite3.connect(\"blog.db\")", "    conn.row_factory = sqlite3.Row", "    return conn",
        "def init_db():", "    with get_db() as db:", "        db.execute('''", "            CREATE TABLE IF NOT EXISTS articles (", "                id INTEGER PRIMARY KEY AUTOINCREMENT,", "                titre TEXT NOT NULL,", "                contenu TEXT NOT NULL,", "                date_creation TIMESTAMP DEFAULT CURRENT_TIMESTAMP", "            )", "        ''')",
        "init_db()"
    ])
    make_code_content(prs,"app.py - Les routes",[
        "@app.route(\"/\")", "def index():", "    with get_db() as db:", "        articles = db.execute(\"SELECT * FROM articles ORDER BY date_creation DESC\").fetchall()", "    return render_template(\"index.html\", articles=articles)",
        "",
        "@app.route(\"/article/<int:id>\")", "def article(id):", "    with get_db() as db:", "        art = db.execute(\"SELECT * FROM articles WHERE id=?\", (id,)).fetchone()", "    return render_template(\"article.html\", art=art)",
        "",
        "@app.route(\"/nouveau\", methods=[\"GET\", \"POST\"])", "def nouveau():", "    if request.method == \"POST\":", "        titre = request.form[\"titre\"]", "        contenu = request.form[\"contenu\"]", "        with get_db() as db:", "            db.execute(\"INSERT INTO articles (titre, contenu) VALUES (?,?)\", (titre, contenu))", "        return redirect(url_for(\"index\"))", "    return render_template(\"edit.html\", art=None)",
        "",
        "@app.route(\"/supprimer/<int:id>\")", "def supprimer(id):", "    with get_db() as db:", "        db.execute(\"DELETE FROM articles WHERE id=?\", (id,))", "    return redirect(url_for(\"index\"))"
    ])
    make_code_content(prs,"base.html et index.html",[
        "<!-- templates/base.html -->", "<!DOCTYPE html>", "<html lang=\"fr\">", "<head>", '    <meta charset=\"UTF-8\">', '    <title>{% block title %}Mon Blog{% endblock %}</title>', '    <link href=\"https://cdn.jsdelivr.net/npm/bootstrap@5.3.0/dist/css/bootstrap.min.css\" rel=\"stylesheet\">', "</head>", "<body class=\"bg-dark text-light\">", "    <nav class=\"navbar navbar-dark bg-primary\">", '        <div class=\"container\">', '            <a class=\"navbar-brand\" href=\"/\">📝 Mon Blog</a>', '        </div>', "    </nav>", "    <div class=\"container mt-4\">", "        {% block content %}{% endblock %}", "    </div>", "</body></html>",
        "",
        "<!-- templates/index.html -->", "{% extends \"base.html\" %}",
        "{% block content %}", '<h1>Articles</h1>', '<a href=\"/nouveau\" class=\"btn btn-success mb-3\">+ Nouvel article</a>',
        "{% for art in articles %}", '<div class=\"card bg-secondary mb-3\">', '    <div class=\"card-body\">', '        <h5 class=\"card-title\">{{ art.titre }}</h5>', '        <p class=\"card-text\">{{ art.contenu[:100] }}...</p>', '        <a href=\"/article/{{ art.id }}\" class=\"btn btn-primary\">Lire</a>', '        <a href=\"/supprimer/{{ art.id }}\" class=\"btn btn-danger\">Suppr</a>', "    </div>", "</div>", "{% endfor %}", "{% endblock %}"
    ])
    make_code_content(prs,"article.html et edit.html",[
        "<!-- templates/article.html -->", "{% extends \"base.html\" %}",
        "{% block content %}", '<h1>{{ art.titre }}</h1>', '<p class=\"text-muted\">{{ art.date_creation }}</p>', '<hr>', '<p>{{ art.contenu }}</p>', '<a href=\"/\" class=\"btn btn-secondary\">← Retour</a>', "{% endblock %}",
        "",
        "<!-- templates/edit.html -->", "{% extends \"base.html\" %}",
        "{% block content %}", '<h1>{% if art %}Modifier{% else %}Nouvel article{% endif %}</h1>',
        '<form method=\"POST\">', '    <div class=\"mb-3\">', '        <label>Titre</label>', '        <input name=\"titre\" class=\"form-control\" value=\"{{ art.titre if art else \"\" }}\">',
        '    </div>', '    <div class=\"mb-3\">', '        <label>Contenu</label>', '        <textarea name=\"contenu\" class=\"form-control\" rows=\"10\">{{ art.contenu if art else \"\" }}</textarea>', '    </div>', '    <button type=\"submit\" class=\"btn btn-primary\">Sauvegarder</button>', '</form>', "{% endblock %}"
    ])
    make_ex(prs,"Exercice : Blog amélioré",[
        "Étapes :\n1. Reproduis le code du blog complet\n2. Lance-le : python app.py\n3. Crée 3 articles\n4. Améliorations possibles :\n   • Modifier un article (route /editer/<id>)\n   • Ajouter une barre de recherche\n   • Ajouter des catégories / tags\n   • Ajouter des commentaires (nouvelle table + formulaire)\n   • Ajouter un compteur de vues\n   • Tronquer le contenu dans la liste (déjà fait avec [:100])\n\nBONUS : Ajoute l'authentification avec flask-login"
    ])
    make_summary(prs,["Blog Flask complet : CRUD d'articles","SQLite pour la persistance","Bootstrap pour le style sans effort CSS","Jinja2 pour les templates dynamiques","Architecture : routes → BD → templates","Un blog est le \"Hello World\" des apps web"])
    save(prs,"46_Mini_Projet_Blog.pptx")
    return prs

def p47_debogage_tests():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Débogage et Tests","Écrire du code fiable","Partie 5 - Projets Avancés")
    make_obj(prs,["Maîtriser le débogueur Python","Écrire des tests unitaires","Utiliser assert et pytest","Comprendre TDD (Test Driven Development)"])
    make_content(prs,"Pourquoi tester son code ?",[
        ("Les bugs arrivent à TOUT le monde, même aux meilleurs développeurs.",15),
        "","Tester son code = VÉRIFIER qu'il fait ce qu'on attend.",
        "","Sans tests :","  • Tu changes une ligne → tu ne sais pas si tout le reste marche","  • Tu ajoutes une fonction → tu découvres 3 bugs cachés","  • Tu passes 2h à trouver un bug que tu aurais pu détecter en 5min",
        "","Avec tests :","  ✅ Confiance pour modifier le code","  ✅ Détection rapide des bugs","  ✅ Documentation vivante (le test montre comment utiliser le code)","  ✅ Meilleure conception (un code testable est mieux conçu)",
        "","Types de tests :","  • Tests UNITAIRE : une fonction spécifique","  • Tests d'INTÉGRATION : plusieurs composants ensemble","  • Tests E2E : parcours complet utilisateur"
    ])
    make_code_content(prs,"Déboguer avec breakpoint()",[
        "# Technique 1 : breakpoint() - Python 3.7+",
        "def calculer_prix(prix_ht, tva=0.2):", "    breakpoint()  # ➡️ pause ici !", "    tva_calculee = prix_ht * tva", "    total = prix_ht + tva_calculee", "    return total",
        "print(calculer_prix(100))",
        "",
        "# Quand le programme arrive à breakpoint() :",
        "# (Pdb) est une console interactive !",
        "# Commandes utiles :",
        "#   p variable    → print(variable)", "#   n            → next (ligne suivante)", "#   s            → step into (entre dans la fonction)", "#   c            → continue (jusqu'au prochain breakpoint)", "#   l            → list (affiche le code autour)", "#   q            → quit",
        "#   p dir(obj)   → inspecte un objet", "#   p locals()   → toutes les variables locales"
    ])
    make_code_content(prs,"Tests unitaires avec unittest",[
        "import unittest", "from mon_module import addition, division",
        "class TestCalculs(unittest.TestCase):", "    def test_addition(self):", "        self.assertEqual(addition(2, 3), 5)", "        self.assertEqual(addition(-1, 1), 0)",
        "    def test_division(self):", "        self.assertEqual(division(10, 2), 5)", "        self.assertAlmostEqual(division(7, 3), 2.333, places=2)",
        "    def test_division_par_zero(self):", "        with self.assertRaises(ValueError):", "            division(5, 0)",
        "    def test_parametres_defaut(self):", "        self.assertEqual(addition(5), 5)  # b=0 par défaut",
        "if __name__ == '__main__':", "    unittest.main()",
        "# Assertions courantes :", "# assertEqual, assertTrue, assertFalse", "# assertIn, assertNotIn, assertIsNone", "# assertRaises, assertAlmostEqual"
    ])
    make_code_content(prs,"Tests avec pytest (plus moderne)",[
        "# pip install pytest",
        "", "# test_calculs.py (pas besoin de classe !)", "from mon_module import addition, division",
        "def test_addition():", "    assert addition(2, 3) == 5", "    assert addition(-1, 1) == 0", "    assert addition(5) == 5",
        "def test_division():", "    assert division(10, 2) == 5", "    assert division(7, 3) == 2.333  # marche pas pour floats",
        "def test_division_par_zero():", "    import pytest", "    with pytest.raises(ValueError):", "        division(5, 0)",
        "def test_parametres():", "    assert addition(5) == 5  # valeur par défaut de b",
        "# Lancer : pytest -v", "# Résultat : tous les ✓ ou ✗ avec explications"
    ])
    make_content(prs,"TDD - Test Driven Development",[
        ("TDD = écrire le test AVANT le code.",15),
        "","Cycle TDD :","  1. 🔴 Écris un test qui ÉCHOUE (car la fonction n'existe pas)","  2. ✅ Écris le code MINIMUM pour faire passer le test","  3. ♻️ Refactore (améliore) le code sans casser le test",
        "","Exemple concret :",
        "  # Étape 1 : j'écris le test", "  def test_factorial():", "      assert factorial(0) == 1", "      assert factorial(5) == 120",
        "  # Étape 2 : j'écris juste assez pour que ça passe", "  def factorial(n):", "      if n == 0: return 1", "      return n * factorial(n-1)",
        "  # Étape 3 : je peux améliorer (ajouter validation)",
        "Avantages TDD :", "  • Code mieux conçu (testable dès le départ)","  • Moins de bugs (les tests sont écrits AVANT)","  • Documentation automatique"
    ])
    make_ex(prs,"Exercice : Tests pour le Projet Pendu",[
        "Reprends le projet Pendu (p40) et écris des tests :\n\n1. Crée test_pendu.py\n2. Teste la classe Pendu :\n   - test_proposer_lettre_trouvee()\n   - test_proposer_lettre_ratee()\n   - test_proposer_deux_fois_la_meme()\n   - test_est_gagne()\n   - test_est_perdu()\n   - test_plusieurs_lettres()\n\n3. Lance : pytest -v test_pendu.py\n4. Ajoute un test qui échoue volontairement pour voir\n\nBONUS : Calcule le % de couverture de code (pip install coverage)"
    ])
    make_summary(prs,["breakpoint() pour déboguer interactivement","Tests unitaires : unittest (intégré) ou pytest (moderne)","Assertions : assertEqual, assertTrue, assertRaises...","TDD : Test d'abord, code après","Les tests donnent CONFIANCE pour modifier le code","Un projet sans tests = un projet fragile"])
    save(prs,"47_Debogage_Tests.pptx")
    return prs

def p48_bonnes_pratiques():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Bonnes Pratiques Python","Écrire du code professionnel","Partie 5 - Projets Avancés")
    make_obj(prs,["Suivre PEP 8 (style de code)","Écrire des docstrings","Respecter DRY, KISS, SOLID","Organiser un projet Python"])
    make_content(prs,"PEP 8 - Le style Python",[
        ("PEP 8 = Python Enhancement Proposal 8 = guide de style officiel.",15),
        "","Règles principales :","",
        "✅ Indentation : 4 espaces (PAS de tabulations)","   def fonction():","       return 42",
        "✅ Limite de ligne : 79 caractères max","   resultat = calcul_complexe(parametre1, parametre2,",
        "                            parametre3)",
        "✅ Espaces autour des opérateurs :","   x = 1 + 2  # ✅", "   y=1+2      # ❌",
        "✅ Noms de variables : snake_case","   nom_utilisateur = \"Alice\"  # ✅", "   nomUtilisateur = \"Alice\"  # ❌ (camelCase en Java)",
        "✅ Noms de classes : PascalCase","   class CompteBancaire:  # ✅", "   class compte_bancaire:  # ❌"
    ])
    make_content(prs,"Règles de nommage",[
        ("Conventions Python :",15),
        "","  Variables / fonctions : snake_case","    nom_utilisateur, calculer_moyenne()",
        "  Constantes : MAJUSCULES","    TAUX_TVA = 0.20, VITESSE_LUMIERE = 3e8",
        "  Classes : PascalCase","    MonException, CompteBancaire",
        "  Méthodes privées : _prefixe (convention)","    def _valider_donnees(self):",
        "  Variables privées : __double_prefixe (name mangling)","    self.__mot_de_passe",
        "  Fichiers : snake_case.py","    mon_module.py, test_calculs.py",
        "  Dossiers : snake_case","    mes_outils/, tests/",
        "","🔑 Consistance > tout : choisis un style et applique-le PARTOUT"
    ])
    make_code_content(prs,"Docstrings - Documenter son code",[
        "# Docstring = string multi-lignes après def/class/module",
        "def calculer_moyenne(notes):", '    """Calcule la moyenne d\'une liste de notes.', "",
        "    Args:", "        notes: Liste de nombres (0-20).",
        "", "    Returns:", "        float: La moyenne des notes.",
        "", "    Raises:", "        ValueError: Si la liste est vide.", '    """',
        "    if not notes:", "        raise ValueError(\"Liste vide\")", "    return sum(notes) / len(notes)",
        "",
        "# Voir la doc :", "help(calculer_moyenne)", "print(calculer_moyenne.__doc__)",
        "",
        "# Docstring de module (en haut du fichier)", '"""Module pour les calculs mathématiques du projet."""',
        "# Docstring de classe", "class Chien:", '    """Représente un chien avec son nom et son âge."""', "    pass"
    ])
    make_content(prs,"Principes DRY, KISS, YAGNI",[
        ("DRY = Don't Repeat Yourself",15),
        "  ❌ NE COPIE PAS le même code à plusieurs endroits","  ✅ Crée une FONCTION réutilisable",
        "","  ❌ (calcul TVA partout)","  ✅ def appliquer_tva(prix): return prix * 1.20",
        "","KISS = Keep It Simple, Stupid","  ❌ Code trop complexe avec 15 conditions imbriquées","  ✅ Simple, lisible, compréhensible",
        "","YAGNI = You Ain't Gonna Need It","  ❌ Ajouter des fonctionnalités \"au cas où\"","  ✅ N'ajoute que ce dont tu as BESOIN MAINTENANT",
        "","  \"Le débogage est deux fois plus dur que l'écriture de code.\"","  \"Donc si tu écris le code le plus malin possible,", "   tu es par définition incapable de le débuguer.\" — Kernighan"
    ])
    make_content(prs,"Structure de projet",[
        ("Structure recommandée pour un projet :",15),
        "","mon_projet/","├── README.md          # Description du projet","├── .gitignore          # Fichiers ignorés par Git","├── requirements.txt    # Dépendances pip","├── .venv/              # Environnement virtuel (ignoré)",
        "├── config.py          # Configuration","├── main.py             # Point d'entrée","├── mon_module/         # Code source","│   ├── __init__.py","│   ├── utils.py","│   └── classes.py",
        "└── tests/            # Tests","    ├── __init__.py","    ├── test_utils.py","    └── test_classes.py",
        "","💡 Un projet bien organisé est plus FACILE à maintenir","  et plus AGRÉABLE à reprendre après 6 mois d'absence"
    ])
    make_ex(prs,"Exercice : Refactoring",[
        'Refactore ce code (rends-le propre) :\n\ndef m(a,b,c):\n    x=a+b+c;y=x/3;print("moy:"+str(y))\n    if y>=10:print("pass");print("bravo")\n    else:print("echec");return y',
        "","Corrige :",
        "1. Renomme la fonction et les variables correctement","2. Ajoute une docstring complète","3. Sépare les responsabilités (calcul ≠ affichage)","4. Utilise f-string au lieu de +","5. Ajoute des annotations de type","6. Ajoute la gestion d'erreurs","7. Écris un test pour la fonction",
        "","def calculer_moyenne(notes: list[float]) -> float:", '    """Calcule..."""', "    ..."
    ])
    make_summary(prs,["PEP 8 : 4 espaces, snake_case, 79 caractères, espaces","Docstrings : documente TOUT ce qui est réutilisable","DRY : ne te répète pas, crée des fonctions","KISS : simple > complexe","YAGNI : pas de fonctionnalités inutiles","Un code propre = un développeur heureux"])
    save(prs,"48_Bonnes_pratiques.pptx")
    return prs

def p49_projet_final():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Projet Final - Guide Complet","Mets en pratique TOUT ce que tu as appris","Partie 5 - Projets Avancés")
    make_obj(prs,["Choisir un projet qui te passionne","Planifier et structurer","Appliquer tout ce qu'on a vu","Publier sur GitHub"])
    make_content(prs,"Le projet final - Pourquoi ?",[
        ("C'est le moment de montrer ce que tu sais faire !",15),
        "","Ce projet doit :","  ✅ Utiliser la PLUPART des concepts vus dans le cours","  ✅ Être quelque chose qui t'INTÉRESSE VRAIMENT","  ✅ Être PUBLIÉ sur GitHub",
        "","Choisis PARMI les idées ci-dessous ou INVENTE LA TIENNE :",
        "","🎮 Jeux : Snake, Tic-Tac-Toe, Démineur, Jeu de rôle","📊 Outils : Gestionnaire de budget, Suivi d'habitudes","🌐 Web : To-do list en Flask, Portfolio, Wiki","📱 Utilitaires : Rappels, Organiseur de fichiers, Convertisseur",
        "","Critères d'évaluation :","  • Fonctionnel (ça marche !)","  • Code propre (PEP 8, docstrings)","  • Organisé (modules, fichiers séparés)","  • Testé (au moins quelques tests)","  • Versionné (Git + GitHub)"
    ])
    make_content(prs,"Idée 1 : Jeu Snake",[
        ("Le classique !",15),
        "","Technologies : pygame (pip install pygame)","Concepts : POO, boucle de jeu, événements clavier, score",
        "","Fonctionnalités :","  • Serpent qui grandit en mangeant","  • Score qui s'affiche","  • Game Over quand le serpent touche le bord/lui-même","  • Sauvegarde du meilleur score (JSON/BD)","  • Difficulté croissante (vitesse augmente)",
        "","Structure :","snake/","├── main.py","├── snake.py (classe Serpent)","├── pomme.py (classe Pomme)","├── jeu.py (boucle principale)","├── scores.py (sauvegarde)","└── config.py (couleurs, vitesse, etc.)",
        "","💡 Pourquoi Snake ? : POO, événements, animation, fun !"
    ])
    make_content(prs,"Idée 2 : Application Web de Quiz",[
        ("Crée un site de quiz interactif avec Flask + SQLite.",15),
        "","Fonctionnalités :","  • Inscription / Connexion (sessions)","  • CRUD des quiz et questions (admin)","  • Jouer : répondre aux questions, timer","  • Scoreboard : classement des joueurs","  • Catégories de quiz",
        "","Base de données :","  • table utilisateurs","  • table quiz (titre, catégorie)","  • table questions (texte, réponses, bonne réponse)","  • table scores (utilisateur, quiz, score, date)",
        "","Pages :","  / → accueil, /quiz → liste, /quiz/<id> → jouer","  /admin → créer/modifier quiz","  /scores → classement",
        "","💡 Pourquoi Quiz ? : Flask + BD + auth + CRUD complet"
    ])
    make_content(prs,"Idée 3 : Gestionnaire de Budget",[
        ("Suis tes dépenses avec des graphiques.",15),
        "","Technologies : Python + SQLite + matplotlib (graphiques)","Concepts : fichiers, BD, dates, statistiques",
        "","Fonctionnalités :","  • Ajouter une dépense (montant, catégorie, date)","  • Voir le solde actuel","  • Graphiques par catégorie (camembert)","  • Graphiques par mois (barres)","  • Exporter en CSV","  • Budget mensuel avec alertes",
        "","Structure :","budget/","├── main.py","├── depenses.py (classes Dépense, Budget)","├── base.py (SQLite CRUD)","├── graphiques.py (matplotlib)","└── export.py (CSV)",
        "","💡 Pourquoi Budget ? : BD + data viz + très utile dans la vie !"
    ])
    make_content(prs,"Plan de travail recommandé",[
        ("Semaine 1 - Planification (ne code pas encore) :",15),
        "  • Choisis ton projet","  • Écris les fonctionnalités sur papier","  • Dessine l'architecture (fichiers, classes, BD)","  • Crée le dépôt GitHub avec README.md",
        "","Semaine 2 - Base :","  • Structure du projet (dossiers, fichiers, venv)","  • Fonctionnalité de base qui marche","  • Premier commit fonctionnel",
        "","Semaine 3 - Améliorations :","  • Ajoute les fonctionnalités secondaires","  • Améliore l'interface (Couleurs, ergonomie)","  • Commits réguliers",
        "","Semaine 4 - Polissage :","  • Tests (au moins les fonctions principales)","  • Gestion des erreurs","  • README.md complet avec instructions","  • Push final sur GitHub",
        "","🎯 L'important : finir ! Même imparfait, un projet fini > parfait abandonné"
    ])
    make_ex(prs,"Checklist finale",[
        "Avant de considérer le projet terminé :\n\n☐ Le programme fonctionne sans erreur\n☐ Le code suit PEP 8\n☐ Les fonctions ont des docstrings\n☐ Les noms de variables sont clairs\n☐ Le projet est organisé en plusieurs fichiers\n☐ Il y a un .gitignore\n☐ Il y a un requirements.txt\n☐ Il y a un README.md avec :\n    • Description du projet\n    • Comment installer / lancer\n    • Comment utiliser\n    • Technologies utilisées\n☐ Le projet est sur GitHub\n☐ Tu es FIER du résultat !  🎉\n\n⚠ Si une case est vide : ce n'est pas grave, mais note-la comme amélioration future"
    ])
    make_summary(prs,["Choisis un projet qui te PASSIONNE","Planifie avant de coder (papier + crayon)","Commence SIMPLE, ajoute des fonctionnalités ensuite","Commit régulièrement, push sur GitHub","Un projet FINI vaut mieux qu'un projet PARFAIT","Le meilleur apprentissage = construire des projets !"])
    save(prs,"49_Projet_Final.pptx")
    return prs

def p50_prochaines_etapes():
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    make_title(prs,"Prochaines Étapes","Après ce cours, où aller ?","Partie 5 - Projets Avancés")
    make_obj(prs,["Récapituler le parcours","Découvrir les prochains sujets","Trouver des ressources pour continuer","S'engager dans la communauté"])
    make_content(prs,"Ton parcours - Rétrospective",[
        ("🏁 FÉLICITATIONS ! Tu es arrivé à la fin de ce cours.",15),
        "",
        "Rappelle-toi d'où tu viens :",
        "","  01 🖥️  Qu'est-ce que l'informatique","  02 ⚙️  Comment marche un ordinateur","  03 🖱️  Systèmes d'exploitation","  04 🌐  Internet et le Web","  05 📐  Algorithmes et Programmation","  06 🔧  Installation Python & VS Code","  07 💻  Terminal / Ligne de commande","  08 👋  Hello World !",
        "  09-12 📦  Variables, Strings, Nombres, Input","  13 🔀  Conditions, if/else","  14 🏋️  Exercices Conditions",
        "  15 🔄  Boucles while","  16 🔁  Boucles for","  17-18 📚  Listes, Dictionnaires","  19 🏋️  Exercices Listes & Boucles","  20 🧮  Mini-Projet : Calculatrice"
    ])
    make_content(prs,"Ton parcours - Suite",[
        ("Ce que tu maîtrises maintenant :",15),
        "","  21-23 🔧  Fonctions (définition, paramètres, portée)","  24 📦  Tuples et Sets","  25 ⚠️  Gestion d'erreurs","  26 📄  Fichiers","  27 📦  Modules et Import",
        "  28 📚  Bibliothèques utiles (random, datetime...)","  29 ✅  Mini-Projet : Gestionnaire de Tâches","  30 🏋️  Exercices de Révision",
        "  31-32 🎯  POO (Classes, Héritage)","  33-34 💡  Comprehensions, Lambda, map/filter",
        "  35 🐛  Exceptions avancées, Débogage","  36 🃏  Mini-Projet : Jeu de Cartes",
        "  37-39 🔗  Git (intro, commandes, branches, GitHub)","  40 🎮  Mini-Projet : Pendu",
        "  41 📦  Environnements virtuels","  42 🌍  APIs HTTP",
        "  43 🌤️  Mini-Projet : Météo","  44 🗄️  SQLite",
        "  45 🌐  Flask (Web)","  46 📝  Mini-Projet : Blog",
        "  47 🧪  Tests et Débogage","  48 ✨  Bonnes pratiques (PEP 8, DRY, KISS)",
        "  49 🏆  Projet Final"
    ])
    make_content(prs,"Prochains sujets à explorer",[
        ("Tu es prêt pour passer au NIVEAU SUPÉRIEUR.",15),
        "","🔜 Frameworks Web avancés :","  • Django (framework web complet, très utilisé)","  • FastAPI (API modernes et rapides)","  • Plus de Flask (extensions, déploiement)",
        "","🔜 Data Science :","  • NumPy (calculs scientifiques)","  • Pandas (analyse de données)","  • Matplotlib / Seaborn (visualisation)","  • Jupyter Notebooks",
        "","🔜 Développement de jeux :","  • Pygame (déjà vu!) → va plus loin","  • Godot (moteur de jeu complet)","  • Ren'Py (visual novels)",
        "","🔜 Spécialisations :","  • Sécurité / Hacking éthique (TryHackMe)","  • Automatisation (scripts, bots)","  • Mobile (Kivy, BeeWare)"
    ])
    make_content(prs,"Ressources pour continuer",[
        ("Sites gratuits pour apprendre :",15),
        "","🇫🇷 En français :","  • OpenClassrooms - cours complets","  • Doc Python officielle (traduite)","  • Grafikart - tutoriels vidéo",
        "","🇬🇧 En anglais (le plus de ressources) :","  • Real Python - articles de qualité","  • Python.org - tuto officiel","  • W3Schools - références interactives","  • freeCodeCamp - cours vidéo gratuits",
        "","📚 Livres recommandés :","  • \"Apprendre Python 3\" (Zed Shaw)","  • \"Automate the Boring Stuff with Python\" (très pratique)","  • \"Python Crash Course\" (complet)",
        "","🎯 Plateformes d'exercices :","  • Codewars (défis progressifs)","  • LeetCode (préparation entretiens)","  • Exercism (avec mentorat gratuit)"
    ])
    make_content(prs,"Rejoindre la communauté",[
        ("La programmation ne s'apprend pas seul !",15),
        "","Rejoins la communauté :",
        "","💬 Discord / Slack :","  • Python Discord (énorme communauté)","  • Serveurs de programmation français",
        "","🌍 GitHub :","  • Regarde le code des autres","  • Contribue à des projets open source","  • Montre TES projets",
        "","🐦 Twitter/X :","  • Suis des développeurs Python","  • Partage tes progrès",
        "","📰 Stack Overflow :","  • Cherche avant de poser une question","  • Apprends à poser des questions CLAIRES",
        "","💡 Astuce : le meilleur moyen d'apprendre est d'ENSEIGNER","  • Explique ce que tu as appris à quelqu'un d'autre","  • Crée un blog ou une vidéo tuto","  • Aide les débutants sur les forums"
    ])
    make_content(prs,"Conseils pour la suite",[
        ("10 conseils pour continuer à progresser :",15),
        "","1. CODE TOUS LES JOURS ! 15 minutes > 5 heures une fois par mois","2. Finis ce que tu commences (même imparfait)","3. Lis le code des autres (GitHub, open source)","4. Fais des projets personnels (pas que des tutos)","5. Ne copie pas, COMPRENDS (refais avec tes mots)","6. Échoue souvent, ça fait partie du processus","7. Reviens sur ton ancien code et améliore-le","8. Note ce que tu apprends (blog, journal, fiches)","9. Enseigne ce que tu sais (meilleur moyen d'apprendre)","10. 🎯 AMUSE-TOI ! La programmation est créative et passionnante",
        "","Rappelle-toi : le meilleur programmeur n'est pas celui qui","connaît toutes les syntaxes, mais celui qui sait CHERCHER","et APPRENDRE par lui-même."
    ])
    make_content(prs,"Merci et Bonne continuation !",[
        ("C'est la fin de ce cours, mais le début de ton aventure.",15),
        "","","Tu as appris :","  ✅ Les bases de l'informatique","  ✅ La programmation Python du débutant à avancé","  ✅ Les outils du développeur (Git, terminal, VS Code)","  ✅ La création d'applications (Flask, BD, APIs)","  ✅ À structurer et tester ton code",
        "","Rien ne peut t'arrêter maintenant. Continue à coder,","continue à apprendre, et surtout... amuse-toi !",
        "","🎉 BRAVO POUR TOUT CE TRAVAIL ! 🎉"
    ])
    make_summary(prs,["50 présentations parcourues du débutant à l'avancé","Python, Git, POO, APIs, Flask, BD, tests","Projets pratiques : Calculatrice, Pendu, Blog, Météo, etc.","Continue avec Django, Data Science, Jeux","Rejoins la communauté, partage ton code","Code tous les jours, finis tes projets, amuse-toi !"])
    save(prs,"50_Prochaines_etapes.pptx")
    return prs

if __name__ == "__main__":
    funcs = [p43_mini_projet_meteo, p44_sqlite, p45_flask_intro,
             p46_mini_projet_blog, p47_debogage_tests, p48_bonnes_pratiques,
             p49_projet_final, p50_prochaines_etapes]
    for f in funcs:
        print(f"\n Génération : {f.__name__}")
        try:
            f()
        except Exception as e:
            print(f"  ❌ Erreur dans {f.__name__}: {e}")
    print(f"\n{'='*50}\n✅ {len(funcs)} présentations générées !\n📁 Dossier : {OUT}\n{'='*50}")
